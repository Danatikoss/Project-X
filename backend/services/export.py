"""
PPTX and PDF export service.

Strategy (Variant C): full slide cloning via python-pptx OPC layer.
For each PPTX-sourced slide:
  1. Open the original PPTX as a Presentation
  2. Deep-copy the slide's spTree (shapes)
  3. Copy ALL related parts (images, video, GIF, audio, charts...)
  4. Remap rId references in slide XML
For PDF-sourced slides: embed thumbnail as full-page image.

Media overlays (GIF/video positioned on top):
  - PPTX: add_picture() on top of the cloned slide (GIF = animated, video = poster frame)
  - PDF: PIL composite onto thumbnail before converting to PDF page
"""
import copy
import io
import json
import logging
import uuid
from pathlib import Path, PurePosixPath
from typing import Optional

from sqlalchemy.orm import Session

from config import settings
from models.assembly import AssembledPresentation
from models.slide import SlideLibraryEntry, SourcePresentation

logger = logging.getLogger(__name__)

_SKIP_RELTYPES = {
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
}

# ── Media overlay helpers ─────────────────────────────────────────────────────

def _media_path_from_url(url: str) -> Optional[Path]:
    """Convert /media-files/{filename} overlay URL to absolute filesystem path."""
    prefix = "/media-files/"
    if not url.startswith(prefix):
        return None
    rel = url[len(prefix):]
    p = Path(settings.upload_dir) / "media" / rel
    return p if p.exists() else None


def _open_as_pil(path: Path):
    """Open an image/GIF as PIL RGBA (first frame). Returns None for videos/unknowns."""
    try:
        from PIL import Image
        img = Image.open(str(path))
        if hasattr(img, 'n_frames') and img.n_frames > 1:
            img.seek(0)
        return img.convert("RGBA")
    except Exception:
        return None


def _get_image_size(path: Path) -> tuple[int, int] | None:
    """Return (width, height) of an image file without fully decoding it."""
    try:
        from PIL import Image
        with Image.open(str(path)) as img:
            return img.size
    except Exception:
        return None


def _extract_video_frame(path: Path) -> "io.BytesIO | None":
    """Extract the first frame of a video file as a PNG BytesIO using ffmpeg.
    Returns None if ffmpeg is unavailable or extraction fails.
    """
    import subprocess
    import shutil
    if not shutil.which("ffmpeg"):
        return None
    try:
        result = subprocess.run(
            [
                "ffmpeg", "-y", "-i", str(path),
                "-vframes", "1", "-f", "image2", "-vcodec", "png", "pipe:1",
            ],
            capture_output=True,
            timeout=15,
        )
        if result.returncode == 0 and result.stdout:
            buf = io.BytesIO(result.stdout)
            buf.seek(0)
            return buf
    except Exception as e:
        logger.debug(f"ffmpeg frame extraction failed: {e}")
    return None


def _contain_rect(
    nat_w: int, nat_h: int,
    cx: int, cy: int, cw: int, ch: int
) -> tuple[int, int, int, int]:
    """
    Compute object-contain placement of an image inside a container.
    Returns (x, y, w, h) — the actual rectangle where the image should be drawn,
    centered inside the container with aspect ratio preserved.
    """
    if nat_w <= 0 or nat_h <= 0 or cw <= 0 or ch <= 0:
        return cx, cy, cw, ch
    img_ar = nat_w / nat_h
    container_ar = cw / ch
    if img_ar > container_ar:
        # wider image → fit by width
        dw = cw
        dh = max(1, round(cw / img_ar))
    else:
        # taller image → fit by height
        dh = ch
        dw = max(1, round(ch * img_ar))
    off_x = (cw - dw) // 2
    off_y = (ch - dh) // 2
    return cx + off_x, cy + off_y, dw, dh


def _add_overlays_pptx(dest_prs, dest_slide, slide_id: str, overlays_map: dict):
    """Add media overlay shapes on top of an already-added PPTX slide."""
    slide_overlays = overlays_map.get(slide_id, [])
    if not slide_overlays:
        return

    sw = dest_prs.slide_width   # EMU
    sh = dest_prs.slide_height  # EMU

    for ov in slide_overlays:
        try:
            left = int(ov["x"] / 100 * sw)
            top  = int(ov["y"] / 100 * sh)
            w    = int(ov["w"] / 100 * sw)
            h    = int(ov["h"] / 100 * sh)

            path = _media_path_from_url(ov.get("url", ""))
            if not path:
                logger.debug(f"Overlay media not found: {ov.get('url')}")
                continue

            file_type = ov.get("file_type", "")

            if file_type in ("gif", "image"):
                # Apply object-contain sizing: place at actual natural-AR dimensions
                size = _get_image_size(path)
                if size:
                    left, top, w, h = _contain_rect(size[0], size[1], left, top, w, h)
                dest_slide.shapes.add_picture(str(path), left, top, w, h)
            elif file_type == "video":
                # Embed video as a proper playable shape using add_movie.
                # Extract poster frame via ffmpeg to get natural dimensions and preview image.
                from pptx.opc.constants import CT as _CT
                frame_buf = _extract_video_frame(path)
                poster_buf = None
                if frame_buf:
                    from PIL import Image as _PIL
                    frame_img = _PIL.open(frame_buf).convert("RGB")
                    nat_w, nat_h = frame_img.width, frame_img.height
                    al, at, aw, ah = _contain_rect(nat_w, nat_h, left, top, w, h)
                    poster_buf = io.BytesIO()
                    frame_img.save(poster_buf, "JPEG")
                    poster_buf.seek(0)
                else:
                    al, at, aw, ah = left, top, w, h

                try:
                    dest_slide.shapes.add_movie(
                        str(path), al, at, aw, ah,
                        poster_frame_image=poster_buf,
                        mime_type=_CT.MP4,
                    )
                except Exception as e2:
                    logger.debug(f"add_movie failed: {e2}")
                    if poster_buf:
                        poster_buf.seek(0)
                        dest_slide.shapes.add_picture(poster_buf, al, at, aw, ah)
                    else:
                        _add_video_placeholder_pptx(dest_slide, left, top, w, h)
        except Exception as e:
            logger.debug(f"Overlay PPTX add failed: {e}")


def _add_video_placeholder_pptx(slide, left, top, width, height):
    """Add a dark rect with a play icon as a stand-in for a video overlay."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    tf = shape.text_frame
    tf.text = "▶"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xff, 0xff, 0xff)


def _composite_overlays_pil(base_img, slide_id: str, overlays_map: dict):
    """Composite media overlays onto a PIL RGBA image in-place. Returns modified image."""
    slide_overlays = overlays_map.get(slide_id, [])
    if not slide_overlays:
        return base_img

    from PIL import Image
    W, H = base_img.size
    result = base_img.copy()

    for ov in slide_overlays:
        try:
            x = int(ov["x"] / 100 * W)
            y = int(ov["y"] / 100 * H)
            w = max(1, int(ov["w"] / 100 * W))
            h = max(1, int(ov["h"] / 100 * H))

            path = _media_path_from_url(ov.get("url", ""))
            if not path:
                continue

            frame = _open_as_pil(path)
            if frame is None:
                # Try extracting first frame from video via ffmpeg
                vid_buf = _extract_video_frame(path)
                if vid_buf:
                    frame = Image.open(vid_buf).convert("RGBA")
            if frame is None:
                # Video placeholder: dark rect
                placeholder = Image.new("RGBA", (w, h), (26, 26, 46, 200))
                result.paste(placeholder, (x, y), placeholder)
                continue

            # Apply object-contain sizing to match what the editor shows
            ax, ay, aw, ah = _contain_rect(frame.width, frame.height, x, y, w, h)
            frame_resized = frame.resize((aw, ah), Image.LANCZOS)
            result.paste(frame_resized, (ax, ay), frame_resized)
        except Exception as e:
            logger.debug(f"Overlay PIL composite failed: {e}")

    return result


def _clone_slide(dest_prs, src_pptx_path: str, slide_index: int) -> bool:
    """
    Clone slide[slide_index] from src_pptx_path into dest_prs.
    Copies shapes + ALL dependent parts (images, video, GIF, audio, charts).
    Returns True on success, False on failure.
    """
    from pptx import Presentation as Prs
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI
    import lxml.etree as etree

    try:
        src_prs = Prs(src_pptx_path)
    except Exception as e:
        logger.warning(f"Cannot open source PPTX {src_pptx_path}: {e}")
        return False

    if slide_index >= len(src_prs.slides):
        logger.warning(f"Slide index {slide_index} out of range in {src_pptx_path}")
        return False

    src_slide = src_prs.slides[slide_index]
    pkg = dest_prs.part.package  # correct way in pptx 1.0

    # ── 1. Add blank slide ────────────────────────────────────────────────────
    blank_layout = dest_prs.slide_layouts[6]
    dest_slide = dest_prs.slides.add_slide(blank_layout)

    # ── 2. Copy shapes tree verbatim ──────────────────────────────────────────
    src_sp_tree = copy.deepcopy(src_slide.shapes._spTree)
    dest_slide.shapes._spTree.clear()
    for child in src_sp_tree:
        dest_slide.shapes._spTree.append(child)

    # ── 3. Copy background if explicitly set ──────────────────────────────────
    try:
        pml = "http://schemas.openxmlformats.org/presentationml/2006/main"
        src_cSld = src_slide._element.find(f"{{{pml}}}cSld")
        dest_cSld = dest_slide._element.find(f"{{{pml}}}cSld")
        if src_cSld is not None and dest_cSld is not None:
            src_bg = src_cSld.find(f"{{{pml}}}bg")
            if src_bg is not None:
                existing_bg = dest_cSld.find(f"{{{pml}}}bg")
                if existing_bg is not None:
                    dest_cSld.remove(existing_bg)
                dest_cSld.insert(0, copy.deepcopy(src_bg))
    except Exception as e:
        logger.debug(f"Background copy skipped: {e}")

    # ── 4. Copy all dependent parts, build rId mapping ────────────────────────
    rId_map: dict[str, str] = {}

    for old_rId, rel in list(src_slide.part.rels.items()):
        if rel.reltype in _SKIP_RELTYPES:
            continue
        try:
            if rel.is_external:
                new_rId = dest_slide.part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                src_part = rel.target_part
                ext = PurePosixPath(str(src_part.partname)).suffix
                uid = uuid.uuid4().hex[:10]
                new_partname = PackURI(f"/ppt/media/cloned_{uid}{ext}")
                # pptx 1.0: Part(partname, content_type, package, blob)
                new_part = Part(
                    new_partname, src_part.content_type, pkg, src_part.blob
                )
                new_rId = dest_slide.part.relate_to(new_part, rel.reltype)
            rId_map[old_rId] = new_rId
        except Exception as e:
            logger.debug(f"Skipping rel {old_rId} ({rel.reltype}): {e}")

    # ── 5. Remap rId references in the copied XML (in-place) ──────────────────
    if rId_map:
        xml_str = etree.tostring(dest_slide._element, encoding="unicode")
        for old_id, new_id in rId_map.items():
            xml_str = xml_str.replace(f'="{old_id}"', f'="{new_id}"')
        new_el = etree.fromstring(xml_str)
        # _element is the root — replace children and attrs in-place
        dest_slide._element[:] = new_el[:]
        for attr in new_el.attrib:
            dest_slide._element.set(attr, new_el.get(attr))

    return True


def _add_thumbnail_slide(dest_prs, slide_entry: SlideLibraryEntry):
    """Add a full-page thumbnail image as a slide (fallback for PDF-sourced slides)."""
    from pptx.util import Inches

    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)
    thumb_path = Path(settings.thumbnail_dir) / (slide_entry.thumbnail_path or "")
    if thumb_path.exists():
        new_slide.shapes.add_picture(
            io.BytesIO(thumb_path.read_bytes()),
            Inches(0), Inches(0),
            dest_prs.slide_width, dest_prs.slide_height,
        )


def export_to_pptx(db: Session, assembly_id: int) -> str:
    """Export assembly as PPTX.

    Each slide is embedded as a full-page PNG image (from the stored thumbnail)
    so the visual output always matches what the user sees on the site.
    Media overlays are added as picture shapes on top.
    """
    from pptx import Presentation
    from pptx.util import Inches

    assembly = db.query(AssembledPresentation).get(assembly_id)
    if not assembly:
        raise ValueError(f"Assembly {assembly_id} not found")

    slide_ids: list[int] = json.loads(assembly.slide_ids_json or "[]")
    if not slide_ids:
        raise ValueError("Презентация не содержит слайдов")

    overlays_map: dict = json.loads(assembly.overlays_json or "{}")

    slides = [s for sid in slide_ids if (s := db.query(SlideLibraryEntry).get(sid))]

    export_dir = Path(settings.export_dir)
    export_dir.mkdir(parents=True, exist_ok=True)
    export_path = export_dir / f"{assembly_id}_{uuid.uuid4().hex[:8]}.pptx"

    dest_prs = Presentation()
    # Standard widescreen 16:9
    dest_prs.slide_width = Inches(13.33)
    dest_prs.slide_height = Inches(7.5)

    for slide_entry in slides:
        # For PPTX slides with an embedded video, clone the original slide
        # so the video is playable in the exported PPTX (not just a black thumbnail).
        cloned = False
        if slide_entry.video_path and slide_entry.source_id:
            try:
                source = slide_entry.source
                if source and source.file_type == "pptx" and source.file_path:
                    cloned = _clone_slide(dest_prs, source.file_path, slide_entry.slide_index)
            except Exception as e:
                logger.debug(f"Video slide clone failed (slide {slide_entry.id}): {e}")

        if not cloned:
            _add_thumbnail_slide(dest_prs, slide_entry)

        dest_slide = dest_prs.slides[-1]
        _add_overlays_pptx(dest_prs, dest_slide, str(slide_entry.id), overlays_map)

    dest_prs.save(str(export_path))

    assembly.export_path = str(export_path)
    assembly.status = "exported"
    db.commit()

    return str(export_path)


def export_to_pdf(db: Session, assembly_id: int) -> str:
    """Export assembly as PDF. Composites media overlays onto each slide thumbnail."""
    import fitz
    from PIL import Image

    assembly = db.query(AssembledPresentation).get(assembly_id)
    if not assembly:
        raise ValueError(f"Assembly {assembly_id} not found")

    slide_ids: list[int] = json.loads(assembly.slide_ids_json or "[]")
    slides = [s for sid in slide_ids if (s := db.query(SlideLibraryEntry).get(sid))]
    overlays_map: dict = json.loads(assembly.overlays_json or "{}")

    export_dir = Path(settings.export_dir)
    export_dir.mkdir(parents=True, exist_ok=True)
    export_path = export_dir / f"{assembly_id}_{uuid.uuid4().hex[:8]}.pdf"

    doc = fitz.open()
    for slide in slides:
        thumb_path = Path(settings.thumbnail_dir) / (slide.thumbnail_path or "")
        slide_id_str = str(slide.id)
        has_overlays = bool(overlays_map.get(slide_id_str))

        if thumb_path.exists():
            if has_overlays:
                # Composite overlays via PIL then insert into PDF
                base = Image.open(str(thumb_path)).convert("RGBA")
                composited = _composite_overlays_pil(base, slide_id_str, overlays_map)
                buf = io.BytesIO()
                composited.convert("RGB").save(buf, "PNG")
                buf.seek(0)
                img_doc = fitz.open("png", buf.read())
                img_pdf = fitz.open("pdf", img_doc.convert_to_pdf())
                img_doc.close()
            else:
                img_doc = fitz.open(str(thumb_path))
                img_pdf = fitz.open("pdf", img_doc.convert_to_pdf())
                img_doc.close()
            doc.insert_pdf(img_pdf)
            img_pdf.close()
        else:
            doc.new_page(width=1920, height=1080)

    doc.save(str(export_path))
    doc.close()

    return str(export_path)
