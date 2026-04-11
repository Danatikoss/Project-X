"""
Presentation Planner Service.

Flow:
  1. extract_text() — pull readable text from PPTX / PDF / DOCX / plain text
  2. plan_presentation() — AI analyses content → returns ordered list of slide blueprints
  3. render_presentation() — render each blueprint → merge into one PPTX file

The planner prompt enforces visual variety:
  - No text walls (max 5 bullets per slide)
  - Chooses icon_grid, chart_bar, key_message, process_flow instead of bullet lists
  - No two consecutive slides with the same layout
"""

import io
import json
import logging
import os
import tempfile
import uuid
from pathlib import Path

from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy.orm import Session

from config import settings

logger = logging.getLogger(__name__)

# ─── Text extraction ──────────────────────────────────────────────────────────

def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation as _Prs
    prs = _Prs(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Слайд {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def extract_text_from_pdf(file_bytes: bytes) -> str:
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    lines = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if not text:
            # Image-based page — try to extract text from embedded images via OCR blocks
            blocks = page.get_text("blocks")
            text = " ".join(b[4].strip() for b in blocks if b[4].strip())
        if text:
            lines.append(f"\n--- Страница {i} ---")
            lines.append(text)

    result = "\n".join(lines).strip()
    if not result:
        raise ValueError(
            "PDF не содержит извлекаемого текста. "
            "Это отсканированный документ — загрузите DOCX, TXT или PPTX, "
            "либо введите текст вручную."
        )


def extract_text_from_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def extract_text(file_bytes: bytes, file_ext: str) -> str:
    ext = file_ext.lower().lstrip(".")
    if ext == "pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    else:
        return file_bytes.decode("utf-8", errors="replace")


# ─── Planning prompt ──────────────────────────────────────────────────────────

_PLAN_SYSTEM = """You are a world-class presentation designer. Your PRIMARY goal is CONTENT DENSITY — pack maximum meaning into minimum slides.

══════════════════════════════════════════════════════
CONTENT DENSITY RULES (highest priority, no exceptions)
══════════════════════════════════════════════════════
1. ONE SLIDE = ONE COMPLETE LOGICAL SECTION. If ideas are semantically related, combine them on one slide using a rich layout.
2. NEVER create a new slide for a single isolated thought. Ask yourself: "Can this merge with the previous or next slide?"
3. AGGREGATION RULE: If you find yourself creating 3+ consecutive slides with short content (1-2 bullets each), STOP and merge them into one icon_grid (up to 4 cards) or two_column slide.
4. The slide count limit given in the user message is a HARD MAXIMUM. Staying under it is mandatory.

══════════════════════════════════════════
LAYOUT COST — choose the cheapest that fits
══════════════════════════════════════════
EXPENSIVE (use sparingly — each costs a full slide):
  ⚠ section_divider — only for major topic transitions; max 1 per 5 slides
  ⚠ title_content   — absolute last resort; only if no visual layout fits; max 1 total

DENSE (preferred — pack more per slide):
  ✓ icon_grid    — 3-4 concepts/features/benefits on one slide → USE INSTEAD OF BULLETS
  ✓ two_column   — 2 related topics with bullet lists on one slide
  ✓ process_flow — 3-5 sequential steps on one slide
  ✓ comparison   — explicit A vs B on one slide
  ✓ chart_bar / chart_pie — data with 3+ values on one slide

CONTEXTUAL (use when content clearly matches):
  ○ key_message — one powerful punchy statement (≤15 words)
  ○ big_stat    — one dramatic number + context
  ○ timeline    — chronological milestones (max 5)
  ○ quote       — a direct quote or testimonial

DECISION TREE (apply in order):
  • 3-4 concepts/items? → icon_grid
  • 2 topics to compare? → two_column or comparison
  • Steps in a process? → process_flow
  • Numbers/metrics ≥3? → chart_bar or chart_pie
  • One big number? → big_stat
  • Powerful single statement? → key_message
  • Quote? → quote
  • Major section break needed? → section_divider (use rarely)
  • Nothing fits? → title_content (≤5 bullets, last resort)

══════════════════════════════
FIELD LIMITS (hard constraints)
══════════════════════════════
- Title: max 60 chars
- icon_grid cards: heading ≤30 chars, text ≤80 chars; always 3-4 cards
- process_flow steps: label ≤25 chars, desc ≤60 chars
- key_message.message: ≤15 words, impactful
- two_column items: ≤7 per column, ≤100 chars each
- No emoji in any field

STRUCTURE:
- First slide: key_message or section_divider (intro)
- Last slide: key_message or big_stat (closing call to action)
- No two consecutive slides with the same layout

Available layouts:
icon_grid | key_message | process_flow | chart_bar | chart_pie | big_stat | two_column | comparison | timeline | quote | section_divider | title_content

JSON schema:
icon_grid:       {"layout":"icon_grid","title":"...","content":{"cards":[{"heading":"...","text":"..."}]},"speaker_notes":"..."}
key_message:     {"layout":"key_message","title":"...","content":{"message":"...","subtext":"..."},"speaker_notes":"..."}
process_flow:    {"layout":"process_flow","title":"...","content":{"steps":[{"label":"...","desc":"..."}]},"speaker_notes":"..."}
chart_bar:       {"layout":"chart_bar","title":"...","content":{"categories":["A","B"],"series":[{"name":"Metric","values":[10,20]}]},"speaker_notes":"..."}
chart_pie:       {"layout":"chart_pie","title":"...","content":{"slices":[{"label":"A","value":60},{"label":"B","value":40}]},"speaker_notes":"..."}
big_stat:        {"layout":"big_stat","title":"...","content":{"value":"...","label":"...","context":["...","..."]},"speaker_notes":"..."}
two_column:      {"layout":"two_column","title":"...","content":{"left":{"heading":"...","items":["..."]},"right":{"heading":"...","items":["..."]}},"speaker_notes":"..."}
comparison:      {"layout":"comparison","title":"...","content":{"left":{"label":"...","items":["..."]},"right":{"label":"...","items":["..."]}},"speaker_notes":"..."}
timeline:        {"layout":"timeline","title":"...","content":{"steps":[{"label":"...","event":"..."}]},"speaker_notes":"..."}
quote:           {"layout":"quote","title":"...","content":{"quote":"...","attribution":"..."},"speaker_notes":"..."}
section_divider: {"layout":"section_divider","title":"...","content":{"subtitle":"..."},"speaker_notes":"..."}
title_content:   {"layout":"title_content","title":"...","content":{"type":"bullets","items":["..."]},"speaker_notes":"..."}

Respond with ONLY a JSON array. No markdown, no wrapper object."""


def _get_client() -> AsyncOpenAI:
    kwargs: dict = {"api_key": settings.openai_api_key}
    if settings.openai_base_url:
        kwargs["base_url"] = settings.openai_base_url
    return AsyncOpenAI(**kwargs)


# ─── Slide count budget ────────────────────────────────────────────────────────

def _compute_target_slide_count(text: str) -> tuple[int, int]:
    """
    Return (min_slides, max_slides) based on raw character count.

    Calibration:
      < 1 000 chars  — very short brief   → 3-4 slides
      1 000-3 000    — short document     → 4-6 slides
      3 000-6 000    — medium document    → 6-9 slides
      6 000-12 000   — long document      → 9-13 slides
      > 12 000       — very long          → 12-16 slides
    """
    n = len(text.strip())
    if n < 1_000:
        return 3, 4
    if n < 3_000:
        return 4, 6
    if n < 6_000:
        return 6, 9
    if n < 12_000:
        return 9, 13
    return 12, 16


# ─── Post-LLM aggregation ─────────────────────────────────────────────────────

def _is_short_slide(bp: dict) -> bool:
    """
    True when a blueprint carries so little content it could be a card
    inside an icon_grid on a larger slide.
    """
    layout = bp.get("layout", "")
    c = bp.get("content", {})

    if layout == "title_content":
        items = c.get("items") or []
        # Short = 1-2 brief bullets
        return len(items) <= 2 and all(len(i) <= 60 for i in items)

    if layout == "key_message":
        words = (c.get("message") or "").split()
        return len(words) <= 8 and not c.get("subtext")

    if layout == "section_divider":
        # Section dividers at non-start/end positions waste a slide
        return True

    return False


def _bp_to_card(bp: dict) -> dict:
    """
    Extract a heading + text pair from a short blueprint for use as an icon_grid card.
    """
    layout = bp.get("layout", "")
    title = (bp.get("title") or "")[:30]
    c = bp.get("content", {})

    if layout == "title_content":
        items = c.get("items") or []
        text = "; ".join(items[:2])[:80]
    elif layout == "key_message":
        text = (c.get("message") or "")[:80]
    elif layout == "section_divider":
        text = (c.get("subtitle") or "")[:80]
    else:
        text = title[:80]
        title = (bp.get("title") or "")[:30]

    return {"heading": title, "text": text}


def _aggregate_short_slides(blueprints: list[dict]) -> list[dict]:
    """
    Scan for runs of 3+ consecutive short slides and collapse them into
    a single icon_grid (up to 4 cards).  Runs at the very start or very
    end of the deck are left intact (intro/closing slides are expected there).

    Returns a new list; does not mutate the input.
    """
    if len(blueprints) < 3:
        return blueprints

    result: list[dict] = []
    i = 0

    while i < len(blueprints):
        # Protect first and last slide from merging
        if i == 0 or i == len(blueprints) - 1:
            result.append(blueprints[i])
            i += 1
            continue

        # Detect run of short slides (excluding first/last)
        run_end = i
        while (
            run_end < len(blueprints) - 1   # don't consume the last slide
            and _is_short_slide(blueprints[run_end])
        ):
            run_end += 1

        run_len = run_end - i
        if run_len >= 3:
            # Merge up to 4 cards into one icon_grid
            cards_bps = blueprints[i : i + min(run_len, 4)]
            cards = [_bp_to_card(bp) for bp in cards_bps]
            merged_title = blueprints[i].get("title", "Обзор")[:60]
            merged = {
                "layout": "icon_grid",
                "title":  merged_title,
                "content": {"cards": cards},
                "speaker_notes": "Объединённый слайд.",
            }
            result.append(merged)
            logger.debug(
                f"Aggregated {len(cards_bps)} short slides → icon_grid: {[b.get('title') for b in cards_bps]}"
            )
            i += len(cards_bps)
            # If there are more short slides in the same run (> 4), continue loop
        else:
            result.append(blueprints[i])
            i += 1

    return result


# ─── Post-LLM hard count enforcement ─────────────────────────────────────────

def _enforce_slide_count(blueprints: list[dict], max_slides: int) -> list[dict]:
    """
    If the model still returned more slides than allowed, aggressively
    drop the least dense slides from the middle of the deck.
    Always preserves the first and last slide.
    """
    if len(blueprints) <= max_slides:
        return blueprints

    logger.warning(
        f"Model returned {len(blueprints)} slides, max={max_slides} — trimming"
    )

    # Score each middle slide by content density (higher = keep)
    def _density(bp: dict) -> int:
        layout = bp.get("layout", "")
        c = bp.get("content", {})
        # Expensive/thin layouts score low
        if layout == "section_divider":
            return 0
        if layout == "title_content":
            return len(c.get("items") or [])
        if layout == "key_message":
            return len((c.get("message") or "").split())
        if layout == "icon_grid":
            return len(c.get("cards") or []) * 10  # high value
        if layout == "two_column":
            l_items = len((c.get("left") or {}).get("items") or [])
            r_items = len((c.get("right") or {}).get("items") or [])
            return (l_items + r_items) * 8
        if layout in ("chart_bar", "chart_pie", "big_stat"):
            return 15  # data slides always dense
        if layout == "process_flow":
            return len(c.get("steps") or []) * 5
        return 5

    first, *middle, last = blueprints
    middle_scored = sorted(enumerate(middle), key=lambda x: _density(x[1]), reverse=True)
    keep_count = max_slides - 2  # -2 for first and last
    kept_indices = {idx for idx, _ in middle_scored[:keep_count]}
    kept_middle = [bp for idx, bp in enumerate(middle) if idx in kept_indices]
    # Restore original order
    kept_middle.sort(key=lambda bp: middle.index(bp))

    return [first] + kept_middle + [last]


def _build_brand_context_section(brand_context: dict | None) -> str:
    """
    Build a brand context appendix for the planning system prompt.
    Returns empty string when no brand context is provided.
    """
    if not brand_context:
        return ""
    lines: list[str] = []
    if brand_context.get("tone_of_voice"):
        lines.append(f"Tone of voice: {brand_context['tone_of_voice']}")
    if brand_context.get("target_audience"):
        lines.append(f"Target audience: {brand_context['target_audience']}")
    if brand_context.get("prohibitions"):
        lines.append(
            f"PROHIBITIONS — never include in any slide text: {brand_context['prohibitions']}"
        )
    if brand_context.get("brand_guidelines_text"):
        lines.append(f"Brand guidelines: {brand_context['brand_guidelines_text']}")
    if not lines:
        return ""
    return (
        "\n\n══════════════════════\n"
        "CLIENT BRAND CONTEXT (mandatory — apply to every slide)\n"
        "══════════════════════\n"
        + "\n".join(lines)
    )


def _build_plan_system_prompt(
    available_layouts: set[str] | None,
    brand_context: dict | None = None,
) -> str:
    """
    Build the system prompt for plan_presentation(), filtering to only
    the layouts actually available in the user's template and injecting
    brand context when provided.
    """
    if not available_layouts:
        return _PLAN_SYSTEM + _build_brand_context_section(brand_context)

    all_layouts = {
        "icon_grid", "key_message", "process_flow", "chart_bar", "chart_pie",
        "big_stat", "two_column", "comparison", "timeline", "quote",
        "section_divider", "title_content",
    }
    removed = all_layouts - available_layouts
    if not removed:
        return _PLAN_SYSTEM + _build_brand_context_section(brand_context)

    prompt = _PLAN_SYSTEM

    # Update the Available layouts line
    layout_list = " | ".join(sorted(available_layouts))
    import re
    prompt = re.sub(
        r"Available layouts:\n.*",
        f"Available layouts:\n{layout_list}",
        prompt,
    )

    # Remove JSON schema lines for unavailable layouts
    for layout in removed:
        prompt = re.sub(rf"^{layout}:.*\n?", "", prompt, flags=re.MULTILINE)

    # Remove DECISION TREE bullets for unavailable layouts
    removal_hints = {
        "timeline":        "chronological milestones",
        "quote":           "direct quote",
        "big_stat":        "one big number",
        "comparison":      "explicit A vs B",
        "section_divider": "section break",
    }
    for layout in removed:
        if layout in removal_hints:
            prompt = re.sub(
                rf"  • .*{removal_hints[layout]}.*→ {layout}.*\n?",
                "",
                prompt,
            )

    return prompt + _build_brand_context_section(brand_context)


async def plan_presentation(
    content_text: str,
    title: str = "",
    language_hint: str = "",
    template_pptx_path: str | None = None,
    brand_context: dict | None = None,
) -> list[dict]:
    """
    Ask AI to produce a full slide plan from the provided text content.
    Returns a list of blueprint dicts ready for render_slide_pptx().
    """
    from services.slide_generator import get_available_layouts

    client = _get_client()

    min_slides, max_slides = _compute_target_slide_count(content_text)

    available_layouts = get_available_layouts(template_pptx_path)
    system_prompt = _build_plan_system_prompt(available_layouts, brand_context)
    logger.info(f"Available layouts for LLM: {sorted(available_layouts)}")
    if brand_context:
        logger.info(
            f"Brand context injected: tone={brand_context.get('tone_of_voice')!r}, "
            f"audience={brand_context.get('target_audience')!r}"
        )

    lang_note = f"\nIMPORTANT: Write all slide text in {language_hint}." if language_hint else ""
    user_msg = (
        f"Presentation title: {title}\n\n"
        f"SLIDE COUNT: Generate STRICTLY between {min_slides} and {max_slides} slides. "
        f"This is a hard limit — do not exceed {max_slides} slides under any circumstances. "
        f"Merge related ideas onto one slide using icon_grid or two_column rather than splitting.\n\n"
        f"Content to convert into slides:\n\n{content_text[:12000]}"
        f"{lang_note}"
    )

    resp = await client.chat.completions.create(
        model=settings.generator_model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_msg},
        ],
        temperature=0.4,   # lower = more disciplined, fewer hallucinated slides
        max_tokens=6000,
    )

    raw = (resp.choices[0].message.content or "[]").strip()

    # Strip markdown fences if present
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.rsplit("```", 1)[0]

    blueprints: list[dict] = json.loads(raw.strip())

    logger.info(f"LLM returned {len(blueprints)} blueprints (budget: {min_slides}-{max_slides})")

    # Post-process 1: merge 3+ consecutive short slides → icon_grid
    blueprints = _aggregate_short_slides(blueprints)
    logger.info(f"After aggregation: {len(blueprints)} blueprints")

    # Post-process 2: hard count enforcement if model still exceeded
    blueprints = _enforce_slide_count(blueprints, max_slides)
    logger.info(f"After count enforcement: {len(blueprints)} blueprints")

    # Reviewer pass — fix consecutive layouts, tone/prohibitions, narrative flow
    from services.reviewer_agent import review_and_fix_blueprints
    blueprints = await review_and_fix_blueprints(
        blueprints=blueprints,
        available_layouts=available_layouts,
        brand_context=brand_context,
        title=title,
    )
    logger.info(f"After reviewer pass: {len(blueprints)} blueprints")

    # Sanitise and trim each blueprint — enforce layout limits before any rendering
    from services.blueprint_validator import validate_and_trim
    for bp in blueprints:
        bp.setdefault("layout", "title_content")
        bp.setdefault("title", "")
        bp.setdefault("content", {})
        bp.setdefault("speaker_notes", "")
        validate_and_trim(bp)

    return blueprints


# ─── Assembly creation from plan ──────────────────────────────────────────────

async def create_assembly_from_plan(
    blueprints: list[dict],
    title: str,
    brand_template_id: int | None,
    user_id: int,
    db: Session,
) -> int:
    """
    Render each blueprint → SlideLibraryEntry (with thumbnail) → AssembledPresentation.
    Returns the assembly ID.
    """
    from services.slide_generator import (
        BrandColors, _extract_brand_colors, save_slide_from_blueprint,
    )
    from models.brand import BrandTemplate
    from models.assembly import AssembledPresentation

    # Load brand colors once for all slides
    colors: BrandColors = BrandColors()
    template_pptx_path: str | None = None

    # Auto-apply default template when none explicitly selected
    if not brand_template_id:
        default_tmpl = db.query(BrandTemplate).filter(
            BrandTemplate.is_default == True,
            BrandTemplate.owner_id == user_id,
        ).first()
        if default_tmpl:
            brand_template_id = default_tmpl.id
            logger.info(f"Auto-applying default brand template id={brand_template_id}")

    if brand_template_id:
        tmpl = db.query(BrandTemplate).filter(BrandTemplate.id == brand_template_id).first()
        if tmpl:
            if tmpl.pptx_path and os.path.exists(tmpl.pptx_path):
                template_pptx_path = tmpl.pptx_path
            stored = json.loads(tmpl.colors_json or "{}")
            if stored:
                colors = BrandColors(**{k: v for k, v in stored.items()
                                        if k in BrandColors.__dataclass_fields__})
            elif template_pptx_path:
                colors = _extract_brand_colors(template_pptx_path)

            # Apply strict brand guidelines — same as generate_slide() in slide_generator.py.
            # Previously this block was missing, so background_image_path, font settings,
            # shape_color, and text-zone positions were never applied for plan-based generation.
            if tmpl.font_family:
                colors.font_family = tmpl.font_family
            if tmpl.title_font_color:
                colors.title_font_color = tmpl.title_font_color
            if tmpl.title_font_size:
                colors.title_font_size = tmpl.title_font_size
            if tmpl.body_font_color:
                colors.body_font_color = tmpl.body_font_color
            if tmpl.body_font_size:
                colors.body_font_size = tmpl.body_font_size
            if tmpl.shape_color:
                colors.shape_color = tmpl.shape_color
                colors.primary = tmpl.shape_color
            if tmpl.shape_opacity is not None:
                colors.shape_opacity = tmpl.shape_opacity
            if tmpl.background_image_path and os.path.exists(tmpl.background_image_path):
                colors.background_image_path = tmpl.background_image_path
            for field in ("title_x", "title_y", "title_w", "title_h",
                          "body_x",  "body_y",  "body_w",  "body_h"):
                val = getattr(tmpl, field, None)
                if val is not None:
                    setattr(colors, field, val)

    # Fixed brand overrides from env — always win over template settings
    if settings.fixed_bg_image and os.path.exists(settings.fixed_bg_image):
        colors.background_image_path = settings.fixed_bg_image
    if settings.fixed_shape_color:
        colors.shape_color = settings.fixed_shape_color
        colors.primary     = settings.fixed_shape_color
    if settings.fixed_title_font_size > 0:
        colors.title_font_size = settings.fixed_title_font_size
    if settings.fixed_body_font_size > 0:
        colors.body_font_size = settings.fixed_body_font_size

    # Render each slide and save to library
    slide_ids: list[int] = []
    for i, bp in enumerate(blueprints):
        entry = await save_slide_from_blueprint(
            db=db,
            blueprint=bp,
            colors=colors,
            template_pptx_path=template_pptx_path,
            user_id=user_id,
            slide_index=i,
        )
        slide_ids.append(entry.id)

    # Create assembly
    assembly = AssembledPresentation(
        owner_id=user_id,
        title=title,
        prompt="(AI-генерация)",
        slide_ids_json=json.dumps(slide_ids),
        status="draft",
        brand_template_id=brand_template_id,
    )
    db.add(assembly)
    db.commit()
    db.refresh(assembly)
    return assembly.id, slide_ids
