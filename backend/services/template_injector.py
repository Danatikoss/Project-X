"""
Template Injector — copies a slide from Libraryslides.pptx and injects
content into named slots (shapes whose name starts with 'slot_').

Key design:
- Preserves ALL visual formatting AND media (background images, pictures)
- Remaps r:embed / r:link rIds so media relationships are valid
- Replaces only text content of each slot, keeping font/size/color
- Multi-paragraph text is split on \\n
"""
import copy
import io
import json
import logging
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn, nsmap
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from services.template_library import TemplateInfo, TEMPLATES_DIR
from services.icon_service import icon_to_png

_CHART_TYPE_MAP = {
    "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":   XL_CHART_TYPE.LINE,
    "pie":    XL_CHART_TYPE.PIE,
}

PPTX_PATH = TEMPLATES_DIR / "Libraryslides.pptx"

logger = logging.getLogger(__name__)

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ── Media copying ─────────────────────────────────────────────────────────────

def _copy_media_and_remap(source_slide_part, new_slide_part, xml_elem) -> etree._Element:
    """
    Deep-copy xml_elem and remap all r:embed / r:link attributes.
    For each rId found in source_slide_part.rels, the media part is added
    to new_slide_part and a fresh rId is assigned.
    Returns the remapped copy.
    """
    elem_copy = copy.deepcopy(xml_elem)

    # Collect all rId values used in the element
    r_embed_tag = f"{{{R_NS}}}embed"
    r_link_tag  = f"{{{R_NS}}}link"

    rId_map: dict[str, str] = {}

    for node in elem_copy.iter():
        for attr in (r_embed_tag, r_link_tag):
            old_rid = node.get(attr)
            if old_rid and old_rid not in rId_map:
                # Look up the part in the source slide
                if old_rid in source_slide_part.rels:
                    rel = source_slide_part.rels[old_rid]
                    if rel.is_external:
                        new_rid = new_slide_part.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                    else:
                        target_part = rel.target_part
                        new_rid = new_slide_part.part.relate_to(target_part, rel.reltype)
                    rId_map[old_rid] = new_rid
                    logger.debug("Remapped rId %s → %s (%s)", old_rid, new_rid, rel.reltype.split("/")[-1])

    # Apply remap
    for node in elem_copy.iter():
        for attr in (r_embed_tag, r_link_tag):
            old_rid = node.get(attr)
            if old_rid and old_rid in rId_map:
                node.set(attr, rId_map[old_rid])

    return elem_copy


# ── Text injection ────────────────────────────────────────────────────────────

def _copy_run_format(source_run_elem):
    rPr = source_run_elem.find(qn("a:rPr"))
    return copy.deepcopy(rPr) if rPr is not None else None


def _make_run(text: str, rPr=None) -> etree._Element:
    r = etree.Element(qn("a:r"))
    if rPr is not None:
        r.append(copy.deepcopy(rPr))
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return r


def _set_shape_text(shape, new_text: str):
    """Replace text in a shape's text frame while preserving per-paragraph formatting."""
    if not shape.has_text_frame:
        return

    txBody = shape.text_frame._txBody

    # Collect formatting from each original paragraph: (pPr, rPr)
    para_formats: list[tuple] = []
    for p in txBody.findall(qn("a:p")):
        pPr_elem = p.find(qn("a:pPr"))
        pPr = copy.deepcopy(pPr_elem) if pPr_elem is not None else None

        # Use the first run's rPr in this paragraph
        first_run = p.find(qn("a:r"))
        rPr = _copy_run_format(first_run) if first_run is not None else None

        para_formats.append((pPr, rPr))

    # Fallback if no paragraphs found
    if not para_formats:
        para_formats = [(None, None)]

    lines = new_text.split("\n")

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    for i, line in enumerate(lines):
        # Use corresponding format, or reuse last one for extra lines
        pPr, rPr = para_formats[min(i, len(para_formats) - 1)]
        p = etree.SubElement(txBody, qn("a:p"))
        if pPr is not None:
            p.append(copy.deepcopy(pPr))
        if line:
            p.append(_make_run(line, rPr))

    logger.debug("Set text on shape %r: %r", shape.name, new_text[:60])


# ── Icon injection ────────────────────────────────────────────────────────────

def _inject_icon(slide, shape, icon_name: str, color: str = "#FFFFFF"):
    """
    Replace a placeholder shape with a Lucide icon PNG at the same position/size.
    Preserves the z-order by inserting the picture element where the placeholder was.
    """
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    png_bytes = icon_to_png(icon_name, color=color, size=512)
    if png_bytes is None:
        logger.warning("Icon %r not found — skipping slot %r", icon_name, shape.name)
        return

    sp_tree = slide.shapes._spTree
    sp_elem = shape._element

    # Record z-position before removing placeholder
    insert_idx = list(sp_tree).index(sp_elem) if sp_elem in sp_tree else None
    if insert_idx is not None:
        sp_tree.remove(sp_elem)

    # Add picture (appended to end of sp_tree by python-pptx)
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)

    # Move the new picture element to the original z-position
    if insert_idx is not None:
        pic_elem = sp_tree[-1]
        sp_tree.remove(pic_elem)
        sp_tree.insert(insert_idx, pic_elem)

    logger.debug("Injected icon %r into slide at (%d, %d)", icon_name, left, top)


# ── Chart injection ───────────────────────────────────────────────────────────

def _inject_chart(slide, shape, chart_data_raw: str):
    """
    Replace a placeholder shape with a native PPTX chart.
    chart_data_raw: JSON string produced by _fill_chart_slot.
    """
    try:
        data = json.loads(chart_data_raw)
    except (json.JSONDecodeError, TypeError):
        logger.warning("Invalid chart JSON for slot %r — skipping", shape.name)
        return

    chart_type = _CHART_TYPE_MAP.get(str(data.get("type", "bar")).lower(), XL_CHART_TYPE.BAR_CLUSTERED)
    categories  = data.get("categories", [])
    series_list = data.get("series", [])
    title       = data.get("title", "")

    if not categories or not series_list:
        logger.warning("Chart missing categories/series for slot %r — skipping", shape.name)
        return

    # Validate all values are numeric; drop non-numeric entries gracefully
    clean_series = []
    for s in series_list:
        try:
            values = [float(v) for v in s.get("values", [])]
        except (TypeError, ValueError):
            logger.warning("Non-numeric values in chart series %r — skipping series", s.get("name"))
            continue
        if len(values) != len(categories):
            logger.warning("Series %r length %d ≠ categories %d — padding/truncating", s.get("name"), len(values), len(categories))
            values = (values + [0.0] * len(categories))[:len(categories)]
        clean_series.append({"name": s.get("name", ""), "values": values})

    if not clean_series:
        logger.warning("No valid series for chart slot %r — skipping", shape.name)
        return

    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    cd = ChartData()
    cd.categories = categories
    for s in clean_series:
        cd.add_series(s["name"], tuple(s["values"]))

    # Remove placeholder, record z-position, then reinsert chart at same index
    sp_tree = slide.shapes._spTree
    sp_elem  = shape._element
    sp_list  = list(sp_tree)
    insert_idx = sp_list.index(sp_elem) if sp_elem in sp_list else None
    if insert_idx is not None:
        sp_tree.remove(sp_elem)

    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, cd)
    chart = chart_frame.chart

    if insert_idx is not None:
        chart_elem = sp_tree[-1]
        sp_tree.remove(chart_elem)
        sp_tree.insert(insert_idx, chart_elem)

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False

    chart.has_legend = len(clean_series) > 1

    logger.debug("Injected %r chart (%d cats, %d series) into slot %r",
                 data.get("type"), len(categories), len(clean_series), shape.name)


# ── Core copy logic ───────────────────────────────────────────────────────────

def _copy_slide_into(source_prs: Presentation, slide_index: int,
                     out_prs: Presentation, slots: dict[str, str]):
    """
    Copy slide[slide_index] from source_prs into out_prs,
    transferring all media parts and remapping rIds.
    Then inject slot text.
    Returns the new slide.
    """
    source_slide = source_prs.slides[slide_index]

    blank_layout = out_prs.slide_layouts[6]
    new_slide = out_prs.slides.add_slide(blank_layout)

    source_part = source_slide.part
    new_part    = new_slide.part

    sp_tree        = new_slide.shapes._spTree
    source_sp_tree = source_slide.shapes._spTree

    # ── Copy background (with media remap) ──────────────────────────────────
    source_cSld = source_slide._element.find(qn("p:cSld"))
    new_cSld    = new_slide._element.find(qn("p:cSld"))

    if source_cSld is not None and new_cSld is not None:
        source_bg = source_cSld.find(qn("p:bg"))
        if source_bg is not None:
            existing_bg = new_cSld.find(qn("p:bg"))
            if existing_bg is not None:
                new_cSld.remove(existing_bg)
            remapped_bg = _copy_media_and_remap(source_part, new_slide, source_bg)
            new_cSld.insert(0, remapped_bg)

    # ── Copy shapes (with media remap) ───────────────────────────────────────
    for elem in source_sp_tree:
        tag = etree.QName(elem.tag).localname
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            remapped = _copy_media_and_remap(source_part, new_slide, elem)
            sp_tree.append(remapped)

    # ── Inject slot text ─────────────────────────────────────────────────────
    for shape in list(new_slide.shapes):
        if shape.name not in slots:
            continue
        if shape.name.startswith("slot_chart_"):
            _inject_chart(new_slide, shape, slots[shape.name])
        elif shape.name.startswith("slot_icon_"):
            _inject_icon(new_slide, shape, icon_name=slots[shape.name])
        elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
            _set_shape_text(shape, slots[shape.name])

    return new_slide


# ── Public API ────────────────────────────────────────────────────────────────

def inject_into_slide(template: TemplateInfo, slots: dict[str, str]) -> Presentation:
    """
    Copy the template slide, inject slots, return a single-slide Presentation.
    """
    source_prs = Presentation(str(template.pptx_path))

    out_prs = Presentation()
    out_prs.slide_width  = source_prs.slide_width
    out_prs.slide_height = source_prs.slide_height

    # Remove default blank slide
    sldIdLst = out_prs.slides._sldIdLst
    for sid in list(sldIdLst):
        sldIdLst.remove(sid)

    _copy_slide_into(source_prs, template.slide_index, out_prs, slots)
    return out_prs


def inject_into_presentation(
    base_prs: Presentation,
    template: TemplateInfo,
    slots: dict[str, str],
    source_cache: dict | None = None,
):
    """
    Append an injected template slide to an existing Presentation in-place.

    source_cache: optional dict {pptx_path_str → Presentation} to avoid
    reloading the same PPTX multiple times (prevents duplicate media entries).
    """
    path_str = str(template.pptx_path)
    if source_cache is not None:
        if path_str not in source_cache:
            source_cache[path_str] = Presentation(path_str)
        source_prs = source_cache[path_str]
    else:
        source_prs = Presentation(path_str)
    return _copy_slide_into(source_prs, template.slide_index, base_prs, slots)
