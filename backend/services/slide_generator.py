"""
Slide generation service — Template-First Architecture.

Flow:
  1. Load brand template PPTX → build BrandColors (for thumbnail only)
  2. Call LLM to generate structured blueprint JSON
  3. validate_and_trim(blueprint) — enforce layout-specific char/count limits
  4. render_slide_pptx(blueprint, template) — Template-First render:
       • Find matching named layout in template master
       • prs.slides.add_slide(layout) → inherits ALL master design
       • Fill placeholders with blueprint text only (no coordinate math)
  5. Save PPTX → generate PNG thumbnail
  6. vision_validate(thumbnail, blueprint, colors) — GPT-4o Vision QA check
       • If issues found → re-render with fixed blueprint (one retry)
  7. Create SlideLibraryEntry in DB with embedding + blueprint_json
  8. Return entry
"""

import json
import logging
import os
import uuid
from dataclasses import dataclass
from pathlib import Path

import fitz  # PyMuPDF
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PPT
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from config import settings
from models.brand import BrandTemplate
from models.slide import SlideLibraryEntry, SourcePresentation
from schemas.brand_context import BrandContext
from services.blueprint_validator import validate_and_trim
from services.embedding import embed_single

logger = logging.getLogger(__name__)

# ─── Slide canvas (kept for fallback blank presentations) ────────────────────

W = Inches(13.333)
H = Inches(7.5)

# ─── Brand colors (used ONLY for thumbnail generation; not for layout) ────────


@dataclass
class BrandColors:
    primary: str = "1E3A8A"
    secondary: str = "3B82F6"
    background: str = "FFFFFF"
    text: str = "0F172A"
    text_body: str = "1E293B"
    text_muted: str = "64748B"
    accent_light: str = "EFF6FF"
    divider: str = "E2E8F0"
    font_family: str = "Montserrat"
    title_font_color: str = "FFFFFF"
    title_font_size: int = 30
    body_font_color: str = "1E293B"
    body_font_size: int = 18
    shape_color: str = "1E3A8A"
    shape_opacity: int = 100
    background_image_path: str | None = None
    # Text zone fractions (kept for backward compat with vision_validator)
    title_x: float = 0.038
    title_y: float = 0.00
    title_w: float = 0.924
    title_h: float = 0.193
    body_x: float = 0.038
    body_y: float = 0.220
    body_w: float = 0.924
    body_h: float = 0.760

    def _rgb(self, h: str) -> RGBColor:
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @property
    def primary_rgb(self): return self._rgb(self.primary)
    @property
    def secondary_rgb(self): return self._rgb(self.secondary)
    @property
    def bg_rgb(self): return self._rgb(self.background)
    @property
    def text_rgb(self): return self._rgb(self.text)
    @property
    def text_body_rgb(self): return self._rgb(self.body_font_color)
    @property
    def text_muted_rgb(self): return self._rgb(self.text_muted)
    @property
    def accent_light_rgb(self): return self._rgb(self.accent_light)
    @property
    def divider_rgb(self): return self._rgb(self.divider)
    @property
    def shape_rgb(self): return self._rgb(self.shape_color)
    @property
    def title_color_rgb(self): return self._rgb(self.title_font_color)

    def content_area(self, header_h_in: float = 1.45):
        if self.background_image_path:
            return (W * self.body_x, H * self.body_y, W * self.body_w, H * self.body_h)
        content_top = Inches(header_h_in + 0.15)
        return (Inches(0.55), content_top, W - Inches(0.75), H - content_top - Inches(0.2))


def _extract_brand_colors(pptx_path: str) -> BrandColors:
    """Pull accent1/dk1/lt1 from the PPTX slide master theme."""
    try:
        prs = Presentation(pptx_path)
        NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        clr_scheme = prs.slide_masters[0].element.find(f".//{{{NS}}}clrScheme")
        if clr_scheme is None:
            return BrandColors()

        def _get(name: str) -> str | None:
            el = clr_scheme.find(f"{{{NS}}}{name}")
            if el is None:
                return None
            for child in el:
                v = child.get("val") or child.get("lastClr")
                if v and len(v) == 6:
                    return v.upper()
            return None

        return BrandColors(
            primary=_get("accent1") or "1E3A8A",
            secondary=_get("accent2") or "3B82F6",
            background=_get("lt1") or "FFFFFF",
            text=_get("dk1") or "0F172A",
            text_body=_get("dk1") or "0F172A",
        )
    except Exception as e:
        logger.warning(f"Color extraction failed: {e}")
        return BrandColors()


# ─── Brand styling ───────────────────────────────────────────────────────────


def apply_brand_styling(slide, brand: BrandContext) -> None:
    """
    Apply brand colors and font family to a rendered slide's placeholders.

    Rules:
      - Title placeholder (TITLE / CENTER_TITLE): font_family + title_color
        (falls back to primary_color when title_color is None)
      - Body placeholder (BODY): font_family + body_color
      - Non-placeholder solid-fill shapes: accent_color on fill
      - Any field that is None on the BrandContext is skipped entirely —
        template defaults are never overridden with None.
    """
    def _rgb(hex_str: str | None) -> RGBColor | None:
        if not hex_str:
            return None
        h = hex_str.lstrip("#")
        if len(h) != 6:
            return None
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _style_text(shape, font_family: str | None, color_hex: str | None) -> None:
        if not shape.has_text_frame:
            return
        rgb = _rgb(color_hex)
        for para in shape.text_frame.paragraphs:
            # Paragraph-level default: covers any future runs added without explicit formatting
            if font_family:
                para.font.name = font_family
            if rgb is not None:
                para.font.color.rgb = rgb
            # Run-level: theme colors defined in the layout/master inherit at run
            # level and beat the paragraph default, so we must also set each run.
            for run in para.runs:
                if font_family:
                    run.font.name = font_family
                if rgb is not None:
                    run.font.color.rgb = rgb

    accent_rgb = _rgb(brand.accent_color)

    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (PPT.TITLE, PPT.CENTER_TITLE):
                _style_text(shape, brand.font_family, brand.title_color or brand.primary_color)
            elif ph_type == PPT.BODY:
                _style_text(shape, brand.font_family, brand.body_color)
        else:
            if accent_rgb is not None:
                try:
                    if shape.fill.type == 1:  # MSO_FILL_TYPE.SOLID
                        shape.fill.fore_color.rgb = accent_rgb
                except Exception:
                    pass


# ─── Template-First: Layout mapping ──────────────────────────────────────────
#
# Maps our 12 layout names to preferred template layout names (case-insensitive).
# First match in the list wins; if none match we fall back to any TITLE+BODY layout.

_LAYOUT_PREFS: dict[str, list[str]] = {
    # Scratch-rendered layouts — any template supports them (no placeholder needed)
    "metrics_grid":    ["TITLE_AND_BODY", "ONE_COLUMN_TEXT", "TITLE_ONLY", "TITLE"],
    "title_content":   ["TITLE_AND_BODY", "ONE_COLUMN_TEXT", "TITLE_AND_TWO_COLUMNS"],
    "two_column":      ["TITLE_AND_TWO_COLUMNS", "COMPARISON", "TITLE_AND_BODY", "ONE_COLUMN_TEXT"],
    "comparison":      ["TITLE_AND_TWO_COLUMNS", "COMPARISON", "TITLE_AND_BODY", "ONE_COLUMN_TEXT"],
    "icon_grid":       ["TITLE_AND_BODY", "ONE_COLUMN_TEXT", "TITLE_AND_TWO_COLUMNS"],
    "process_flow":    ["TITLE_AND_BODY", "ONE_COLUMN_TEXT"],
    "timeline":        ["TITLE_AND_BODY", "ONE_COLUMN_TEXT"],
    "key_message":     ["MAIN_POINT", "SECTION_HEADER", "TITLE_ONLY", "TITLE_AND_BODY"],
    "big_stat":        ["BIG_NUMBER", "TITLE_AND_BODY", "MAIN_POINT"],
    "quote":           ["TITLE_AND_BODY", "ONE_COLUMN_TEXT", "SECTION_TITLE_AND_DESCRIPTION"],
    "section_divider": ["TITLE", "SECTION_HEADER", "SECTION_TITLE_AND_DESCRIPTION", "MAIN_POINT"],
    "chart_bar":       ["TITLE_AND_BODY", "ONE_COLUMN_TEXT"],
    "chart_pie":       ["TITLE_AND_BODY", "ONE_COLUMN_TEXT"],
}


def get_available_layouts(template_pptx_path: str | None) -> set[str]:
    """
    Return the set of our layout names that the given template can serve.
    Called by presentation_planner before the LLM prompt is built so the
    model is restricted to layouts the template can actually render.
    """
    if not template_pptx_path or not os.path.exists(template_pptx_path):
        return set(_LAYOUT_PREFS.keys())
    try:
        prs = Presentation(template_pptx_path)
        tmpl_names = {lay.name.upper() for lay in prs.slide_layouts}
        available: set[str] = set()
        for our_layout, prefs in _LAYOUT_PREFS.items():
            for pref in prefs:
                if pref.upper() in tmpl_names:
                    available.add(our_layout)
                    break
        return available or set(_LAYOUT_PREFS.keys())
    except Exception as e:
        logger.warning(f"get_available_layouts: {e}")
        return set(_LAYOUT_PREFS.keys())


def _find_layout(prs: Presentation, our_layout: str):
    """Return the best matching SlideLayout from the template."""
    prefs = _LAYOUT_PREFS.get(our_layout, ["TITLE_AND_BODY"])
    name_map = {lay.name.upper(): lay for lay in prs.slide_layouts}

    for pref in prefs:
        if pref.upper() in name_map:
            return name_map[pref.upper()]

    # Fallback: any layout with TITLE + BODY
    for lay in prs.slide_layouts:
        types = {ph.placeholder_format.type for ph in lay.placeholders}
        if PPT.TITLE in types and PPT.BODY in types:
            return lay

    return prs.slide_layouts[0]


# ─── Template-First: Placeholder filling ─────────────────────────────────────


def _set_ph_text(ph, text: str) -> None:
    """Write text into a placeholder, one line per paragraph.
    Does NOT set font/color/size — those are inherited from the master."""
    tf = ph.text_frame
    tf.clear()
    if not text:
        return
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.add_run().text = line


def _set_ph_lines(ph, lines: list[str], bold_indices: set[int] | None = None) -> None:
    """Write structured lines into a placeholder.
    Indices listed in bold_indices are rendered bold (e.g. section headings).
    All other formatting is inherited from the master slide."""
    tf = ph.text_frame
    tf.clear()
    bold_indices = bold_indices or set()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if i in bold_indices:
            run.font.bold = True


def _fill_slide_content(slide, blueprint: dict) -> None:
    """
    Fill slide placeholders from blueprint content.

    Rules:
    - NEVER add new shapes (no add_textbox, no add_picture for text).
    - ONLY write into placeholders that already exist on the slide (inherited
      from the slide layout / master).
    - All visual styling (font, color, size, position, background) is owned
      by the template master — we only insert plain text runs.
    - Placeholder discovery uses idx (0=title, 1=primary body, 2=secondary body)
      and placeholder_format.type as a fallback. The template's placeholder names
      are auto-generated Google Slides IDs and are NOT used for matching.
    """
    layout_name = blueprint.get("layout", "title_content")
    content     = blueprint.get("content", {})
    title_text  = blueprint.get("title", "")

    # ── Placeholder discovery ─────────────────────────────────────────────────
    # idx=0  → title / center-title
    # idx=1  → primary body (or subtitle on TITLE layout)
    # idx=2  → secondary body (right column on two-column layouts)
    # idx=12 → slide number — skip
    ph_by_idx: dict[int, object] = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    title_ph = (ph_by_idx.get(0)
                or next((ph for ph in slide.placeholders
                         if ph.placeholder_format.type in (PPT.TITLE, PPT.CENTER_TITLE)), None))

    body_phs = [
        ph for idx, ph in sorted(ph_by_idx.items())
        if idx >= 1
        and ph.placeholder_format.type in (PPT.BODY, PPT.SUBTITLE, PPT.OBJECT)
    ]

    logger.debug(
        "_fill_slide_content: layout=%s  title_ph=%s  body_phs=%d",
        layout_name, title_ph is not None, len(body_phs),
    )

    # ── key_message: hero text goes into the title placeholder ───────────────
    # (MAIN_POINT layout has only a title placeholder, no body)
    if layout_name == "key_message":
        msg = content.get("message") or title_text
        if title_ph:
            _set_ph_text(title_ph, msg)
        if body_phs and content.get("subtext"):
            _set_ph_text(body_phs[0], content["subtext"])
        return

    # ── All other layouts: title text → idx=0 placeholder ────────────────────
    if title_ph:
        _set_ph_text(title_ph, title_text)

    if not body_phs:
        # Layout has no body placeholder (e.g. TITLE_ONLY, SECTION_HEADER).
        # Nothing more to fill — design is complete as-is.
        logger.debug("_fill_slide_content: no body placeholder for layout=%s", layout_name)
        return

    if layout_name == "title_content":
        items = content.get("items", [])
        _set_ph_lines(body_phs[0], [f"\u2022 {i}" for i in items])

    elif layout_name in ("two_column", "comparison"):
        lc = content.get("left", {})
        rc = content.get("right", {})
        l_label = lc.get("heading") or lc.get("label", "")
        r_label = rc.get("heading") or rc.get("label", "")
        l_lines = ([l_label] if l_label else []) + [f"\u2022 {i}" for i in lc.get("items", [])]
        r_lines = ([r_label] if r_label else []) + [f"\u2022 {i}" for i in rc.get("items", [])]

        if len(body_phs) >= 2:
            _set_ph_lines(body_phs[0], l_lines, bold_indices={0} if l_label else set())
            _set_ph_lines(body_phs[1], r_lines, bold_indices={0} if r_label else set())
        else:
            # Template has only one body placeholder — merge both columns into it
            merged = l_lines + ["", "\u2500" * 20, ""] + r_lines
            bold = ({0} if l_label else set()) | (
                {len(l_lines) + 3} if r_label else set()
            )
            _set_ph_lines(body_phs[0], merged, bold_indices=bold)

    elif layout_name == "icon_grid":
        cards = content.get("cards", [])
        icons = ["\u2460", "\u2461", "\u2462", "\u2463"]  # ①②③④
        lines: list[str] = []
        bold_idx: set[int] = set()
        for i, card in enumerate(cards):
            icon = icons[i] if i < len(icons) else f"({i + 1})"
            bold_idx.add(len(lines))
            lines.append(f"{icon} {card.get('heading', '')}")
            if card.get("text"):
                lines.append(card["text"])
            if i < len(cards) - 1:
                lines.append("")
        _set_ph_lines(body_phs[0], lines, bold_indices=bold_idx)

    elif layout_name == "process_flow":
        steps = content.get("steps", [])
        lines = []
        bold_idx: set[int] = set()
        for i, step in enumerate(steps):
            bold_idx.add(len(lines))
            lines.append(f"{i + 1}. {step.get('label', '')}")
            if step.get("desc"):
                lines.append(f"   {step['desc']}")
            if i < len(steps) - 1:
                lines.append("")
        _set_ph_lines(body_phs[0], lines, bold_indices=bold_idx)

    elif layout_name == "timeline":
        steps = content.get("steps", [])
        lines = []
        for i, step in enumerate(steps):
            lines.append(f"{step.get('label', '')}  \u203a  {step.get('event', '')}")
            if i < len(steps) - 1:
                lines.append("")
        _set_ph_lines(body_phs[0], lines)

    elif layout_name == "big_stat":
        value   = content.get("value", "")
        lbl     = content.get("label", "")
        ctx     = content.get("context", [])
        lines   = []
        bold_idx: set[int] = set()
        if value:
            bold_idx.add(0)
            lines.append(value)
        if lbl:
            lines.append(lbl)
        if ctx:
            lines.append("")
            lines.extend(f"\u2022 {item}" for item in ctx)
        _set_ph_lines(body_phs[0], lines, bold_indices=bold_idx)

    elif layout_name == "section_divider":
        subtitle = content.get("subtitle", "")
        sub_ph = next(
            (ph for ph in body_phs if ph.placeholder_format.type == PPT.SUBTITLE),
            body_phs[0],
        )
        _set_ph_text(sub_ph, subtitle)

    elif layout_name == "quote":
        q    = content.get("quote", "")
        attr = content.get("attribution", "")
        lines = []
        if q:
            lines.append(f"\u201C{q}\u201D")
        if attr:
            lines += ["", f"\u2014 {attr}"]
        _set_ph_lines(body_phs[0], lines)

    elif layout_name in ("chart_bar", "chart_pie"):
        # Body placeholder is intentionally left empty — the chart shape is added
        # as a floating graphic below by _add_chart_shape().
        # We clear the body so no placeholder prompt text leaks through.
        _set_ph_text(body_phs[0], "")

    # section_divider, quote, timeline, big_stat — handled above; nothing else needed.


def _add_chart_shape(slide, blueprint: dict, slide_w, slide_h) -> None:
    """Add a bar or pie chart as a floating shape positioned below the title."""
    layout_name = blueprint.get("layout", "")
    content     = blueprint.get("content", {})

    # Find where the title placeholder ends
    title_bottom = int(slide_h * 0.22)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title_bottom = ph.top + ph.height
            break

    chart_left   = int(slide_w * 0.04)
    chart_top    = int(title_bottom)
    chart_width  = int(slide_w * 0.92)
    chart_height = int(slide_h - title_bottom - slide_h * 0.04)

    try:
        if layout_name == "chart_bar":
            categories  = content.get("categories", [])
            series_data = content.get("series", [])
            if not categories or not series_data:
                return
            chart_data = ChartData()
            chart_data.categories = categories
            for s in series_data:
                chart_data.add_series(s.get("name", ""), tuple(s.get("values", [])))
            frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                chart_left, chart_top, chart_width, chart_height,
                chart_data,
            )
            frame.chart.has_legend = len(series_data) > 1

        elif layout_name == "chart_pie":
            slices = content.get("slices", [])
            if not slices:
                return
            chart_data = ChartData()
            chart_data.categories = [s.get("label", f"Item {i+1}") for i, s in enumerate(slices)]
            chart_data.add_series("", tuple(s.get("value", 0) for s in slices))
            frame = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                chart_left, chart_top, chart_width, chart_height,
                chart_data,
            )
            chart = frame.chart
            chart.has_legend = True
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.show_percentage = True
            plot.data_labels.show_category_name = False

    except Exception as e:
        logger.warning(f"_add_chart_shape ({layout_name}): {e}")


# ─── Scratch renderer (metrics_grid, section_divider) ────────────────────────

_SCRATCH_LAYOUTS = {"metrics_grid", "section_divider"}

_DEFAULT_PRIMARY = "1A3C8F"   # fallback when BrandContext has no primary_color


def _scratch_rgb(hex_str: str | None, default: str = _DEFAULT_PRIMARY) -> RGBColor:
    h = (hex_str or default).lstrip("#")
    if len(h) != 6:
        h = default
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _clear_placeholders(slide) -> None:
    """Remove all placeholder shapes so we start with a blank canvas."""
    from lxml import etree
    sp_tree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        try:
            sp_tree.remove(ph._element)
        except Exception:
            pass


def _add_textbox(slide, x: float, y: float, w: float, h: float,
                 text: str, font_size: int, bold: bool,
                 color: RGBColor, slide_w: int, slide_h: int,
                 align_center: bool = False) -> None:
    """
    Add a text box using normalised (0–1) coordinates.
    x, y, w, h are fractions of slide_w / slide_h.
    """
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN

    left   = int(x * slide_w)
    top    = int(y * slide_h)
    width  = int(w * slide_w)
    height = int(h * slide_h)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    para = tf.paragraphs[0]
    if align_center:
        para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = color


def _render_metrics_grid(slide, blueprint: dict, brand: BrandContext) -> None:
    """
    Render a metrics_grid slide from scratch.
    Cards are laid out in rows of 3–4, evenly spaced.
    """
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    sw = slide.shapes._spTree.getparent().getparent()   # not used — use slide dims
    # Grab slide dimensions from the presentation object attached to the slide part
    try:
        slide_w = slide.part.package.presentation.slide_width
        slide_h = slide.part.package.presentation.slide_height
    except Exception:
        slide_w = int(W)
        slide_h = int(H)

    # Explicit colors — never rely on theme/template defaults.
    # White cards: value in white (readable on any brand-colored background),
    # label in light gray.
    value_rgb   = RGBColor(0xFF, 0xFF, 0xFF)          # white
    label_rgb   = RGBColor(0xCC, 0xCC, 0xCC)          # light gray
    muted_rgb   = RGBColor(0xAA, 0xAA, 0xAA)          # dimmer gray for sublabel
    title_rgb   = RGBColor(0xFF, 0xFF, 0xFF)          # white title on brand bg

    # Card background uses primary brand color so white text is always legible
    card_bg_rgb = _scratch_rgb(brand.primary_color if brand else None)

    _clear_placeholders(slide)

    title_text = blueprint.get("title", "")
    if title_text:
        _add_textbox(
            slide, x=0.04, y=0.04, w=0.92, h=0.12,
            text=title_text, font_size=28, bold=True,
            color=title_rgb, slide_w=slide_w, slide_h=slide_h,
        )

    metrics = (blueprint.get("content") or {}).get("metrics", [])
    if not metrics:
        return

    n = len(metrics)
    cols = 4 if n >= 5 else 3
    rows = (n + cols - 1) // cols

    # Layout constants (normalised)
    margin_x = 0.04
    margin_y_top = 0.20   # below title
    margin_y_bot = 0.04
    gap_x = 0.02
    gap_y = 0.04

    total_w = 1.0 - 2 * margin_x - gap_x * (cols - 1)
    card_w  = total_w / cols
    avail_h = 1.0 - margin_y_top - margin_y_bot - gap_y * (rows - 1)
    card_h  = avail_h / rows

    print(f"[metrics_grid] slide={slide_w}x{slide_h} EMU, cards={len(metrics)}")

    for idx, metric in enumerate(metrics):
        col = idx % cols
        row = idx // cols

        cx = margin_x + col * (card_w + gap_x)
        cy = margin_y_top + row * (card_h + gap_y)

        left   = int(cx * slide_w)
        top    = int(cy * slide_h)
        width  = int(card_w * slide_w)
        height = int(card_h * slide_h)

        print(f"  card {idx}: left={left} top={top} width={width} height={height}")

        # Brand-colored rounded rectangle card (white text stays legible on any brand color)
        card = slide.shapes.add_shape(
            1,   # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
            left, top, width, height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg_rgb
        card.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.width = Pt(0.5)
        # Adjust corner rounding
        try:
            card._element.spPr.prstGeom.avLst.clear()
            from lxml import etree
            etree.SubElement(
                card._element.spPr.prstGeom.avLst,
                qn("a:gd"), name="adj", fmla="val 20000",
            )
        except Exception:
            pass

        # Value text (large, bold, white) — top 55% of card
        val_h = card_h * 0.55
        _add_textbox(
            slide,
            x=cx + 0.01, y=cy + 0.04,
            w=card_w - 0.02, h=val_h - 0.04,
            text=metric.get("value", ""),
            font_size=36, bold=True,
            color=value_rgb,
            slide_w=slide_w, slide_h=slide_h,
            align_center=True,
        )

        # Label text (small, light gray) — below value
        label_y = cy + val_h
        label_h = card_h * 0.28
        _add_textbox(
            slide,
            x=cx + 0.01, y=label_y,
            w=card_w - 0.02, h=label_h,
            text=metric.get("label", ""),
            font_size=12, bold=False,
            color=label_rgb,
            slide_w=slide_w, slide_h=slide_h,
            align_center=True,
        )

        # Sublabel (optional, dimmer gray)
        sublabel = metric.get("sublabel") or ""
        if sublabel:
            sublabel_y = label_y + label_h
            sublabel_h = card_h * 0.18
            _add_textbox(
                slide,
                x=cx + 0.01, y=sublabel_y,
                w=card_w - 0.02, h=sublabel_h,
                text=sublabel,
                font_size=9, bold=False,
                color=muted_rgb,
                slide_w=slide_w, slide_h=slide_h,
                align_center=True,
            )


def _render_section_divider_scratch(slide, blueprint: dict, brand: BrandContext) -> None:
    """
    Render a section_divider slide from scratch: title centered on blank canvas.
    """
    try:
        slide_w = slide.part.package.presentation.slide_width
        slide_h = slide.part.package.presentation.slide_height
    except Exception:
        slide_w = int(W)
        slide_h = int(H)

    primary_rgb = _scratch_rgb(brand.primary_color if brand else None)

    _clear_placeholders(slide)

    title_text = blueprint.get("title", "")
    if title_text:
        _add_textbox(
            slide,
            x=0.08, y=0.35, w=0.84, h=0.30,
            text=title_text, font_size=40, bold=True,
            color=primary_rgb,
            slide_w=slide_w, slide_h=slide_h,
            align_center=True,
        )


def render_from_scratch(slide, blueprint: dict, brand: BrandContext | None) -> None:
    """
    Entry point for layouts that build their own shapes rather than filling
    placeholders from the PPTX template.

    Currently handles: metrics_grid, section_divider.
    """
    try:
        layout = blueprint.get("layout", "")
        if layout == "metrics_grid":
            _render_metrics_grid(slide, blueprint, brand)
        elif layout == "section_divider":
            _render_section_divider_scratch(slide, blueprint, brand)
    except Exception as e:
        import traceback
        print(f"[metrics_grid] RENDER ERROR: {e}")
        traceback.print_exc()


# ─── Core renderer ────────────────────────────────────────────────────────────


def render_slide_pptx(
    blueprint: dict,
    colors: "BrandColors",          # kept in signature; not used for layout
    template_pptx_path: str | None = None,
    brand_context: BrandContext | None = None,
) -> Presentation:
    """
    Template-First rendering pipeline:
      1. Load template PPTX — master contains background, fonts, colors, positions
      2. Remove all existing slides (master + layouts stay intact)
      3. Find the slide layout whose name best matches our layout type
      4. add_slide(layout) → entire master design is inherited automatically
      5. Fill only the placeholders with blueprint text
      6. For chart layouts, add a floating chart shape
    No coordinate math. No hardcoded colors. Design = 100% template-driven.
    """
    layout_name = blueprint.get("layout", "title_content")

    if template_pptx_path and os.path.exists(template_pptx_path):
        prs = Presentation(template_pptx_path)
        # Remove pre-existing slides while keeping master + layouts intact.
        _r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        sldIdLst = prs.slides._sldIdLst
        rids = [sld.get(_r_ns) for sld in list(sldIdLst)]
        for sld_id in list(sldIdLst):
            sldIdLst.remove(sld_id)
        for rId in rids:
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
    else:
        prs = Presentation()
        prs.slide_width  = W
        prs.slide_height = H

    slide_layout = _find_layout(prs, layout_name)
    slide        = prs.slides.add_slide(slide_layout)

    if layout_name in _SCRATCH_LAYOUTS:
        render_from_scratch(slide, blueprint, brand_context)
    else:
        _fill_slide_content(slide, blueprint)
        if brand_context is not None:
            apply_brand_styling(slide, brand_context)
        if layout_name in ("chart_bar", "chart_pie"):
            _add_chart_shape(slide, blueprint, prs.slide_width, prs.slide_height)

    logger.debug(
        "render_slide_pptx: layout=%s → tmpl_layout=%r  template=%s",
        layout_name, slide_layout.name, bool(template_pptx_path),
    )
    return prs


# ─── Thumbnail ───────────────────────────────────────────────────────────────
#
# Two strategies:
#   • Branded (background_image_path set): PIL compositor — resize bg image to
#     1920×1080, draw semi-transparent title bar + title text.  Bypasses fitz
#     which cannot render embedded PPTX background images (only ~2% non-white).
#   • Plain (no background): fitz renders the PPTX directly.

_THUMB_W = 1920
_THUMB_H = 1080


def _render_thumbnail(pptx_path: str, out_dir: str) -> str:
    """
    PPTX → PNG via PyMuPDF (fitz).
    fitz reports 13.333"×7.5" slides as 400×600 pt (portrait); we compensate
    by computing per-axis scale factors that map back to 1920×1080.
    """
    from PIL import Image

    SLIDE_PT_W = 960.0   # 13.333" × 72 pt/in
    SLIDE_PT_H = 540.0   # 7.5"    × 72 pt/in

    doc  = fitz.open(pptx_path)
    page = doc[0]
    rect = page.rect

    sx = (_THUMB_W / rect.width)  * (rect.width  / SLIDE_PT_W)
    sy = (_THUMB_H / rect.height) * (rect.height / SLIDE_PT_H)

    pix = page.get_pixmap(matrix=fitz.Matrix(sx, sy), alpha=False)
    doc.close()

    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    if img.size != (_THUMB_W, _THUMB_H):
        img = img.resize((_THUMB_W, _THUMB_H), Image.LANCZOS)

    out = os.path.join(out_dir, "0.png")
    img.save(out, "PNG")
    return out


def _blueprint_body_preview(blueprint: dict) -> list[str]:
    """
    Extract up to 3 preview lines from a blueprint's content for thumbnail display.
    These are shown below the title bar in the branded thumbnail.
    """
    layout  = blueprint.get("layout", "")
    content = blueprint.get("content", {})
    lines: list[str] = []

    if layout == "icon_grid":
        for card in (content.get("cards") or [])[:4]:
            lines.append(f"• {card.get('heading', '')}")
    elif layout in ("process_flow",):
        for i, step in enumerate((content.get("steps") or [])[:4], 1):
            lines.append(f"{i}. {step.get('label', '')}")
    elif layout == "key_message":
        if content.get("message"):
            lines.append(content["message"])
        if content.get("subtext"):
            lines.append(content["subtext"])
    elif layout == "big_stat":
        if content.get("value"):
            lines.append(content["value"])
        if content.get("label"):
            lines.append(content["label"])
    elif layout == "section_divider":
        if content.get("subtitle"):
            lines.append(content["subtitle"])
    elif layout == "quote":
        if content.get("quote"):
            lines.append(f"\u201C{content['quote'][:80]}\u201D")
    elif layout in ("two_column", "comparison"):
        lc = content.get("left") or {}
        rc = content.get("right") or {}
        if lc.get("heading") or lc.get("label"):
            lines.append((lc.get("heading") or lc.get("label", "")) + " / " +
                         (rc.get("heading") or rc.get("label", "")))
        for item in (lc.get("items") or [])[:2]:
            lines.append(f"• {item}")
    elif layout == "timeline":
        for step in (content.get("steps") or [])[:4]:
            lines.append(f"{step.get('label', '')}  \u2192  {step.get('event', '')[:40]}")
    elif layout in ("chart_bar", "chart_pie"):
        lines.append("[ chart ]")
    else:  # title_content
        for item in (content.get("items") or [])[:4]:
            lines.append(f"• {item}")

    return [l for l in lines if l.strip()][:4]


def _thumb_font(size: int):
    """Load Helvetica or fall back to default PIL font."""
    from PIL import ImageFont
    for path in (
        "/System/Library/Fonts/Helvetica.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ):
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _thumb_font_bold(size: int):
    """Load Helvetica Bold or fall back."""
    from PIL import ImageFont
    for path in (
        "/System/Library/Fonts/Helvetica.ttc",   # index 1 for bold, but truetype handles it
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ):
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _draw_body_thumbnail(draw, blueprint: dict, body_rgb: tuple, bar_h: int) -> None:
    """
    Draw layout-specific body content on the thumbnail image.
    Each layout gets a visually distinct representation.
    """
    from PIL import ImageDraw as _ID

    layout  = blueprint.get("layout", "title_content")
    content = blueprint.get("content", {})

    # Body area dimensions
    PAD   = 52
    top   = bar_h + 30
    left  = PAD
    right = _THUMB_W - PAD
    width = right - left
    avail = _THUMB_H - top - 30   # available height for body

    # ── icon_grid: 2×2 card grid ─────────────────────────────────────────────
    if layout == "icon_grid":
        cards = (content.get("cards") or [])[:4]
        if not cards:
            return
        cols = 2
        rows = (len(cards) + 1) // 2
        card_w = (width - PAD) // cols
        card_h = min(180, (avail - PAD * (rows - 1)) // rows)
        card_font   = _thumb_font(max(22, int(_THUMB_H * 0.028)))
        heading_font = _thumb_font_bold(max(26, int(_THUMB_H * 0.032)))
        for i, card in enumerate(cards):
            col = i % cols
            row = i // cols
            cx = left + col * (card_w + PAD)
            cy = top  + row * (card_h + 20)
            # Card box
            draw.rounded_rectangle(
                [(cx, cy), (cx + card_w, cy + card_h)],
                radius=12,
                fill=(255, 255, 255, 25),
                outline=(255, 255, 255, 60),
                width=1,
            )
            # Heading
            heading = (card.get("heading") or "")[:28]
            draw.text((cx + 16, cy + 14), heading, fill=body_rgb, font=heading_font)
            # Body text
            text = (card.get("text") or "")[:60]
            if text:
                draw.text((cx + 16, cy + 14 + 36), text, fill=(body_rgb[0], body_rgb[1], body_rgb[2]), font=card_font)

    # ── big_stat: huge number centered ────────────────────────────────────────
    elif layout == "big_stat":
        value = str(content.get("value") or "")[:12]
        label = str(content.get("label") or "")[:50]
        ctx   = (content.get("context") or [])[:3]

        if value:
            stat_font = _thumb_font_bold(max(120, int(_THUMB_H * 0.15)))
            bbox = draw.textbbox((0, 0), value, font=stat_font)
            stat_w = bbox[2] - bbox[0]
            stat_x = (_THUMB_W - stat_w) // 2
            stat_y = top + max(10, (avail - 200) // 3)
            draw.text((stat_x, stat_y), value, fill=body_rgb, font=stat_font)

            if label:
                lbl_font = _thumb_font(max(28, int(_THUMB_H * 0.036)))
                lbl_bbox = draw.textbbox((0, 0), label, font=lbl_font)
                lbl_x = (_THUMB_W - (lbl_bbox[2] - lbl_bbox[0])) // 2
                draw.text((lbl_x, stat_y + 160), label, fill=body_rgb, font=lbl_font)

        ctx_font = _thumb_font(max(22, int(_THUMB_H * 0.026)))
        cy = top + avail - len(ctx) * 38 - 10
        for item in ctx:
            item_text = f"• {item}"[:70]
            draw.text((left, cy), item_text, fill=body_rgb, font=ctx_font)
            cy += 38

    # ── two_column / comparison: two panels ───────────────────────────────────
    elif layout in ("two_column", "comparison"):
        lc = content.get("left") or {}
        rc = content.get("right") or {}
        l_label = str(lc.get("heading") or lc.get("label", ""))[:30]
        r_label = str(rc.get("heading") or rc.get("label", ""))[:30]
        l_items = (lc.get("items") or [])[:4]
        r_items = (rc.get("items") or [])[:4]

        col_w     = (width - 60) // 2
        sep_x     = left + col_w + 30
        head_font = _thumb_font_bold(max(26, int(_THUMB_H * 0.032)))
        item_font = _thumb_font(max(22, int(_THUMB_H * 0.026)))

        # Vertical divider
        draw.line([(sep_x, top), (sep_x, top + avail - 20)],
                  fill=(255, 255, 255, 80), width=2)

        for col_idx, (label, items) in enumerate([(l_label, l_items), (r_label, r_items)]):
            cx = left if col_idx == 0 else sep_x + 30
            cy = top
            if label:
                draw.text((cx, cy), label, fill=body_rgb, font=head_font)
                cy += 44
            for item in items:
                draw.text((cx, cy), f"• {item[:38]}", fill=body_rgb, font=item_font)
                cy += 34

    # ── process_flow: numbered steps ──────────────────────────────────────────
    elif layout == "process_flow":
        steps = (content.get("steps") or [])[:5]
        step_font  = _thumb_font_bold(max(26, int(_THUMB_H * 0.032)))
        desc_font  = _thumb_font(max(22, int(_THUMB_H * 0.026)))
        step_h     = min(160, avail // max(len(steps), 1))
        num_colors = [(0x3B, 0x82, 0xF6), (0x06, 0xB6, 0xD4), (0x10, 0xB9, 0x81),
                      (0xF5, 0x9E, 0x0B), (0xEF, 0x44, 0x44)]

        for i, step in enumerate(steps):
            cy  = top + i * step_h
            nc  = num_colors[i % len(num_colors)]
            # Circle number
            r = 26
            draw.ellipse([(left, cy), (left + r * 2, cy + r * 2)], fill=nc)
            num_font = _thumb_font_bold(max(22, int(_THUMB_H * 0.026)))
            draw.text((left + r - 8, cy + r - 14), str(i + 1), fill=(255, 255, 255), font=num_font)
            # Label + desc
            label = str(step.get("label", ""))[:28]
            desc  = str(step.get("desc", ""))[:50]
            draw.text((left + r * 2 + 16, cy + 2), label, fill=body_rgb, font=step_font)
            if desc:
                draw.text((left + r * 2 + 16, cy + 38), desc, fill=body_rgb, font=desc_font)

            # Arrow between steps (except last)
            if i < len(steps) - 1:
                ax = left + r
                ay = cy + r * 2 + 4
                draw.polygon([(ax - 10, ay), (ax + 10, ay), (ax, ay + 14)],
                             fill=(255, 255, 255, 80))

    # ── key_message: large centered message ───────────────────────────────────
    elif layout == "key_message":
        message = str(content.get("message") or "")[:80]
        subtext = str(content.get("subtext") or "")[:100]
        msg_font = _thumb_font_bold(max(52, int(_THUMB_H * 0.065)))
        sub_font = _thumb_font(max(28, int(_THUMB_H * 0.034)))

        # Wrap message across lines (~30 chars each)
        words, lines, line = message.split(), [], []
        for w in words:
            line.append(w)
            if len(" ".join(line)) > 32:
                lines.append(" ".join(line[:-1]))
                line = [w]
        if line:
            lines.append(" ".join(line))

        total_h = len(lines) * 70 + (40 if subtext else 0)
        cy = top + max(20, (avail - total_h) // 2)
        for msg_line in lines[:3]:
            bbox = draw.textbbox((0, 0), msg_line, font=msg_font)
            tx = (_THUMB_W - (bbox[2] - bbox[0])) // 2
            draw.text((tx, cy), msg_line, fill=body_rgb, font=msg_font)
            cy += 70

        if subtext:
            bbox = draw.textbbox((0, 0), subtext[:80], font=sub_font)
            tx = (_THUMB_W - (bbox[2] - bbox[0])) // 2
            draw.text((tx, cy + 10), subtext[:80], fill=body_rgb, font=sub_font)

    # ── quote: quotation marks + attribution ──────────────────────────────────
    elif layout == "quote":
        quote = str(content.get("quote") or "")[:120]
        attr  = str(content.get("attribution") or "")[:60]
        q_font = _thumb_font(max(32, int(_THUMB_H * 0.040)))
        a_font = _thumb_font(max(24, int(_THUMB_H * 0.028)))

        # Big quotation mark
        qm_font = _thumb_font_bold(max(100, int(_THUMB_H * 0.12)))
        draw.text((left, top - 20), "\u201C", fill=(255, 255, 255, 100), font=qm_font)

        # Wrap quote
        words, lines, line = quote.split(), [], []
        for w in words:
            line.append(w)
            if len(" ".join(line)) > 45:
                lines.append(" ".join(line[:-1]))
                line = [w]
        if line:
            lines.append(" ".join(line))

        cy = top + 60
        for qline in lines[:4]:
            draw.text((left + 20, cy), qline, fill=body_rgb, font=q_font)
            cy += 46

        if attr:
            draw.text((left + 20, cy + 20), f"\u2014 {attr}", fill=body_rgb, font=a_font)

    # ── chart_bar: simple bar chart visual ────────────────────────────────────
    elif layout == "chart_bar":
        cats   = (content.get("categories") or [])[:6]
        series = (content.get("series") or [])[:1]
        if cats and series:
            values = [float(v) for v in (series[0].get("values") or [])[:len(cats)]]
            max_v  = max(values) if values else 1
            bar_colors = [(0x3B, 0x82, 0xF6), (0x06, 0xB6, 0xD4), (0x10, 0xB9, 0x81),
                          (0xF5, 0x9E, 0x0B), (0xEF, 0x44, 0x44), (0x8B, 0x5C, 0xF6)]
            chart_h   = avail - 60
            bar_w     = max(40, (width - 40) // max(len(cats), 1) - 16)
            cat_font  = _thumb_font(max(18, int(_THUMB_H * 0.022)))
            baseline  = top + chart_h

            for i, (cat, val) in enumerate(zip(cats, values)):
                bx = left + 20 + i * (bar_w + 16)
                bh = int(chart_h * 0.85 * (val / max_v)) if max_v > 0 else 20
                bc = bar_colors[i % len(bar_colors)]
                draw.rounded_rectangle(
                    [(bx, baseline - bh), (bx + bar_w, baseline)],
                    radius=6, fill=bc,
                )
                # Category label
                cat_label = str(cat)[:10]
                draw.text((bx, baseline + 6), cat_label, fill=body_rgb, font=cat_font)
            # Baseline
            draw.line([(left, baseline), (right, baseline)], fill=(255, 255, 255, 60), width=2)

    # ── chart_pie: simple pie/donut visual ────────────────────────────────────
    elif layout == "chart_pie":
        slices = (content.get("slices") or [])[:6]
        if slices:
            total  = sum(float(s.get("value", 0)) for s in slices) or 1
            colors_list = [
                (0x3B, 0x82, 0xF6), (0x06, 0xB6, 0xD4), (0x10, 0xB9, 0x81),
                (0xF5, 0x9E, 0x0B), (0xEF, 0x44, 0x44), (0x8B, 0x5C, 0xF6),
            ]
            import math
            cx_pie = _THUMB_W // 2
            cy_pie = top + avail // 2
            radius = min(avail // 2 - 20, 220)
            angle  = -90.0
            for i, sl in enumerate(slices):
                sweep = float(sl.get("value", 0)) / total * 360
                col   = colors_list[i % len(colors_list)]
                draw.pieslice(
                    [(cx_pie - radius, cy_pie - radius),
                     (cx_pie + radius, cy_pie + radius)],
                    start=angle, end=angle + sweep,
                    fill=col,
                )
                angle += sweep

            # Legend
            leg_font = _thumb_font(max(20, int(_THUMB_H * 0.024)))
            lx = right - 280
            ly = top + 20
            for i, sl in enumerate(slices):
                col = colors_list[i % len(colors_list)]
                draw.rectangle([(lx, ly + i * 32), (lx + 20, ly + i * 32 + 20)], fill=col)
                label = str(sl.get("label", ""))[:18]
                val   = float(sl.get("value", 0))
                draw.text((lx + 28, ly + i * 32), f"{label} {val:.0f}%", fill=body_rgb, font=leg_font)

    # ── timeline: milestones with line ────────────────────────────────────────
    elif layout == "timeline":
        steps    = (content.get("steps") or [])[:5]
        t_font   = _thumb_font_bold(max(24, int(_THUMB_H * 0.028)))
        e_font   = _thumb_font(max(20, int(_THUMB_H * 0.024)))
        step_gap = avail // max(len(steps), 1)
        line_y   = top + avail // 2

        # Horizontal timeline line
        draw.line([(left + 20, line_y), (right - 20, line_y)],
                  fill=(0x3B, 0x82, 0xF6), width=4)

        for i, step in enumerate(steps):
            sx = left + 20 + i * ((right - left - 40) // max(len(steps) - 1, 1))
            # Dot on line
            draw.ellipse([(sx - 10, line_y - 10), (sx + 10, line_y + 10)],
                         fill=(0x3B, 0x82, 0xF6))
            label = str(step.get("label", ""))[:14]
            event = str(step.get("event", ""))[:30]
            if i % 2 == 0:  # above line
                draw.text((sx - 40, line_y - 60), label, fill=body_rgb, font=t_font)
                draw.text((sx - 40, line_y - 30), event, fill=body_rgb, font=e_font)
            else:            # below line
                draw.text((sx - 40, line_y + 20), label, fill=body_rgb, font=t_font)
                draw.text((sx - 40, line_y + 46), event, fill=body_rgb, font=e_font)

    # ── section_divider: subtitle centered ────────────────────────────────────
    elif layout == "section_divider":
        subtitle = str(content.get("subtitle") or "")[:100]
        sub_font = _thumb_font(max(36, int(_THUMB_H * 0.044)))
        if subtitle:
            bbox = draw.textbbox((0, 0), subtitle, font=sub_font)
            tx = (_THUMB_W - (bbox[2] - bbox[0])) // 2
            ty = top + (avail - (bbox[3] - bbox[1])) // 2
            draw.text((tx, ty), subtitle, fill=body_rgb, font=sub_font)

    # ── title_content / default: bullet list ──────────────────────────────────
    else:
        items     = (content.get("items") or [])[:5]
        item_font = _thumb_font(max(28, int(_THUMB_H * 0.034)))
        cy        = top + 10
        for item in items:
            text = f"• {item}"[:70]
            draw.text((left, cy), text, fill=body_rgb, font=item_font)
            cy += max(38, int(_THUMB_H * 0.046))
            if cy > _THUMB_H - 40:
                break


def _render_thumbnail_branded(
    pptx_path: str,
    out_dir: str,
    colors: "BrandColors",
    blueprint: dict,
) -> str:
    """
    PIL-based thumbnail for branded slides (background_image_path is set).
    Composes: background image → semi-transparent title bar → title text →
    layout-specific body visuals (cards, charts, big numbers, columns, etc.).
    Falls back to fitz if PIL fails.
    """
    from PIL import Image, ImageDraw

    try:
        bg   = Image.open(colors.background_image_path).convert("RGB")
        bg   = bg.resize((_THUMB_W, _THUMB_H), Image.LANCZOS)
        draw = ImageDraw.Draw(bg, "RGBA")

        # ── Title bar ────────────────────────────────────────────────────────
        sc = colors.shape_color.lstrip("#")
        shape_r, shape_g, shape_b = int(sc[0:2], 16), int(sc[2:4], 16), int(sc[4:6], 16)
        bar_h     = int(_THUMB_H * 0.18)
        bar_alpha = int(max(0, min(100, colors.shape_opacity)) / 100 * 200)
        draw.rectangle([(0, 0), (_THUMB_W, bar_h)],
                       fill=(shape_r, shape_g, shape_b, bar_alpha))

        # Accent divider line
        acc = colors.secondary.lstrip("#")
        draw.rectangle([(0, bar_h), (_THUMB_W, bar_h + 6)],
                       fill=(int(acc[0:2], 16), int(acc[2:4], 16), int(acc[4:6], 16), 200))

        # ── Title text ────────────────────────────────────────────────────────
        title = (blueprint.get("title") or "")[:80]
        if title:
            tc         = colors.title_font_color.lstrip("#")
            title_rgb  = (int(tc[0:2], 16), int(tc[2:4], 16), int(tc[4:6], 16))
            font_size  = max(40, min(72, int(_THUMB_H * 0.055)))
            title_font = _thumb_font(font_size)
            draw.text((48, int(bar_h * 0.15)), title, fill=title_rgb, font=title_font)

        # ── Body: auto-detect text color against background ───────────────────
        # Parse background image average brightness in the body area
        # Use white text if background is dark (< 128), dark text if light
        try:
            import numpy as np
            body_region = bg.crop((0, bar_h + 10, _THUMB_W, _THUMB_H))
            avg_brightness = float(np.array(body_region).mean())
        except Exception:
            avg_brightness = 50.0  # assume dark

        if avg_brightness < 128:
            body_rgb = (220, 228, 240)  # near-white for dark backgrounds
        else:
            bc       = colors.body_font_color.lstrip("#")
            body_rgb = (int(bc[0:2], 16), int(bc[2:4], 16), int(bc[4:6], 16))

        # ── Layout-specific body drawing ──────────────────────────────────────
        _draw_body_thumbnail(draw, blueprint, body_rgb, bar_h)

        out = os.path.join(out_dir, "0.png")
        bg.save(out, "PNG")
        logger.debug("_render_thumbnail_branded: composed → %s (layout=%s)",
                     out, blueprint.get("layout"))
        return out

    except Exception as e:
        logger.warning("_render_thumbnail_branded failed (%s) — fitz fallback", e)
        return _render_thumbnail(pptx_path, out_dir)


# ─── Vision QA helper ────────────────────────────────────────────────────────


async def _render_with_vision_qa(
    blueprint: dict,
    colors: BrandColors,
    template_pptx_path: str | None,
    pptx_path: str,
    thumb_abs: str,
    gen_dir: Path,
    thumb_dir: Path,
    label: str = "",
    brand_context: BrandContext | None = None,
) -> dict:
    """
    Run GPT-4o Vision on the rendered thumbnail; if issues found, attempt
    exactly ONE re-render with the Vision-proposed fix.

    Contract:
    - Max 1 Vision call + 1 re-render (no loops, no recursion).
    - validate_and_trim() applied to Vision's output before re-rendering.
    - No-op guard: skip re-render if JSON diff is identical after trimming.
    - Returns the blueprint used for the final PPTX on disk.
    """
    from services.vision_validator import vision_validate
    import json as _json

    try:
        ok, fixed_bp = await vision_validate(thumb_abs, blueprint, colors)
    except Exception as e:
        logger.warning("Vision QA (%s): API error — skipping: %s", label, e)
        return blueprint

    if ok or not fixed_bp:
        return blueprint

    validate_and_trim(fixed_bp)

    if _json.dumps(fixed_bp, sort_keys=True) == _json.dumps(blueprint, sort_keys=True):
        logger.info("Vision QA (%s): fix identical after trim — skipping", label)
        return blueprint

    logger.info(
        "Vision QA (%s): re-rendering  layout: %s→%s  title: %r→%r",
        label,
        blueprint.get("layout"), fixed_bp.get("layout"),
        blueprint.get("title"), fixed_bp.get("title"),
    )

    pptx_fixed: str | None = None
    try:
        prs_fixed  = render_slide_pptx(fixed_bp, colors, template_pptx_path, brand_context)
        pptx_fixed = str(gen_dir / f"gen_{uuid.uuid4()}.pptx")
        prs_fixed.save(pptx_fixed)

        thumb_tmp = os.path.join(str(thumb_dir), "0_vision_fix.png")
        try:
            if colors.background_image_path and os.path.exists(colors.background_image_path):
                _render_thumbnail_branded(pptx_fixed, str(thumb_dir), colors, fixed_bp)
                os.rename(os.path.join(str(thumb_dir), "0.png"), thumb_tmp)
            else:
                _render_thumbnail(pptx_fixed, str(thumb_dir))
                os.rename(os.path.join(str(thumb_dir), "0.png"), thumb_tmp)
        except Exception as te:
            logger.error("Vision QA (%s): fixed thumbnail render failed: %s", label, te)
            Path(pptx_fixed).unlink(missing_ok=True)
            return blueprint

        os.replace(pptx_fixed, pptx_path)
        os.replace(thumb_tmp,  thumb_abs)
        logger.info("Vision QA (%s): re-render applied ✓", label)
        return fixed_bp

    except Exception as e:
        logger.warning("Vision QA (%s): re-render error — keeping original: %s", label, e)
        if pptx_fixed:
            Path(pptx_fixed).unlink(missing_ok=True)

    return blueprint


# ─── Single-slide LLM generation ─────────────────────────────────────────────

_SYSTEM_PROMPT = """You are an expert presentation designer creating visually stunning slides.

CRITICAL ANTI-TEXT-WALL RULES (no exceptions):
1. NEVER use title_content with more than 5 bullet points
2. title_content is the LAST RESORT — use visual layouts whenever possible
3. If content has 3-4 concepts/features/benefits → use icon_grid
4. If content has a key insight or strong statement → use key_message
5. If content has a process or sequential steps → use process_flow
6. If content has 3+ numbers/metrics → use chart_bar or chart_pie
7. Split any content that would produce >5 bullets into multiple slides

Available layouts and when to use them:
- icon_grid      — 3-4 concepts/features/benefits (PREFERRED over bullet lists)
- key_message    — one powerful statement that dominates the slide
- process_flow   — sequential steps (3-5) with descriptions
- chart_bar      — bar/column chart when you have 3+ comparable numeric values
- chart_pie      — pie chart for proportions/shares that sum to 100%
- big_stat       — one large dramatic number + label + 2-3 context bullets
- two_column     — two related topics side by side with bullet lists
- comparison     — explicit A vs B or before/after panels
- timeline       — milestones or historical sequence (max 5 points)
- quote          — powerful quote or testimonial with attribution
- section_divider — full-color section divider (for transitions between topics)
- title_content  — ONLY if none of the above fits AND max 5 short bullets

Respond with ONLY valid JSON (no markdown, no commentary):

icon_grid:       {"layout":"icon_grid","title":"...","content":{"cards":[{"heading":"...","text":"..."},{"heading":"...","text":"..."},{"heading":"...","text":"..."},{"heading":"...","text":"..."}]},"speaker_notes":"..."}
key_message:     {"layout":"key_message","title":"...","content":{"message":"One powerful statement — max 15 words","subtext":"Optional supporting detail or source"},"speaker_notes":"..."}
process_flow:    {"layout":"process_flow","title":"...","content":{"steps":[{"label":"Step name","desc":"Short desc"},{"label":"...","desc":"..."},{"label":"...","desc":"..."}]},"speaker_notes":"..."}
chart_bar:       {"layout":"chart_bar","title":"...","content":{"categories":["A","B","C","D"],"series":[{"name":"Metric","values":[10,20,15,30]}]},"speaker_notes":"..."}
chart_pie:       {"layout":"chart_pie","title":"...","content":{"slices":[{"label":"Category A","value":45},{"label":"Category B","value":30},{"label":"Category C","value":25}]},"speaker_notes":"..."}
title_content:   {"layout":"title_content","title":"...","content":{"type":"bullets","items":["...","..."]},"speaker_notes":"..."}
two_column:      {"layout":"two_column","title":"...","content":{"left":{"heading":"...","items":["..."]},"right":{"heading":"...","items":["..."]}},"speaker_notes":"..."}
big_stat:        {"layout":"big_stat","title":"...","content":{"value":"...","label":"...","context":["...","..."]},"speaker_notes":"..."}
section_divider: {"layout":"section_divider","title":"...","content":{"subtitle":"..."},"speaker_notes":"..."}
quote:           {"layout":"quote","title":"...","content":{"quote":"...","attribution":"..."},"speaker_notes":"..."}
comparison:      {"layout":"comparison","title":"...","content":{"left":{"label":"...","items":["..."]},"right":{"label":"...","items":["..."]}},"speaker_notes":"..."}
timeline:        {"layout":"timeline","title":"...","content":{"steps":[{"label":"...","event":"..."}]},"speaker_notes":"..."}

Rules:
- Title: max 60 chars, punchy
- icon_grid cards: 3 or 4 cards; heading max 30 chars; text max 80 chars
- key_message: message max 15 words — make it impactful
- process_flow: 3-5 steps; label max 25 chars; desc max 60 chars
- chart_bar/pie: use real numbers from context when available
- Match the language of the user's prompt
- speaker_notes: one sentence of presenter guidance"""


async def generate_blueprint(prompt: str, context: str = "") -> dict:
    kwargs: dict = {"api_key": settings.openai_api_key}
    if settings.openai_base_url:
        kwargs["base_url"] = settings.openai_base_url
    client = AsyncOpenAI(**kwargs)
    user_msg = f"Context: {context}\n\nSlide topic: {prompt}" if context else prompt
    resp = await client.chat.completions.create(
        model=settings.generator_model,
        messages=[
            {"role": "system", "content": _SYSTEM_PROMPT},
            {"role": "user",   "content": user_msg},
        ],
        temperature=0.7,
        max_tokens=1000,
    )
    raw = (resp.choices[0].message.content or "{}").strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())


# ─── Main entry point ─────────────────────────────────────────────────────────


async def generate_slide(
    db: Session,
    prompt: str,
    template_id: int | None,
    user_id: int | None,
    context: str = "",
) -> SlideLibraryEntry:
    """Full single-slide pipeline → returns saved SlideLibraryEntry."""

    # 1. Load brand template
    colors: BrandColors = BrandColors()
    template_pptx_path: str | None = None

    if template_id:
        tmpl = db.query(BrandTemplate).filter(BrandTemplate.id == template_id).first()
        if tmpl:
            if tmpl.pptx_path and os.path.exists(tmpl.pptx_path):
                template_pptx_path = tmpl.pptx_path
                stored = json.loads(tmpl.colors_json or "{}")
                if stored:
                    colors = BrandColors(**{k: v for k, v in stored.items()
                                           if k in BrandColors.__dataclass_fields__})
                else:
                    colors = _extract_brand_colors(template_pptx_path)

            if tmpl.font_family:       colors.font_family       = tmpl.font_family
            if tmpl.title_font_color:  colors.title_font_color  = tmpl.title_font_color
            if tmpl.title_font_size:   colors.title_font_size   = tmpl.title_font_size
            if tmpl.body_font_color:   colors.body_font_color   = tmpl.body_font_color
            if tmpl.body_font_size:    colors.body_font_size    = tmpl.body_font_size
            if tmpl.shape_color:
                colors.shape_color = tmpl.shape_color
                colors.primary     = tmpl.shape_color
            if tmpl.shape_opacity is not None:
                colors.shape_opacity = tmpl.shape_opacity
            if tmpl.background_image_path and os.path.exists(tmpl.background_image_path):
                colors.background_image_path = tmpl.background_image_path
            for field in ("title_x", "title_y", "title_w", "title_h",
                          "body_x",  "body_y",  "body_w",  "body_h"):
                val = getattr(tmpl, field, None)
                if val is not None:
                    setattr(colors, field, val)

    if settings.fixed_bg_image and os.path.exists(settings.fixed_bg_image):
        colors.background_image_path = settings.fixed_bg_image
    if settings.fixed_shape_color:
        colors.shape_color = settings.fixed_shape_color
        colors.primary     = settings.fixed_shape_color
    if settings.fixed_title_font_size > 0:
        colors.title_font_size = settings.fixed_title_font_size
    if settings.fixed_body_font_size > 0:
        colors.body_font_size = settings.fixed_body_font_size

    # Load structured brand context for PPTX styling
    brand_context: BrandContext | None = None
    if template_id:
        from services.brand_context import load_brand_context
        brand_context = load_brand_context(template_id, db)

    # 2. Generate + validate blueprint
    blueprint = await generate_blueprint(prompt, context)
    blueprint.setdefault("layout",  "title_content")
    blueprint.setdefault("title",   prompt[:60])
    blueprint.setdefault("content", {"type": "bullets", "items": []})
    validate_and_trim(blueprint)

    # 3. Render PPTX
    prs = render_slide_pptx(blueprint, colors, template_pptx_path, brand_context)

    # 4. Save PPTX
    gen_dir = Path(settings.upload_dir) / "generated"
    gen_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = str(gen_dir / f"gen_{uuid.uuid4()}.pptx")
    prs.save(pptx_path)

    title = blueprint.get("title", prompt)[:80]

    source = SourcePresentation(
        owner_id=user_id,
        filename=f"[AI] {title}.pptx",
        file_path=pptx_path,
        file_type="pptx",
        slide_count=1,
        status="done",
        is_ai_source=True,
    )
    db.add(source)
    db.flush()

    # 5. Thumbnail
    thumb_dir = Path(settings.thumbnail_dir) / str(source.id)
    thumb_dir.mkdir(parents=True, exist_ok=True)
    try:
        if colors.background_image_path and os.path.exists(colors.background_image_path):
            thumb_abs = _render_thumbnail_branded(pptx_path, str(thumb_dir), colors, blueprint)
        else:
            thumb_abs = _render_thumbnail(pptx_path, str(thumb_dir))
    except Exception as thumb_err:
        logger.error("generate_slide thumbnail failed: %s", thumb_err)
        from services.thumbnail import _make_placeholder_thumbnail
        thumb_abs = str(thumb_dir / "0.png")
        Path(thumb_abs).write_bytes(_make_placeholder_thumbnail(title, 0))

    thumbnail_path = f"{source.id}/0.png"

    # 6. Vision QA (max 1 re-render)
    if settings.vision_model:
        blueprint = await _render_with_vision_qa(
            blueprint=blueprint, colors=colors,
            template_pptx_path=template_pptx_path,
            pptx_path=pptx_path, thumb_abs=thumb_abs,
            gen_dir=gen_dir, thumb_dir=thumb_dir,
            label=f"generate_slide/{title[:30]}",
            brand_context=brand_context,
        )

    # 7. XML blob
    xml_blob: str | None = None
    try:
        from lxml import etree
        prs_reload = Presentation(pptx_path)
        if prs_reload.slides:
            xml_blob = etree.tostring(prs_reload.slides[0]._element, encoding="unicode")
    except Exception:
        pass

    # 8. Embedding
    summary = _blueprint_to_summary(blueprint)
    tags    = _blueprint_to_tags(blueprint)
    embedding: list[float] | None = None
    try:
        from services.embedding import build_slide_embed_text
        embedding = await embed_single(build_slide_embed_text(title, summary, tags))
    except Exception as e:
        logger.warning("Embedding failed: %s", e)

    # 9. Save SlideLibraryEntry
    entry = SlideLibraryEntry(
        source_id      = source.id,
        slide_index    = 0,
        thumbnail_path = thumbnail_path,
        xml_blob       = xml_blob,
        title          = title,
        summary        = summary,
        tags_json      = json.dumps(tags, ensure_ascii=False),
        layout_type    = blueprint.get("layout", "title_content"),
        language       = _detect_language(prompt),
        embedding_json = json.dumps(embedding) if embedding else None,
        has_media      = False,
        blueprint_json = json.dumps(blueprint, ensure_ascii=False),
    )
    db.add(entry)
    db.commit()
    db.refresh(entry)
    db.refresh(source)
    entry.source = source
    return entry


async def save_slide_from_blueprint(
    db: Session,
    blueprint: dict,
    colors: "BrandColors",
    template_pptx_path: str | None,
    user_id: int | None,
    slide_index: int = 0,
    brand_context: BrandContext | None = None,
) -> "SlideLibraryEntry":
    """
    Render a pre-made blueprint → save as SlideLibraryEntry with thumbnail.
    Skips AI generation (blueprint already provided by presentation_planner).
    """
    blueprint.setdefault("layout", "title_content")
    blueprint.setdefault("title", "")
    blueprint.setdefault("content", {})

    validate_and_trim(blueprint)
    title = (blueprint.get("title") or "")[:80]

    prs = render_slide_pptx(blueprint, colors, template_pptx_path, brand_context)

    gen_dir = Path(settings.upload_dir) / "generated"
    gen_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = str(gen_dir / f"gen_{uuid.uuid4()}.pptx")
    prs.save(pptx_path)

    source = SourcePresentation(
        owner_id=user_id,
        filename=f"[AI] {title}.pptx",
        file_path=pptx_path,
        file_type="pptx",
        slide_count=1,
        status="done",
        is_ai_source=True,
    )
    db.add(source)
    db.flush()

    thumb_dir = Path(settings.thumbnail_dir) / str(source.id)
    thumb_dir.mkdir(parents=True, exist_ok=True)
    try:
        if colors.background_image_path and os.path.exists(colors.background_image_path):
            thumb_abs = _render_thumbnail_branded(pptx_path, str(thumb_dir), colors, blueprint)
        else:
            thumb_abs = _render_thumbnail(pptx_path, str(thumb_dir))
    except Exception as thumb_err:
        logger.error("save_slide_from_blueprint thumbnail failed: %s", thumb_err)
        from services.thumbnail import _make_placeholder_thumbnail
        thumb_abs = str(thumb_dir / "0.png")
        Path(thumb_abs).write_bytes(_make_placeholder_thumbnail(title, slide_index))

    thumbnail_path = f"{source.id}/0.png"

    if settings.vision_model:
        blueprint = await _render_with_vision_qa(
            blueprint=blueprint, colors=colors,
            template_pptx_path=template_pptx_path,
            pptx_path=pptx_path, thumb_abs=thumb_abs,
            gen_dir=gen_dir, thumb_dir=thumb_dir,
            label=f"slide_{slide_index}/{title[:30]}",
            brand_context=brand_context,
        )

    xml_blob: str | None = None
    try:
        from lxml import etree
        prs_reload = Presentation(pptx_path)
        if prs_reload.slides:
            xml_blob = etree.tostring(prs_reload.slides[0]._element, encoding="unicode")
    except Exception:
        pass

    summary = _blueprint_to_summary(blueprint)
    tags    = _blueprint_to_tags(blueprint)
    embedding: list[float] | None = None
    try:
        from services.embedding import build_slide_embed_text
        embedding = await embed_single(build_slide_embed_text(title, summary, tags))
    except Exception as e:
        logger.warning("Embedding failed: %s", e)

    entry = SlideLibraryEntry(
        source_id      = source.id,
        slide_index    = slide_index,
        thumbnail_path = thumbnail_path,
        xml_blob       = xml_blob,
        title          = title,
        summary        = summary,
        tags_json      = json.dumps(tags, ensure_ascii=False),
        layout_type    = blueprint.get("layout", "title_content"),
        language       = _detect_language(title),
        embedding_json = json.dumps(embedding) if embedding else None,
        has_media      = False,
        is_generated   = True,
        blueprint_json = json.dumps(blueprint, ensure_ascii=False),
    )
    db.add(entry)
    db.commit()
    db.refresh(entry)
    db.refresh(source)
    entry.source = source
    return entry


# ─── Helpers ─────────────────────────────────────────────────────────────────


def _blueprint_to_summary(bp: dict) -> str:
    layout  = bp.get("layout", "")
    content = bp.get("content", {})
    if layout == "big_stat":
        v, l = content.get("value", ""), content.get("label", "")
        return f"{v} — {l}".strip(" —")
    if layout == "quote":
        return (content.get("quote", "") or "")[:150]
    if layout == "section_divider":
        return content.get("subtitle", "")
    items = content.get("items", [])
    if not items:
        items  = content.get("left",  {}).get("items", [])
        items += content.get("right", {}).get("items", [])
    if not items:
        items = [s.get("event", "") for s in content.get("steps", [])]
    return "; ".join(items[:3])[:200]


def _blueprint_to_tags(bp: dict) -> list[str]:
    stop = {"и", "в", "на", "за", "по", "с", "а", "the", "a", "in", "of", "for", "and", "to", "is", "are"}
    tags = [bp.get("layout", "").replace("_", "-")]
    words = bp.get("title", "").lower().split()
    tags += [w for w in words if len(w) > 3 and w not in stop][:3]
    return [t for t in tags if t][:5]


def _detect_language(text: str) -> str:
    cyr = sum(1 for c in text if "\u0400" <= c <= "\u04FF")
    return "ru" if cyr > len(text) * 0.25 else "en"
