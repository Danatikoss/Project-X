"""
Thumbnail generation for PPTX and PDF slides.
Primary: PyMuPDF (fitz) — can open both formats.
Fallback: Pillow ImageDraw — white background + title text.
"""
import io
import os
import json
import logging
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class SlideData:
    index: int
    thumbnail_bytes: bytes
    text: str
    xml_blob: Optional[str] = None       # PPTX only — full slide XML for export
    slide_json: Optional[str] = None     # PDF only — page dict serialized
    has_video: bool = False              # slide contains embedded video
    has_gif: bool = False                # slide contains animated GIF
    gif_bytes: bytes | None = None      # animated GIF extracted from PPTX
    gif_rect: dict | None = None        # {x,y,w,h} as fractions of slide dimensions
    video_bytes: bytes | None = None    # video file (mp4/mov) extracted from PPTX
    video_ext: str | None = None        # 'mp4', 'mov', etc.
    pptx_title: str | None = None       # raw title from PPTX title shape (preferred over AI)


def _make_placeholder_thumbnail(title: str, slide_index: int) -> bytes:
    """Pillow fallback: white background + centered title text."""
    from PIL import Image, ImageDraw, ImageFont

    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), "#FFFFFF")
    draw = ImageDraw.Draw(img)

    label = title if title else f"Слайд {slide_index + 1}"
    # Try to use a system font; fall back to default
    try:
        font = ImageFont.truetype("arial.ttf", 64)
    except Exception:
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 64)
        except Exception:
            font = ImageFont.load_default()

    # Draw slide number watermark
    draw.rectangle([0, 0, W, H], fill="#F8FAFC")
    draw.rectangle([0, 0, W, 10], fill="#1E3A8A")

    # Center the title
    bbox = draw.textbbox((0, 0), label, font=font)
    text_w = bbox[2] - bbox[0]
    text_h = bbox[3] - bbox[1]
    x = (W - text_w) / 2
    y = (H - text_h) / 2
    draw.text((x, y), label, fill="#1E3A8A", font=font)

    # Slide index badge
    idx_label = str(slide_index + 1)
    draw.rectangle([20, 20, 100, 90], fill="#1E3A8A")
    draw.text((32, 32), idx_label, fill="#FFFFFF", font=font)

    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


_VIDEO_RELS = {
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video",
    "http://schemas.microsoft.com/office/2007/relationships/media",
}

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _get_slide_dimensions(zip_file) -> tuple[int, int]:
    """Read slide canvas size (EMUs) from ppt/presentation.xml. Default: 16:9 standard."""
    try:
        import xml.etree.ElementTree as ET
        prs_xml = zip_file.read("ppt/presentation.xml")
        root = ET.fromstring(prs_xml)
        sld_sz = root.find(f"{{{_P_NS}}}sldSz")
        if sld_sz is not None:
            return int(sld_sz.get("cx", 9144000)), int(sld_sz.get("cy", 5143500))
    except Exception:
        pass
    return 9144000, 5143500


def _extract_pic_xfrm(pic_el) -> tuple[int, int, int, int] | None:
    """Extract (x, y, cx, cy) in EMU from a p:pic element's spPr/a:xfrm."""
    sp_pr = pic_el.find(f"{{{_P_NS}}}spPr")
    if sp_pr is None:
        return None
    xfrm = sp_pr.find(f"{{{_A_NS}}}xfrm")
    if xfrm is None:
        return None
    off = xfrm.find(f"{{{_A_NS}}}off")
    ext = xfrm.find(f"{{{_A_NS}}}ext")
    if off is None or ext is None:
        return None
    return (
        int(off.get("x", 0)),
        int(off.get("y", 0)),
        int(ext.get("cx", 0)),
        int(ext.get("cy", 0)),
    )


def _find_gif_rect_in_slide_xml(
    zip_file, slide_xml_name: str, gif_rId: str, slide_cx: int, slide_cy: int
) -> dict | None:
    """
    Parse slide XML to find the p:pic whose blip matches gif_rId.
    Handles GIFs nested inside group shapes (p:grpSp) with correct coordinate transform.
    Returns {x, y, w, h} as fractions [0..1] of the slide canvas, or None.
    """
    import xml.etree.ElementTree as ET

    def _pic_rId(pic_el) -> str | None:
        """Return the r:embed rId of a p:pic element, or None."""
        bf = pic_el.find(f"{{{_P_NS}}}blipFill")
        if bf is None:
            return None
        blip = bf.find(f"{{{_A_NS}}}blip")
        if blip is None:
            return None
        return blip.get(f"{{{_R_NS}}}embed")

    def _search(container, parent_transform=None):
        """
        Recursively search container for the pic with gif_rId.
        parent_transform: None (slide coords) or (off_x, off_y, ext_cx, ext_cy, chOff_x, chOff_y, chExt_cx, chExt_cy)
        """
        for child in container:
            tag = child.tag

            if tag == f"{{{_P_NS}}}pic":
                if _pic_rId(child) != gif_rId:
                    continue
                coords = _extract_pic_xfrm(child)
                if coords is None:
                    continue
                x, y, cx, cy = coords

                # Apply group coordinate transform if inside a group
                if parent_transform:
                    gx, gy, gcx, gcy, chx, chy, chcx, chcy = parent_transform
                    if chcx > 0 and chcy > 0:
                        sx = gcx / chcx
                        sy = gcy / chcy
                        x = int(gx + (x - chx) * sx)
                        y = int(gy + (y - chy) * sy)
                        cx = int(cx * sx)
                        cy = int(cy * sy)

                x_frac = x / slide_cx
                y_frac = y / slide_cy
                # Skip GIFs fully outside the slide (off-screen start positions for animations)
                if x_frac >= 1.0 or y_frac >= 1.0:
                    return None
                return {
                    "x": max(0.0, x_frac),
                    "y": max(0.0, y_frac),
                    "w": min(1.0, cx / slide_cx),
                    "h": min(1.0, cy / slide_cy),
                }

            elif tag == f"{{{_P_NS}}}grpSp":
                # Extract group's coordinate transform
                grp_transform = None
                gsp_pr = child.find(f"{{{_P_NS}}}grpSpPr")
                if gsp_pr is not None:
                    xfrm = gsp_pr.find(f"{{{_A_NS}}}xfrm")
                    if xfrm is not None:
                        off = xfrm.find(f"{{{_A_NS}}}off")
                        ext = xfrm.find(f"{{{_A_NS}}}ext")
                        choff = xfrm.find(f"{{{_A_NS}}}chOff")
                        chext = xfrm.find(f"{{{_A_NS}}}chExt")
                        if off and ext and choff and chext:
                            grp_transform = (
                                int(off.get("x", 0)), int(off.get("y", 0)),
                                int(ext.get("cx", 1)), int(ext.get("cy", 1)),
                                int(choff.get("x", 0)), int(choff.get("y", 0)),
                                int(chext.get("cx", 1)), int(chext.get("cy", 1)),
                            )
                result = _search(child, grp_transform or parent_transform)
                if result:
                    return result

        return None

    try:
        slide_xml = zip_file.read(slide_xml_name)
        root = ET.fromstring(slide_xml)
        cSld = root.find(f"{{{_P_NS}}}cSld")
        if cSld is None:
            return None
        spTree = cSld.find(f"{{{_P_NS}}}spTree")
        if spTree is None:
            return None
        return _search(spTree)
    except Exception as e:
        logger.debug(f"GIF rect extraction failed for {slide_xml_name}: {e}")
    return None


def _detect_media_in_pptx_slide(
    zip_file, slide_xml_name: str, slide_cx: int = 9144000, slide_cy: int = 5143500
) -> tuple[bool, bool, bytes | None, dict | None, bytes | None, str | None]:
    """
    Detect and extract media from a slide's relationships using ElementTree (reliable).
    Returns (has_video, has_gif, gif_bytes, gif_rect, video_bytes, video_ext).
    """
    import xml.etree.ElementTree as ET

    rels_name = slide_xml_name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
    has_video = False
    has_gif = False
    gif_bytes: bytes | None = None
    gif_rect: dict | None = None
    video_bytes: bytes | None = None
    video_ext: str | None = None

    if rels_name not in zip_file.namelist():
        return has_video, has_gif, gif_bytes, gif_rect, video_bytes, video_ext

    try:
        rels_root = ET.fromstring(zip_file.read(rels_name))
        namelist = set(zip_file.namelist())

        # Collect all GIF candidates to pick the most prominent one later
        gif_candidates: list[tuple[str, bytes]] = []  # [(rId, data), ...]

        for rel in rels_root:
            rId = rel.get("Id", "")
            rel_type = rel.get("Type", "")
            target = rel.get("Target", "")

            if target.startswith("../"):
                target = "ppt/" + target[3:]

            target_lower = target.lower()
            is_image = "relationships/image" in rel_type
            is_video = (
                "relationships/video" in rel_type
                or "relationships/media" in rel_type
                or "office/2007/relationships/media" in rel_type
            )

            if is_image and target_lower.endswith(".gif"):
                has_gif = True
                if target in namelist:
                    try:
                        data = zip_file.read(target)
                        gif_candidates.append((rId, data))
                    except Exception:
                        pass

            elif is_video:
                ext = target_lower.rsplit(".", 1)[-1] if "." in target_lower else ""
                if ext in ("mp4", "mov", "m4v"):
                    has_video = True
                    if target in namelist and video_bytes is None:
                        try:
                            raw = zip_file.read(target)
                            if len(raw) <= 150 * 1024 * 1024:
                                video_bytes = raw
                                video_ext = ext
                        except Exception:
                            pass
                elif ext in ("wmv", "avi"):
                    has_video = True

        # Pick the GIF with the largest on-screen area (most prominent)
        if gif_candidates:
            best_bytes: bytes | None = None
            best_rect: dict | None = None
            best_area = -1.0
            for rId, data in gif_candidates:
                rect = _find_gif_rect_in_slide_xml(zip_file, slide_xml_name, rId, slide_cx, slide_cy)
                if rect:
                    area = rect["w"] * rect["h"]
                    if area > best_area:
                        best_area = area
                        best_bytes = data
                        best_rect = rect
            if best_bytes is None:
                # No rect found for any GIF — just use first one
                best_bytes = gif_candidates[0][1]
                best_rect = _find_gif_rect_in_slide_xml(
                    zip_file, slide_xml_name, gif_candidates[0][0], slide_cx, slide_cy
                )
            gif_bytes = best_bytes
            gif_rect = best_rect

    except Exception as e:
        logger.debug(f"Media detection failed for {slide_xml_name}: {e}")

    return has_video, has_gif, gif_bytes, gif_rect, video_bytes, video_ext


def _add_media_overlay(img_bytes: bytes, has_video: bool, has_gif: bool) -> bytes:
    """Draw a play button or GIF badge overlay on the thumbnail."""
    from PIL import Image, ImageDraw, ImageFont
    import io

    img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    W, H = img.size
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    if has_video:
        # Semi-transparent dark circle with play triangle
        cx, cy = W // 2, H // 2
        r = min(W, H) // 8
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(0, 0, 0, 160))
        # Play triangle
        tri_size = r // 2
        triangle = [
            (cx - tri_size // 2, cy - tri_size),
            (cx - tri_size // 2, cy + tri_size),
            (cx + tri_size, cy),
        ]
        draw.polygon(triangle, fill=(255, 255, 255, 230))

    if has_gif:
        # GIF badge in bottom-left corner
        badge_w, badge_h = 120, 50
        draw.rectangle([10, H - badge_h - 10, 10 + badge_w, H - 10], fill=(50, 150, 50, 200))
        try:
            font = ImageFont.truetype("arial.ttf", 28)
        except Exception:
            font = ImageFont.load_default()
        draw.text((20, H - badge_h - 2), "GIF", fill=(255, 255, 255, 255), font=font)

    result = Image.alpha_composite(img, overlay).convert("RGB")
    buf = io.BytesIO()
    result.save(buf, "PNG")
    return buf.getvalue()


def _extract_text_from_pptx_slide(slide) -> str:
    """Extract all text from a python-pptx slide object."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text.strip())
    return " ".join(texts)


def _get_bg_image_bytes(obj) -> bytes | None:
    """Extract background picture fill image bytes from a slide, layout, or master element."""
    try:
        from pptx.oxml.ns import qn as _qn
        _R_EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        bg_elem = obj.background._bg
        bgPr = bg_elem.find(_qn("p:bgPr"))
        if bgPr is None:
            return None
        blipFill = bgPr.find(_qn("a:blipFill"))
        if blipFill is None:
            return None
        blip = blipFill.find(_qn("a:blip"))
        if blip is None:
            return None
        rId = blip.get(_R_EMBED)
        if not rId:
            return None
        return obj.part.related_parts[rId].blob
    except Exception:
        return None


def _render_pptx_slide_with_pillow(prs, slide_index: int, slide_cx: int, slide_cy: int) -> bytes:
    """
    Pure-Python renderer using python-pptx + Pillow.
    Extracts background image/color from slide, layout, or master.
    Renders text shapes at correct positions with Montserrat font when available.
    """
    from PIL import Image, ImageDraw, ImageFont

    TARGET_W, TARGET_H = 3840, 2160
    scale_x = TARGET_W / slide_cx
    scale_y = TARGET_H / slide_cy

    slide = prs.slides[slide_index]

    # Background: try slide → layout → master for picture fill, then solid fill
    bg_img_bytes: bytes | None = None
    bg_color = (255, 255, 255)
    for obj in [slide, slide.slide_layout, slide.slide_layout.slide_master]:
        img_bytes = _get_bg_image_bytes(obj)
        if img_bytes:
            bg_img_bytes = img_bytes
            break
        try:
            fill = obj.background.fill
            if fill.type is not None:
                rgb = fill.fore_color.rgb
                bg_color = (rgb[0], rgb[1], rgb[2])
                break
        except Exception:
            pass

    if bg_img_bytes:
        try:
            bg_img = Image.open(io.BytesIO(bg_img_bytes)).convert("RGB")
            img = bg_img.resize((TARGET_W, TARGET_H), Image.LANCZOS)
        except Exception:
            img = Image.new("RGB", (TARGET_W, TARGET_H), bg_color)
    else:
        img = Image.new("RGB", (TARGET_W, TARGET_H), bg_color)

    draw = ImageDraw.Draw(img)

    def _font(size_pt: float, bold: bool = False) -> ImageFont.FreeTypeFont:
        size_px = max(8, int(size_pt * 96 / 72 * scale_y))
        candidates = (
            [
                "/usr/share/fonts/montserrat/Montserrat.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            ] if bold else []
        ) + [
            "/usr/share/fonts/montserrat/Montserrat.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
            "/Library/Fonts/Arial.ttf",
            "arial.ttf",
        ]
        for path in candidates:
            try:
                return ImageFont.truetype(path, size_px)
            except Exception:
                continue
        try:
            return ImageFont.load_default(size=size_px)
        except Exception:
            return ImageFont.load_default()

    def _draw_shapes(shapes):
        for shape in shapes:
            try:
                if shape.left is None or shape.top is None:
                    continue
                left = int(shape.left * scale_x)
                top = int(shape.top * scale_y)
                width = int((shape.width or 0) * scale_x)
                height = int((shape.height or 0) * scale_y)
                right, bottom = left + width, top + height

                # Solid shape fill
                try:
                    f = shape.fill
                    if str(f.type) == "SOLID (1)":
                        try:
                            rgb = f.fore_color.rgb
                            draw.rectangle([left, top, right, bottom], fill=(rgb[0], rgb[1], rgb[2]))
                        except Exception:
                            # Theme color — try to get srgbClr or sysClr from XML
                            from pptx.oxml.ns import qn as _qn
                            spPr = shape._element.find(_qn("p:spPr"))
                            if spPr is None:
                                spPr = shape._element
                            solidFill = spPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
                            if solidFill is not None:
                                srgb = solidFill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
                                if srgb is not None:
                                    val = srgb.get("val", "")
                                    if len(val) == 6:
                                        r2 = int(val[0:2], 16)
                                        g2 = int(val[2:4], 16)
                                        b2 = int(val[4:6], 16)
                                        draw.rectangle([left, top, right, bottom], fill=(r2, g2, b2))
                except Exception:
                    pass

                # Picture placeholder
                if shape.shape_type == 13:
                    draw.rectangle([left, top, right, bottom], fill=(215, 215, 220), outline=(170, 170, 180), width=2)
                    cx2, cy2 = (left + right) // 2, (top + bottom) // 2
                    r = min(width, height) // 6
                    draw.line([cx2 - r, cy2 - r, cx2 + r, cy2 + r], fill=(160, 160, 170), width=max(2, r // 4))
                    draw.line([cx2 + r, cy2 - r, cx2 - r, cy2 + r], fill=(160, 160, 170), width=max(2, r // 4))
                    continue

                # Text frame
                if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
                    continue

                y = top + max(2, int(4 * scale_y))
                for para in shape.text_frame.paragraphs:
                    line_text = para.text
                    if not line_text:
                        y += max(4, int(6 * scale_y))
                        continue

                    font_pt, text_color, bold = 18.0, (30, 30, 30), False
                    try:
                        if para.runs:
                            run = para.runs[0]
                            if run.font.size:
                                font_pt = run.font.size.pt
                            bold = bool(run.font.bold)
                            if run.font.color and run.font.color.type is not None:
                                rgb = run.font.color.rgb
                                text_color = (rgb[0], rgb[1], rgb[2])
                    except Exception:
                        pass

                    font = _font(font_pt, bold)

                    words, line_buf = line_text.split(), ""
                    for word in words:
                        candidate = (line_buf + " " + word).strip()
                        lw = draw.textbbox((0, 0), candidate, font=font)[2]
                        if lw > width - 8 and line_buf:
                            if y < bottom:
                                draw.text((left + 4, y), line_buf, fill=text_color, font=font)
                            lh = draw.textbbox((0, 0), line_buf, font=font)[3]
                            y += int(lh * 1.15) + 1
                            line_buf = word
                        else:
                            line_buf = candidate
                    if line_buf and y < bottom:
                        draw.text((left + 4, y), line_buf, fill=text_color, font=font)
                        lh = draw.textbbox((0, 0), line_buf, font=font)[3]
                        y += int(lh * 1.15) + 4

                    if y > bottom:
                        break
            except Exception:
                continue

    # Render layout shapes first (background), then slide shapes on top
    try:
        _draw_shapes(slide.slide_layout.shapes)
    except Exception:
        pass
    _draw_shapes(slide.shapes)

    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def render_single_slide_thumbnail(pptx_bytes: bytes, slide_index: int = 0) -> bytes:
    """
    Render one slide to a 1920×1080 PNG.
    Strategy: LibreOffice (high quality) → Pillow renderer (no system deps).
    Use this everywhere a single-slide thumbnail needs to be regenerated.
    """
    import tempfile, zipfile as _zf
    from pptx import Presentation as _Prs
    import io as _io

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(pptx_bytes)
        tmp_path = f.name

    try:
        # Strategy 1: LibreOffice → PDF → PyMuPDF
        pdf_path = _pptx_to_pdf_via_libreoffice(tmp_path)
        if pdf_path:
            try:
                import fitz
                doc = fitz.open(pdf_path)
                idx = min(slide_index, doc.page_count - 1)
                pix = doc[idx].get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                from PIL import Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img = img.resize((3840, 2160), Image.LANCZOS)
                doc.close()
                buf = _io.BytesIO()
                img.save(buf, "PNG")
                return buf.getvalue()
            finally:
                try:
                    import os as _os
                    _os.unlink(pdf_path)
                    import shutil as _sh
                    _sh.rmtree(os.path.dirname(pdf_path), ignore_errors=True)
                except Exception:
                    pass

        # Strategy 2: Pillow renderer
        prs = _Prs(_io.BytesIO(pptx_bytes))
        slide_cx, slide_cy = 9144000, 5143500
        try:
            with _zf.ZipFile(tmp_path, "r") as zf:
                slide_cx, slide_cy = _get_slide_dimensions(zf)
        except Exception:
            pass
        idx = min(slide_index, len(prs.slides) - 1)
        return _render_pptx_slide_with_pillow(prs, idx, slide_cx, slide_cy)
    finally:
        try:
            import os as _os
            _os.unlink(tmp_path)
        except Exception:
            pass


def _pptx_to_pdf_via_libreoffice(pptx_path: str) -> str | None:
    """
    Convert PPTX to PDF using LibreOffice headless.
    Returns path to the generated PDF, or None on failure.
    """
    import subprocess
    import tempfile
    import shutil

    soffice = shutil.which("soffice")
    if not soffice:
        logger.warning("soffice not found in PATH — falling back to direct PyMuPDF render")
        return None

    tmp_dir = tempfile.mkdtemp()
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, pptx_path],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0:
            logger.warning(f"LibreOffice conversion failed: {result.stderr[:200]}")
            return None

        base_name = Path(pptx_path).stem + ".pdf"
        pdf_path = os.path.join(tmp_dir, base_name)
        if os.path.exists(pdf_path):
            return pdf_path

        logger.warning(f"LibreOffice ran but PDF not found at {pdf_path}")
        return None
    except Exception as e:
        logger.warning(f"LibreOffice conversion error: {e}")
        return None


def extract_pptx_slides(file_path: str) -> list[SlideData]:
    """
    Extract slides from a PPTX file.
    Thumbnails: LibreOffice → PDF → PyMuPDF (high quality) with Pillow fallback (no system deps).
    XML blobs via python-pptx for lossless re-export.
    Detects video/GIF and adds visual overlay.
    """
    import fitz
    import zipfile
    import shutil
    from pptx import Presentation

    slides: list[SlideData] = []

    # --- Read slide dimensions first (needed for both render paths) ---
    slide_cx, slide_cy = 9144000, 5143500
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            slide_cx, slide_cy = _get_slide_dimensions(zf)
    except Exception:
        pass

    # --- Thumbnail rendering: LibreOffice → PDF → PyMuPDF, else Pillow ---
    thumbnail_map: dict[int, bytes] = {}
    pdf_tmp_dir = None
    try:
        pdf_path = _pptx_to_pdf_via_libreoffice(file_path)
        if pdf_path:
            pdf_tmp_dir = os.path.dirname(pdf_path)
            doc = fitz.open(pdf_path)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                from PIL import Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img = img.resize((3840, 2160), Image.LANCZOS)
                buf = io.BytesIO()
                img.save(buf, "PNG")
                thumbnail_map[i] = buf.getvalue()
            doc.close()
    except Exception as e:
        logger.warning(f"LibreOffice thumbnail render failed: {e}")
    finally:
        if pdf_tmp_dir and os.path.isdir(pdf_tmp_dir):
            try:
                shutil.rmtree(pdf_tmp_dir)
            except Exception:
                pass

    # Pillow fallback: render all slides that LibreOffice didn't produce
    if not thumbnail_map:
        logger.info("LibreOffice not available — using Pillow renderer for thumbnails")
        try:
            _prs_tmp = Presentation(file_path)
            for i in range(len(_prs_tmp.slides)):
                try:
                    thumbnail_map[i] = _render_pptx_slide_with_pillow(_prs_tmp, i, slide_cx, slide_cy)
                except Exception as e:
                    logger.warning(f"Pillow render failed for slide {i}: {e}")
        except Exception as e:
            logger.warning(f"Pillow thumbnail fallback failed: {e}")

    # --- XML blob + text extraction via python-pptx ---
    # python-pptx iterates slides in PRESENTATION ORDER (not file number order).
    # slide.part.partname gives the actual XML path for that slide — e.g. "/ppt/slides/slide3.xml"
    try:
        prs = Presentation(file_path)
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i, slide in enumerate(prs.slides):
                text = _extract_text_from_pptx_slide(slide)
                xml_blob = slide._element.xml

                # Extract title from PPTX title shape (preferred over AI-generated title)
                pptx_title: str | None = None
                try:
                    ts = slide.shapes.title
                    if ts is not None and hasattr(ts, "text"):
                        t = ts.text.strip()
                        # Accept only non-trivial titles (not just digits or very short)
                        if t and len(t) > 2 and not t.isdigit():
                            pptx_title = t
                except Exception:
                    pass

                # Get the actual XML file path for THIS slide (correct presentation order)
                slide_xml_name = str(slide.part.partname).lstrip("/")
                # e.g. "ppt/slides/slide3.xml" — NOT necessarily "ppt/slides/slide{i+1}.xml"

                # Detect video/GIF using the correct XML file
                has_video, has_gif = False, False
                gif_bytes: bytes | None = None
                gif_rect: dict | None = None
                video_bytes: bytes | None = None
                video_ext: str | None = None
                if slide_xml_name in zf.namelist():
                    has_video, has_gif, gif_bytes, gif_rect, video_bytes, video_ext = \
                        _detect_media_in_pptx_slide(zf, slide_xml_name, slide_cx, slide_cy)

                # Get thumbnail (from PyMuPDF map or fallback)
                if i in thumbnail_map:
                    thumb = thumbnail_map[i]
                    # Validate it's not blank (check if image has some content)
                    try:
                        import numpy as np
                        from PIL import Image
                        img = Image.open(io.BytesIO(thumb)).convert("RGB")
                        # Sample 1000 random pixels — avoids loading all 2M pixels into memory
                        arr = np.array(img).reshape(-1, 3)
                        rng = np.random.default_rng(42)
                        idx = rng.choice(len(arr), size=min(1000, len(arr)), replace=False)
                        sampled = arr[idx]
                        white_count = int(np.sum(np.all(sampled > 240, axis=1)))
                        if white_count / len(sampled) > 0.95:
                            logger.info(f"Slide {i} appears blank, re-rendering with Pillow")
                            try:
                                thumb = _render_pptx_slide_with_pillow(prs, i, slide_cx, slide_cy)
                            except Exception:
                                thumb = _make_placeholder_thumbnail(text[:60] if text else "", i)
                    except Exception:
                        pass  # keep original thumb
                else:
                    try:
                        thumb = _render_pptx_slide_with_pillow(prs, i, slide_cx, slide_cy)
                    except Exception:
                        thumb = _make_placeholder_thumbnail(text[:60] if text else "", i)

                # Add video/GIF overlay to thumbnail
                if has_video or has_gif:
                    try:
                        thumb = _add_media_overlay(thumb, has_video, has_gif)
                    except Exception as e:
                        logger.warning(f"Media overlay failed for slide {i}: {e}")

                slides.append(SlideData(
                    index=i,
                    thumbnail_bytes=thumb,
                    text=text,
                    xml_blob=xml_blob,
                    has_video=has_video,
                    has_gif=has_gif,
                    gif_bytes=gif_bytes,
                    gif_rect=gif_rect,
                    video_bytes=video_bytes,
                    video_ext=video_ext,
                    pptx_title=pptx_title,
                ))
    except Exception as e:
        logger.error(f"python-pptx extraction failed for {file_path}: {e}")
        raise

    return slides


def extract_pdf_slides(file_path: str) -> list[SlideData]:
    """Extract pages from a PDF file. Thumbnails via PyMuPDF."""
    import fitz

    slides: list[SlideData] = []
    doc = fitz.open(file_path)

    for i, page in enumerate(doc):
        # Render thumbnail
        mat = fitz.Matrix(3.0, 3.0)
        pix = page.get_pixmap(matrix=mat)
        from PIL import Image
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img = img.resize((1920, 1080), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, "PNG")
        thumb = buf.getvalue()

        # Extract text
        text = page.get_text("text").strip()

        # Store page dict as JSON for reference
        try:
            page_dict = page.get_text("dict")
            slide_json = json.dumps({"page_number": i, "width": page.rect.width,
                                     "height": page.rect.height}, ensure_ascii=False)
        except Exception:
            slide_json = json.dumps({"page_number": i})

        slides.append(SlideData(
            index=i,
            thumbnail_bytes=thumb,
            text=text,
            slide_json=slide_json,
        ))

    doc.close()
    return slides


def save_thumbnail(thumbnail_bytes: bytes, source_id: int, slide_index: int,
                   thumbnail_dir: str) -> str:
    """Save thumbnail to disk, return relative path."""
    source_dir = Path(thumbnail_dir) / str(source_id)
    source_dir.mkdir(parents=True, exist_ok=True)
    file_path = source_dir / f"{slide_index}.png"
    file_path.write_bytes(thumbnail_bytes)
    return f"{source_id}/{slide_index}.png"


def save_media(
    data: bytes, source_id: int, slide_index: int, thumbnail_dir: str, ext: str
) -> str:
    """Save extracted media (GIF/video) to disk. Returns relative path."""
    source_dir = Path(thumbnail_dir) / str(source_id)
    source_dir.mkdir(parents=True, exist_ok=True)
    file_path = source_dir / f"{slide_index}.{ext}"
    file_path.write_bytes(data)
    return f"{source_id}/{slide_index}.{ext}"
