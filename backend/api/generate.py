"""
Template-based slide generation API.

POST /generate/presentation       — full presentation from a prompt
POST /generate/slide              — single slide (returns PPTX)
POST /generate/create-assembly    — generate full presentation + save to editor
POST /generate/create-assembly-single — generate single slide + save to editor
GET  /generate/templates          — list available templates
POST /generate/templates/upload   — upload a new PPTX template
DELETE /generate/templates/{id}   — remove a template
"""
import io
import json
import logging
import uuid
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from typing import Optional
from sqlalchemy.orm import Session

from api.deps import get_current_user
from database import get_db
from models.user import User
from models.slide import SourcePresentation, SlideLibraryEntry
from models.company_profile import get_company_context
from models.assembly import AssembledPresentation
from services.template_generator import generate_presentation_plan, fill_single_slide
from services.template_library import load_catalog, get_template_by_id, get_title_slides, list_themes, CATALOG_PATH, TEMPLATES_DIR
from services.template_injector import inject_into_slide, inject_into_presentation
from services.thumbnail import extract_pptx_slides, save_thumbnail
from pptx import Presentation as PptxPresentation
from config import settings

logger = logging.getLogger(__name__)
router = APIRouter()


# ── Request / Response models ─────────────────────────────────────────────────

class GeneratePlanRequest(BaseModel):
    prompt: str = Field(..., min_length=5, max_length=4000)
    theme: str = "default"
    title_template_id: Optional[str] = None
    has_media: bool = False


class SlideInPlan(BaseModel):
    template_id: str
    slots: dict[str, str]
    has_media: bool = False
    slide_type: str = "template"        # "template" | "library"
    library_slide_id: Optional[int] = None
    library_thumbnail_url: Optional[str] = None
    library_title: Optional[str] = None


class PresentationPlan(BaseModel):
    title: str
    slides: list[SlideInPlan]
    theme: str = "default"
    title_template_id: Optional[str] = None
    plan_elapsed_seconds: Optional[float] = None


class GenerateSlideRequest(BaseModel):
    description: str = Field(..., min_length=5, max_length=1000)
    template_id: Optional[str] = None
    theme: str = "default"


class TemplateSlotInfo(BaseModel):
    id: str
    name: str
    description: str
    slots: dict[str, str]
    scenario_tags: list[str]
    theme: str = "default"
    layout_role: str = "content"
    elapsed_seconds: Optional[float] = None


# ── Helpers ───────────────────────────────────────────────────────────────────

def _source_has_slots(file_path: str, _cache: dict = {}) -> bool:
    """Return True if the PPTX at file_path has at least one slot_* named shape."""
    if file_path in _cache:
        return _cache[file_path]
    try:
        from pptx import Presentation as _Prs
        prs = _Prs(file_path)
        has = any(
            shape.name.startswith("slot_")
            for slide in prs.slides
            for shape in slide.shapes
        )
    except Exception:
        has = False
    _cache[file_path] = has
    return has


async def _enrich_plan_with_library_slides(
    slides: list[SlideInPlan],
    db: Session,
    user_id: int,
    threshold: float = 0.90,
) -> list[SlideInPlan]:
    """
    For each slide in plan, check if a relevant *user-uploaded* library slide exists.
    If similarity >= threshold, substitute with the library slide.
    Falls back to the original template slide if nothing matches.
    AI-generated slides are excluded from substitution to avoid circular references.
    Only library slides from slot-named PPTX sources are used (ensures visual consistency).
    """
    from services.embedding import embed_single
    from services.vector_search import hybrid_search

    result = []
    for slide in slides:
        query = " ".join(v for v in slide.slots.values() if v)[:500]
        if not query:
            result.append(slide)
            continue

        try:
            embedding = await embed_single(query)
            candidates = hybrid_search(db, query_embedding=embedding, query_text=query, top_k=10, user_id=user_id)
        except Exception:
            result.append(slide)
            continue

        # Only consider slides from user-uploaded presentations (not AI-generated)
        # AND only from sources that use slot_* named shapes (ensures visual consistency with templates)
        human_candidates = [
            (entry, score) for entry, score in candidates
            if entry.source
            and not entry.source.is_ai_source
            and entry.source.file_path
            and _source_has_slots(entry.source.file_path)
        ]

        if human_candidates and human_candidates[0][1] >= threshold:
            best, score = human_candidates[0]
            thumb_url = f"/thumbnails/{best.thumbnail_path}" if best.thumbnail_path else None
            result.append(SlideInPlan(
                template_id=slide.template_id,
                slots=slide.slots,
                has_media=slide.has_media,
                slide_type="library",
                library_slide_id=best.id,
                library_thumbnail_url=thumb_url,
                library_title=best.title or slide.slots.get("slot_title") or f"Слайд {best.id}",
            ))
        else:
            result.append(slide)

    return result


def _pptx_to_stream(prs: PptxPresentation) -> io.BytesIO:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _make_streaming_response(buf: io.BytesIO, filename: str) -> StreamingResponse:
    from urllib.parse import quote
    ascii_name = filename.encode("ascii", "ignore").decode("ascii") or "presentation.pptx"
    encoded_name = quote(filename, safe="")
    disposition = f"attachment; filename=\"{ascii_name}\"; filename*=UTF-8''{encoded_name}"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": disposition},
    )


# ── Endpoints ─────────────────────────────────────────────────────────────────

@router.get("/templates", response_model=list[TemplateSlotInfo])
def list_templates(current_user: User = Depends(get_current_user)):
    """Return all available slide templates with their slot schemas."""
    catalog = load_catalog()
    return [
        TemplateSlotInfo(
            id=t.id,
            name=t.name,
            description=t.description,
            slots=t.slots,
            scenario_tags=t.scenario_tags,
            theme=t.theme,
            layout_role=t.layout_role,
        )
        for t in catalog
    ]


@router.get("/themes", response_model=list[str])
def list_available_themes(current_user: User = Depends(get_current_user)):
    """Return all distinct themes present in the catalog."""
    return list_themes()


@router.get("/title-slides", response_model=list[TemplateSlotInfo])
def list_title_slides(
    theme: str = "default",
    current_user: User = Depends(get_current_user),
):
    """Return title slides available for a given theme."""
    slides = get_title_slides(theme=theme)
    return [
        TemplateSlotInfo(
            id=t.id,
            name=t.name,
            description=t.description,
            slots=t.slots,
            scenario_tags=t.scenario_tags,
            theme=t.theme,
            layout_role=t.layout_role,
        )
        for t in slides
    ]


@router.post("/plan", response_model=PresentationPlan)
async def create_plan(
    body: GeneratePlanRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """
    Step 1: Ask LLM to create a presentation plan.
    Returns structured plan (template IDs + slots) for preview before rendering.
    """
    # Guard: must have at least one content slide with a valid embedding
    full_catalog = load_catalog()
    content_slides = [t for t in full_catalog if t.layout_role == "content"]
    if not content_slides:
        raise HTTPException(
            status_code=400,
            detail="Нет шаблонов контентных слайдов. Загрузите шаблоны и нажмите Reindex."
        )
    has_embeddings = any(
        t.embedding and sum(abs(x) for x in t.embedding) > 1e-6
        for t in content_slides
    )
    if not has_embeddings:
        raise HTTPException(
            status_code=400,
            detail="Шаблоны загружены, но не проиндексированы. Нажмите Reindex в библиотеке шаблонов."
        )

    import time as _time
    _t0 = _time.perf_counter()

    try:
        company_context = get_company_context(db)
        plan = await generate_presentation_plan(prompt=body.prompt, theme=body.theme, has_media=body.has_media, company_context=company_context)
    except Exception as e:
        logger.error("Plan generation failed: %s", e)
        raise HTTPException(status_code=502, detail=f"Ошибка генерации плана: {e}")

    slides = plan.get("slides", [])
    if not slides:
        raise HTTPException(status_code=502, detail="LLM не вернул слайды")

    base_slides = [SlideInPlan(template_id=s["template_id"], slots=s.get("slots", {})) for s in slides]
    enriched = await _enrich_plan_with_library_slides(base_slides, db=db, user_id=current_user.id)

    elapsed = round(_time.perf_counter() - _t0, 2)
    logger.info("Plan generated: %d slides in %.2fs (user=%d)", len(enriched), elapsed, current_user.id)
    try:
        from api.admin import log_generation
        log_generation(db, "plan", elapsed, user_id=current_user.id, slide_count=len(enriched))
    except Exception:
        pass

    return PresentationPlan(
        title=plan.get("title", "Презентация"),
        slides=enriched,
        theme=body.theme,
        title_template_id=body.title_template_id,
        plan_elapsed_seconds=elapsed,
    )


@router.post("/download")
async def download_presentation(
    body: PresentationPlan,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """
    Step 2: Render a confirmed plan to PPTX and return as a download.
    Accepts the plan returned by /plan (possibly edited by user).
    """
    import time as _time
    _t0 = _time.perf_counter()

    body.slides = await _resolve_media_slides(body.slides)
    out_prs = _build_pptx_from_plan(body, db=db)

    if len(out_prs.slides) == 0:
        raise HTTPException(status_code=502, detail="Не удалось собрать ни одного слайда")

    elapsed = round(_time.perf_counter() - _t0, 2)
    logger.info("Presentation rendered: %d slides in %.2fs (user=%d)", len(out_prs.slides), elapsed, current_user.id)
    try:
        from api.admin import log_generation
        log_generation(db, "download", elapsed, user_id=current_user.id, slide_count=len(out_prs.slides))
    except Exception:
        pass

    filename = f"{body.title}.pptx"
    buf = _pptx_to_stream(out_prs)
    response = _make_streaming_response(buf, filename)
    response.headers["X-Generation-Time"] = str(elapsed)
    return response


async def _resolve_media_slides(slides: list[SlideInPlan]) -> list[SlideInPlan]:
    """
    For slides marked has_media=True that don't already use a media template,
    re-match to a media template and re-fill slots using current slot values as context.
    Falls back to the original slide silently if no media templates exist.
    """
    from services.template_generator import fill_single_slide

    full_catalog = load_catalog()
    result = []
    for slide in slides:
        if not slide.has_media:
            result.append(slide)
            continue

        # Already using a media template — nothing to do
        try:
            tmpl = get_template_by_id(slide.template_id, full_catalog)
            if any(k.startswith("slot_media_") for k in tmpl.slots):
                result.append(slide)
                continue
        except ValueError:
            pass

        # Build a content description from existing slots for re-filling
        content_desc = " ".join(v for v in slide.slots.values() if v)[:500]
        try:
            new_plan = await fill_single_slide(
                slide_description=content_desc or "Медиа слайд",
                has_media=True,
            )
            result.append(SlideInPlan(
                template_id=new_plan["template_id"],
                slots=new_plan["slots"],
                has_media=True,
            ))
            logger.info("Re-matched slide to media template %r", new_plan["template_id"])
        except Exception as e:
            logger.warning("Could not re-match slide to media template: %s — keeping original", e)
            result.append(slide)

    return result


def _build_pptx_from_plan(body: "PresentationPlan", db: Session | None = None) -> "PptxPresentation":
    """Shared helper: render a PresentationPlan into an in-memory PPTX."""
    from pptx import Presentation as PptxPrs
    from services.template_injector import PPTX_PATH as TPL_PATH
    from services.export import _clone_slide

    catalog = load_catalog()

    # Determine slide dimensions from the first template in the plan
    # (template files may differ in size from Libraryslides.pptx)
    slide_w, slide_h = None, None
    for s in (body.slides or []):
        if s.slide_type == "template":
            try:
                tmpl = get_template_by_id(s.template_id, catalog)
                dim_prs = PptxPrs(str(tmpl.pptx_path))
                slide_w, slide_h = dim_prs.slide_width, dim_prs.slide_height
                break
            except Exception:
                pass
    if slide_w is None:
        # Fallback to Libraryslides
        fallback = PptxPrs(str(TPL_PATH))
        slide_w, slide_h = fallback.slide_width, fallback.slide_height

    out_prs = PptxPrs()
    out_prs.slide_width = slide_w
    out_prs.slide_height = slide_h
    sldIdLst = out_prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)

    source_cache: dict = {}

    # Prepend title slide if selected
    if body.title_template_id:
        try:
            title_tmpl = get_template_by_id(body.title_template_id, catalog)
            title_slots = {k: (body.title if "title" in k or "name" in k else "") for k in title_tmpl.slots}
            inject_into_presentation(out_prs, title_tmpl, title_slots, source_cache=source_cache)
            logger.info("Prepended title slide %r", body.title_template_id)
        except Exception as e:
            logger.warning("Could not inject title slide %r: %s", body.title_template_id, e)

    for i, slide in enumerate(body.slides):
        # Library slide: clone directly from source PPTX
        if slide.slide_type == "library" and slide.library_slide_id and db is not None:
            entry = db.query(SlideLibraryEntry).get(slide.library_slide_id)
            if entry and entry.source and entry.source.file_path:
                try:
                    if _clone_slide(out_prs, entry.source.file_path, entry.slide_index):
                        logger.info("Cloned library slide %d into position %d", entry.id, i)
                        continue
                except Exception as e:
                    logger.warning("Library slide clone failed for id=%d: %s, falling back to template", entry.id, e)

        # Template slide (default path)
        try:
            tmpl = get_template_by_id(slide.template_id, catalog)
        except ValueError:
            logger.warning("Unknown template_id %r at slide %d, skipping", slide.template_id, i)
            continue
        try:
            inject_into_presentation(out_prs, tmpl, slide.slots, source_cache=source_cache)
        except Exception as e:
            logger.error("Inject failed for slide %d: %s", i, e)

    return out_prs


def _save_pptx_and_create_assembly(
    out_prs: "PptxPresentation",
    plan_title: str,
    plan_slides: list,
    owner_id: int,
    db: Session,
) -> int:
    """
    Save the rendered PPTX to disk, extract thumbnails, create library entries
    and an assembly. Returns assembly_id.
    """
    from datetime import datetime, timezone

    upload_dir = Path(settings.upload_dir)
    upload_dir.mkdir(parents=True, exist_ok=True)
    file_id = uuid.uuid4().hex
    safe_title = "".join(c for c in plan_title if c.isalnum() or c in " _-")[:80] or "presentation"
    filename = f"{safe_title}.pptx"
    file_path = upload_dir / f"{file_id}.pptx"

    buf = _pptx_to_stream(out_prs)
    file_path.write_bytes(buf.read())

    # Create source record
    source = SourcePresentation(
        owner_id=owner_id,
        filename=filename,
        file_path=str(file_path.resolve()),
        file_type="pptx",
        status="pending",
        is_ai_source=True,
    )
    db.add(source)
    db.commit()
    db.refresh(source)

    # Extract slide thumbnails (synchronous LibreOffice → PNG)
    try:
        slide_data_list = extract_pptx_slides(str(file_path))
    except Exception as e:
        source.status = "error"
        db.commit()
        raise HTTPException(status_code=502, detail=f"Ошибка извлечения слайдов: {e}")

    # Create SlideLibraryEntry records
    slide_ids: list[int] = []
    for i, sd in enumerate(slide_data_list):
        thumb_rel = save_thumbnail(sd.thumbnail_bytes, source.id, sd.index, settings.thumbnail_dir)

        # Pull a reasonable title from plan slide (library or template)
        plan_slide = plan_slides[i] if i < len(plan_slides) else None
        lib_title = getattr(plan_slide, "library_title", None) if getattr(plan_slide, "slide_type", "template") == "library" else None
        slots: dict = getattr(plan_slide, "slots", {}) or {}
        title = (
            lib_title
            or slots.get("slot_title")
            or slots.get("slot_product_name")
            or (slots.get("slot_main_card", "").split("\n")[0] if slots.get("slot_main_card") else "")
            or f"Слайд {i + 1}"
        )

        entry = SlideLibraryEntry(
            source_id=source.id,
            slide_index=sd.index,
            thumbnail_path=thumb_rel,
            xml_blob=getattr(sd, "xml_blob", None),
            slide_json=getattr(sd, "slide_json", None),
            title=title[:200],
            summary="",
            tags_json="[]",
            layout_type="content",
            language="ru",
            text_content=(sd.text[:5000] if sd.text else None),
            is_generated=False,
        )
        db.add(entry)
        db.flush()
        slide_ids.append(entry.id)

    source.status = "done"
    source.slide_count = len(slide_ids)
    source.indexed_at = datetime.now(timezone.utc)
    db.commit()

    # Create Assembly
    assembly = AssembledPresentation(
        owner_id=owner_id,
        title=plan_title,
        prompt="(AI генерация)",
        slide_ids_json=json.dumps(slide_ids),
        status="draft",
    )
    db.add(assembly)
    db.commit()
    db.refresh(assembly)

    return assembly.id


@router.post("/create-assembly")
async def create_assembly_from_plan(
    body: PresentationPlan,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """
    Generate a presentation from a confirmed plan, save slides to the library,
    and create an Assembly in the editor. Returns { assembly_id }.
    """
    import time as _time
    _t0 = _time.perf_counter()

    body.slides = await _resolve_media_slides(body.slides)
    out_prs = _build_pptx_from_plan(body, db=db)
    if len(out_prs.slides) == 0:
        raise HTTPException(status_code=502, detail="Не удалось собрать ни одного слайда")

    assembly_id = _save_pptx_and_create_assembly(
        out_prs=out_prs,
        plan_title=body.title,
        plan_slides=body.slides,
        owner_id=current_user.id,
        db=db,
    )
    elapsed = round(_time.perf_counter() - _t0, 2)
    logger.info("Assembly created: id=%d, %d slides in %.2fs (user=%d)", assembly_id, len(body.slides), elapsed, current_user.id)
    try:
        from api.admin import log_generation
        log_generation(db, "assembly", elapsed, user_id=current_user.id, slide_count=len(body.slides))
    except Exception:
        pass
    return {"assembly_id": assembly_id, "elapsed_seconds": elapsed}


class CreateAssemblySingleRequest(BaseModel):
    description: str = Field(..., min_length=5, max_length=1000)
    template_id: Optional[str] = None
    has_media: bool = False


@router.post("/create-assembly-single")
async def create_assembly_single(
    body: CreateAssemblySingleRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """
    Generate a single slide from a description, save to library,
    and create an Assembly in the editor. Returns { assembly_id }.
    """
    try:
        slide_plan = await fill_single_slide(
            slide_description=body.description,
            template_id=body.template_id,
            has_media=body.has_media,
        )
    except Exception as e:
        logger.error("Single slide generation failed: %s", e)
        raise HTTPException(status_code=502, detail=f"Ошибка генерации: {e}")

    template_id = slide_plan.get("template_id")
    slots = slide_plan.get("slots", {})

    try:
        tmpl = get_template_by_id(template_id)
    except ValueError:
        raise HTTPException(status_code=502, detail=f"Неизвестный шаблон: {template_id!r}")

    try:
        prs = inject_into_slide(tmpl, slots)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка рендера: {e}")

    # Wrap in PresentationPlan-like structure for the helper
    title = slots.get("slot_title") or slots.get("slot_product_name") or "Слайд"

    class _FakeSlide:
        def __init__(self, s: dict):
            self.slots = s

    assembly_id = _save_pptx_and_create_assembly(
        out_prs=prs,
        plan_title=title,
        plan_slides=[_FakeSlide(slots)],
        owner_id=current_user.id,
        db=db,
    )
    return {"assembly_id": assembly_id}


@router.post("/extract-file")
async def extract_file(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
):
    """
    Extract key content from a PDF or DOCX file.
    Returns a clean text summary suitable for use as a generation prompt.
    """
    from services.template_generator import extract_file_content

    filename = file.filename or "document"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in ("pdf", "docx", "doc"):
        raise HTTPException(status_code=400, detail="Поддерживаются только PDF и DOCX файлы")

    content = await file.read()
    if len(content) > 20 * 1024 * 1024:  # 20MB limit
        raise HTTPException(status_code=400, detail="Файл слишком большой (максимум 20MB)")

    try:
        summary = await extract_file_content(content, filename)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error("File extraction failed: %s", e)
        raise HTTPException(status_code=502, detail=f"Ошибка обработки файла: {e}")

    return {"summary": summary, "filename": filename}


@router.post("/slide")
async def generate_single_slide(
    body: GenerateSlideRequest,
    current_user: User = Depends(get_current_user),
):
    """
    Generate a single slide PPTX from a description.
    Optionally specify template_id; otherwise LLM picks one.
    """
    try:
        slide_plan = await fill_single_slide(
            slide_description=body.description,
            template_id=body.template_id,
        )
    except Exception as e:
        logger.error("Single slide generation failed: %s", e)
        raise HTTPException(status_code=502, detail=f"Ошибка генерации: {e}")

    template_id = slide_plan.get("template_id")
    slots = slide_plan.get("slots", {})

    try:
        tmpl = get_template_by_id(template_id)
    except ValueError:
        raise HTTPException(status_code=502, detail=f"Неизвестный шаблон: {template_id!r}")

    try:
        prs = inject_into_slide(tmpl, slots)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка рендера: {e}")

    buf = _pptx_to_stream(prs)
    return _make_streaming_response(buf, "slide.pptx")


# ── Template management ───────────────────────────────────────────────────────

def _save_catalog(catalog_data: list[dict]) -> None:
    with open(CATALOG_PATH, "w", encoding="utf-8") as f:
        json.dump(catalog_data, f, ensure_ascii=False, indent=2)


async def _analyze_slide_metadata(
    slot_names: list[str],
    text_previews: list[str],
) -> dict:
    """
    Auto-generate name, description, scenario_tags, and ai_description for a slide template.
    Uses gpt-4o-mini, text only, no vision.
    Returns dict: {name, description, scenario_tags, ai_description}
    """
    from openai import AsyncOpenAI

    _oai_kwargs: dict = {"api_key": settings.openai_api_key}
    if settings.openai_base_url:
        _oai_kwargs["base_url"] = settings.openai_base_url
    client = AsyncOpenAI(**_oai_kwargs)

    slots_str = ", ".join(slot_names)

    prompt = (
        f"You are classifying a PowerPoint slide TEMPLATE by its structure.\n"
        f"Slot names (these define the layout, not the content): {slots_str}\n\n"
        f"Your job: figure out what KIND of slide this is based purely on its slot structure.\n"
        f"Ignore any specific topic — this is a reusable blank template.\n\n"
        f"Examples of good names: '4 ключевых метрики', '2 колонки с текстом', 'До и после', "
        f"'3 шага процесса', 'Цель и результат', 'Титульный слайд', 'Статистика + описание'\n\n"
        f"Return JSON with exactly these fields:\n"
        f'  "name": structural template name in Russian, 2-5 words — describe the LAYOUT, not a topic\n'
        f'  "description": one sentence in Russian: what kind of content fits this layout?\n'
        f'  "scenario_tags": list of 4-6 English keywords describing layout/use case (e.g. ["key metrics", "kpi", "dashboard"])\n'
        f'  "ai_description": 2 sentences in English describing this template structure and when to use it\n'
    )

    response = await client.chat.completions.create(
        model="gpt-4o-mini",
        response_format={"type": "json_object"},
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2,
        max_tokens=200,
    )
    raw = json.loads(response.choices[0].message.content or "{}")
    return {
        "name": str(raw.get("name", "Шаблон слайда")),
        "description": str(raw.get("description", "")),
        "scenario_tags": raw.get("scenario_tags", []) if isinstance(raw.get("scenario_tags"), list) else [],
        "ai_description": str(raw.get("ai_description", "")),
    }


async def _auto_detect_slots(
    pptx_path: Path,
    slide_index: int,
    user_description: str,
    user_tags: str,
) -> tuple[dict[str, str], str, list[str]]:
    """
    Use GPT-4o vision to auto-detect content slots in a slide.
    Renames shapes in the PPTX file in-place.
    Returns (slots_dict, suggested_description, suggested_tags).
    """
    import base64
    import io as _io
    from openai import AsyncOpenAI
    from pptx import Presentation as PptxPrs
    from pptx.util import Emu
    from services.thumbnail import extract_pptx_slides

    # Render slide thumbnail
    try:
        slide_data_list = extract_pptx_slides(str(pptx_path))
        thumbnail_bytes = slide_data_list[slide_index].thumbnail_bytes if slide_index < len(slide_data_list) else None
    except Exception:
        thumbnail_bytes = None

    # Collect all text shapes info
    prs = PptxPrs(str(pptx_path))
    slide = prs.slides[slide_index]
    slide_w = prs.slide_width or Emu(9144000)
    slide_h = prs.slide_height or Emu(5143500)

    shapes_info = []
    for i, shape in enumerate(slide.shapes):
        if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()[:100]
        left_pct = round(shape.left / slide_w * 100, 1) if shape.left is not None else 0
        top_pct = round(shape.top / slide_h * 100, 1) if shape.top is not None else 0
        w_pct = round(shape.width / slide_w * 100, 1) if shape.width is not None else 0
        h_pct = round(shape.height / slide_h * 100, 1) if shape.height is not None else 0
        shapes_info.append({
            "index": i,
            "current_name": shape.name,
            "text_preview": text,
            "left%": left_pct,
            "top%": top_pct,
            "width%": w_pct,
            "height%": h_pct,
        })

    if not shapes_info:
        return {}, user_description, [t.strip() for t in user_tags.split(",") if t.strip()]

    shapes_json = json.dumps(shapes_info, ensure_ascii=False, indent=2)

    prompt = f"""You are analyzing a PowerPoint slide template to identify content slots.

Text shapes in this slide (positions as % of slide dimensions):
{shapes_json}

Task:
1. Assign a descriptive slot name (slot_XXXX) to each text shape that should be a content placeholder.
   Use names like: slot_title, slot_subtitle, slot_body, slot_description,
   slot_metric_1, slot_metric_2, slot_step_1..slot_step_N, slot_label, etc.
   Skip shapes that are logos, decorative, or non-content (e.g. copyright, tiny labels).
2. Suggest a short description of when to use this template (1 sentence, in Russian).
3. Suggest 4-6 scenario_tags (keywords in English, e.g. "key metrics", "product overview").

Return ONLY valid JSON:
{{
  "slots": {{"current_shape_name": "slot_suggested_name", ...}},
  "description": "...",
  "scenario_tags": ["tag1", "tag2", ...]
}}"""

    _oai_kwargs: dict = {"api_key": settings.openai_api_key}
    if settings.openai_base_url:
        _oai_kwargs["base_url"] = settings.openai_base_url
    client = AsyncOpenAI(**_oai_kwargs)

    user_content: list = []
    if thumbnail_bytes:
        small = thumbnail_bytes
        try:
            from PIL import Image
            img = Image.open(_io.BytesIO(thumbnail_bytes))
            if img.width > 768:
                ratio = 768 / img.width
                img = img.resize((768, int(img.height * ratio)), Image.LANCZOS)
            buf = _io.BytesIO()
            img.save(buf, "JPEG", quality=85)
            small = buf.getvalue()
        except Exception:
            pass
        user_content.append({
            "type": "image_url",
            "image_url": {
                "url": f"data:image/jpeg;base64,{base64.b64encode(small).decode()}",
                "detail": "low",
            },
        })
    user_content.append({"type": "text", "text": prompt})

    response = await client.chat.completions.create(
        model=settings.assembly_model,
        response_format={"type": "json_object"},
        messages=[{"role": "user", "content": user_content}],
        temperature=0.1,
        max_tokens=1000,
    )
    raw = json.loads(response.choices[0].message.content or "{}")

    slot_map: dict[str, str] = raw.get("slots", {})
    suggested_desc: str = raw.get("description", user_description) or user_description
    suggested_tags: list[str] = raw.get("scenario_tags", [])
    if not isinstance(suggested_tags, list):
        suggested_tags = []

    # Rename shapes in PPTX in-place
    if slot_map:
        for shape in slide.shapes:
            if shape.name in slot_map:
                shape.name = slot_map[shape.name]
        prs.save(str(pptx_path))

    # Build final slots dict: slot_name → placeholder description
    final_slots: dict[str, str] = {}
    prs2 = PptxPrs(str(pptx_path))
    slide2 = prs2.slides[slide_index]
    for shape in slide2.shapes:
        if shape.name.startswith("slot_") and hasattr(shape, "has_text_frame") and shape.has_text_frame:
            final_slots[shape.name] = f"Слот {shape.name}"

    # Merge user-provided tags with suggestions
    if user_tags.strip():
        user_tag_list = [t.strip() for t in user_tags.split(",") if t.strip()]
        for t in user_tag_list:
            if t not in suggested_tags:
                suggested_tags.insert(0, t)

    return final_slots, suggested_desc, suggested_tags


class SlotPreview(BaseModel):
    slide_index: int
    shapes: list[dict]  # [{original_name, suggested_slot, text_preview, position}]
    description: str
    scenario_tags: list[str]


@router.post("/templates/preview-slots", response_model=SlotPreview)
async def preview_slots(
    file: UploadFile = File(...),
    slide_index: int = Form(default=0),
    current_user: User = Depends(get_current_user),
):
    """
    Dry-run: анализирует PPTX и возвращает предполагаемые slot-имена без сохранения файла.
    Используй перед загрузкой, чтобы проверить точность автодетекции.
    """
    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Только администратор")

    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Файл должен быть .pptx")

    import tempfile, shutil
    content = await file.read()
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)

    try:
        from pptx import Presentation as PptxPrs
        import io as _io
        prs = PptxPrs(_io.BytesIO(content))
        if slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Слайд {slide_index} не существует")

        # Read original shape names before AI renames them
        slide = prs.slides[slide_index]
        original_names = {i: shape.name for i, shape in enumerate(slide.shapes)}

        ai_slots, ai_desc, ai_tags = await _auto_detect_slots(tmp_path, slide_index, "", "")

        # Map back: original_name → suggested slot name
        prs2 = PptxPrs(str(tmp_path))
        slide2 = prs2.slides[slide_index]
        from pptx.util import Emu
        slide_w = prs2.slide_width or Emu(9144000)
        slide_h = prs2.slide_height or Emu(5143500)

        shapes_out = []
        for i, shape in enumerate(slide2.shapes):
            if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                continue
            shapes_out.append({
                "original_name": original_names.get(i, shape.name),
                "suggested_slot": shape.name if shape.name.startswith("slot_") else None,
                "text_preview": shape.text_frame.text.strip()[:80],
                "position": {
                    "left%": round(shape.left / slide_w * 100, 1) if shape.left else 0,
                    "top%": round(shape.top / slide_h * 100, 1) if shape.top else 0,
                },
            })

        return SlotPreview(
            slide_index=slide_index,
            shapes=shapes_out,
            description=ai_desc,
            scenario_tags=ai_tags,
        )
    finally:
        tmp_path.unlink(missing_ok=True)


@router.post("/templates/upload", response_model=TemplateSlotInfo)
async def upload_template(
    file: UploadFile = File(...),
    name: str = Form(...),
    description: str = Form(default=""),
    scenario_tags: str = Form(default=""),  # comma-separated
    slide_index: int = Form(default=0),
    theme: str = Form(default="default"),
    layout_role: str = Form(default="content"),  # "title" | "content"
    current_user: User = Depends(get_current_user),
):
    """
    Upload a new PPTX slide template.
    If shapes are already named slot_*, they are used directly.
    Otherwise AI auto-detects and renames shapes automatically.
    Only admins can upload templates (they are shared across all users).
    """
    import time as _time
    _t0 = _time.perf_counter()

    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Только администратор может загружать шаблоны")

    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Файл должен быть в формате .pptx")

    # Save uploaded file
    uploads_dir = TEMPLATES_DIR / "uploads"
    uploads_dir.mkdir(exist_ok=True)
    file_id = uuid.uuid4().hex[:12]
    pptx_filename = f"uploads/{file_id}.pptx"
    pptx_path = TEMPLATES_DIR / pptx_filename

    content = await file.read()
    pptx_path.write_bytes(content)

    # Extract slot names from the specified slide
    try:
        from pptx import Presentation as PptxPrs
        import io as _io
        prs = PptxPrs(_io.BytesIO(content))
        if slide_index >= len(prs.slides):
            pptx_path.unlink(missing_ok=True)
            raise HTTPException(status_code=400, detail=f"Слайд {slide_index} не существует в файле (всего {len(prs.slides)})")
        slide = prs.slides[slide_index]
        slots = {
            shape.name: f"Слот {shape.name}"
            for shape in slide.shapes
            if shape.name.startswith("slot_") and hasattr(shape, "has_text_frame") and shape.has_text_frame
        }
        if not slots:
            # AI auto-detects and renames shapes in-place
            try:
                ai_slots, ai_desc, ai_tags = await _auto_detect_slots(
                    pptx_path, slide_index, description, scenario_tags
                )
                slots = ai_slots
                if not description and ai_desc:
                    description = ai_desc
                if not scenario_tags and ai_tags:
                    scenario_tags = ",".join(ai_tags)
            except Exception as ai_err:
                pptx_path.unlink(missing_ok=True)
                raise HTTPException(
                    status_code=400,
                    detail=f"Не удалось автоматически определить слоты: {ai_err}"
                )
        if not slots:
            pptx_path.unlink(missing_ok=True)
            raise HTTPException(
                status_code=400,
                detail="ИИ не смог найти текстовые блоки в слайде. Убедись, что слайд содержит текстовые элементы."
            )
    except HTTPException:
        raise
    except Exception as e:
        pptx_path.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"Не удалось прочитать PPTX: {e}")

    # Build template ID from name
    template_id = "custom_" + "".join(c if c.isalnum() else "_" for c in name.lower())[:30] + "_" + file_id[:6]

    # Parse scenario_tags
    tags = [t.strip() for t in scenario_tags.split(",") if t.strip()]

    theme_clean = theme.strip() or "default"
    role_clean = layout_role.strip() if layout_role.strip() in ("title", "content") else "content"

    # Generate AI description from slot names and tags (text-only, no vision)
    ai_description = ""
    try:
        ai_description = await _generate_ai_description(name, list(slots.keys()), tags)
        logger.info("Generated ai_description for template %r: %r", template_id, ai_description[:80])
    except Exception as e:
        logger.warning("Could not generate ai_description for %r: %s", template_id, e)

    new_entry = {
        "id": template_id,
        "slide_index": slide_index,
        "name": name,
        "description": description,
        "scenario_tags": tags,
        "slots": slots,
        "pptx_file": pptx_filename,
        "theme": theme_clean,
        "layout_role": role_clean,
        "ai_description": ai_description,
        "embedding": [],
    }

    # Append to catalog
    with open(CATALOG_PATH, encoding="utf-8") as f:
        catalog_data = json.load(f)
    catalog_data.append(new_entry)
    _save_catalog(catalog_data)

    elapsed = round(_time.perf_counter() - _t0, 2)
    logger.info("Uploaded new template %r in %.2fs (theme=%r, role=%r) by user %d", template_id, elapsed, theme_clean, role_clean, current_user.id)
    return TemplateSlotInfo(
        id=template_id,
        name=name,
        description=description,
        slots=slots,
        scenario_tags=tags,
        theme=theme_clean,
        layout_role=role_clean,
        elapsed_seconds=elapsed,
    )


class BatchUploadResult(BaseModel):
    created: int
    templates: list[TemplateSlotInfo]
    elapsed_seconds: float = 0.0


@router.post("/templates/upload-batch", response_model=BatchUploadResult)
async def upload_templates_batch(
    file: UploadFile = File(...),
    layout_role: str = Form(default="content"),
    current_user: User = Depends(get_current_user),
):
    """
    Upload a PPTX with one or more slides. Each slide that has slot_* shapes becomes
    a separate template. AI auto-generates name, description, and tags for each slide.
    """
    import time as _time
    _t0 = _time.perf_counter()

    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Только администратор может загружать шаблоны")

    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Файл должен быть в формате .pptx")

    uploads_dir = TEMPLATES_DIR / "uploads"
    uploads_dir.mkdir(exist_ok=True)
    file_id = uuid.uuid4().hex[:12]
    pptx_filename = f"uploads/{file_id}.pptx"
    pptx_path = TEMPLATES_DIR / pptx_filename

    content = await file.read()
    pptx_path.write_bytes(content)

    role_clean = layout_role.strip() if layout_role.strip() in ("title", "content") else "content"

    try:
        from pptx import Presentation as PptxPrs
        import io as _io
        prs = PptxPrs(_io.BytesIO(content))
    except Exception as e:
        pptx_path.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"Не удалось прочитать PPTX: {e}")

    new_entries: list[dict] = []
    for slide_index, slide in enumerate(prs.slides):
        # Read slot_* shapes — use actual shape text as hint (shows expected format)
        slots: dict[str, str] = {}
        for shape in slide.shapes:
            if (
                shape.name.startswith("slot_")
                and hasattr(shape, "has_text_frame")
                and shape.has_text_frame
            ):
                raw_text = shape.text_frame.text.strip()
                slots[shape.name] = raw_text if raw_text else f"Слот {shape.name}"
        if not slots:
            # AI auto-detects and renames shapes in-place
            try:
                ai_slots, ai_desc, ai_tags = await _auto_detect_slots(pptx_path, slide_index, "", "")
                slots = ai_slots
            except Exception as e:
                logger.warning("Auto-detect slots failed for slide %d: %s, skipping", slide_index, e)
            if not slots:
                continue

        # Collect visible text for AI context (non-slot shapes only, to avoid noise)
        text_previews = [
            shape.text_frame.text.strip()[:60]
            for shape in slide.shapes
            if (
                hasattr(shape, "has_text_frame")
                and shape.has_text_frame
                and shape.text_frame.text.strip()
                and not shape.name.startswith("slot_")
            )
        ]

        try:
            meta = await _analyze_slide_metadata(list(slots.keys()), text_previews)
        except Exception as e:
            logger.warning("AI metadata failed for slide %d: %s", slide_index, e)
            meta = {
                "name": f"Шаблон слайда {slide_index + 1}",
                "description": "",
                "scenario_tags": [],
                "ai_description": "",
            }

        # Generate embedding: concatenation of all semantic fields
        embed_text = " ".join(filter(None, [
            meta["name"],
            meta["description"],
            " ".join(meta["scenario_tags"]),
            meta["ai_description"],
        ]))
        try:
            from services.embedding import embed_single
            embedding = await embed_single(embed_text)
        except Exception as e:
            logger.warning("Embedding failed for slide %d: %s", slide_index, e)
            embedding = []

        safe_name = "".join(c if c.isalnum() else "_" for c in meta["name"].lower())[:30]
        template_id = f"custom_{safe_name}_{file_id[:6]}_{slide_index}"

        new_entries.append({
            "id": template_id,
            "slide_index": slide_index,
            "name": meta["name"],
            "description": meta["description"],
            "scenario_tags": meta["scenario_tags"],
            "slots": slots,
            "pptx_file": pptx_filename,
            "theme": "default",
            "layout_role": role_clean,
            "ai_description": meta["ai_description"],
            "embedding": embedding,
        })

    if not new_entries:
        pptx_path.unlink(missing_ok=True)
        raise HTTPException(
            status_code=400,
            detail="Ни один слайд не содержит слотов. Переименуй текстовые блоки в PowerPoint: названия должны начинаться с 'slot_' (например slot_title, slot_body)."
        )

    with open(CATALOG_PATH, encoding="utf-8") as f:
        catalog_data = json.load(f)
    catalog_data.extend(new_entries)
    _save_catalog(catalog_data)

    elapsed = round(_time.perf_counter() - _t0, 2)
    logger.info(
        "Batch upload: %d templates from %r by user %d in %.2fs",
        len(new_entries), file.filename, current_user.id, elapsed
    )
    return BatchUploadResult(
        created=len(new_entries),
        elapsed_seconds=elapsed,
        templates=[
            TemplateSlotInfo(
                id=e["id"],
                name=e["name"],
                description=e["description"],
                slots=e["slots"],
                scenario_tags=e["scenario_tags"],
                theme=e["theme"],
                layout_role=e["layout_role"],
            )
            for e in new_entries
        ],
    )


@router.post("/templates/reindex", status_code=200)
async def reindex_templates(
    current_user: User = Depends(get_current_user),
):
    """
    Regenerate embeddings for all templates that have an empty or near-zero embedding.
    Safe to run multiple times — skips templates that already have valid embeddings.
    """
    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Только администратор может запускать переиндексацию")

    from services.embedding import embed_single

    with open(CATALOG_PATH, encoding="utf-8") as f:
        catalog_data = json.load(f)

    def _needs_reindex(entry: dict) -> bool:
        emb = entry.get("embedding", [])
        if not emb:
            return True
        return sum(abs(x) for x in emb) < 1e-6

    updated = 0
    for entry in catalog_data:
        if not _needs_reindex(entry):
            continue
        embed_text = " ".join(filter(None, [
            entry.get("name", ""),
            entry.get("description", ""),
            " ".join(entry.get("scenario_tags", [])),
            entry.get("ai_description", ""),
        ]))
        try:
            entry["embedding"] = await embed_single(embed_text)
            updated += 1
            logger.info("Reindexed template %r", entry["id"])
        except Exception as e:
            logger.warning("Embedding failed for template %r: %s", entry["id"], e)

    _save_catalog(catalog_data)
    logger.info("Reindex complete: %d templates updated", updated)
    return {"updated": updated, "total": len(catalog_data)}


@router.delete("/templates/{template_id}", status_code=204)
def delete_template(
    template_id: str,
    current_user: User = Depends(get_current_user),
):
    """Remove a template from the catalog. Built-in templates cannot be deleted."""
    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Только администратор может удалять шаблоны")

    with open(CATALOG_PATH, encoding="utf-8") as f:
        catalog_data = json.load(f)

    entry = next((e for e in catalog_data if e["id"] == template_id), None)
    if not entry:
        raise HTTPException(status_code=404, detail="Шаблон не найден")

    # Remove file (only uploaded files, built-in pptx files are shared — don't delete them)
    pptx_file = entry.get("pptx_file", "")
    if pptx_file.startswith("uploads/"):
        pptx_path = TEMPLATES_DIR / pptx_file
        pptx_path.unlink(missing_ok=True)

    # Remove from catalog
    catalog_data = [e for e in catalog_data if e["id"] != template_id]
    _save_catalog(catalog_data)

    logger.info("Deleted template %r by user %d", template_id, current_user.id)
    return None


@router.delete("/templates", status_code=200)
def delete_all_custom_templates(
    current_user: User = Depends(get_current_user),
):
    """Remove all templates from the catalog. Uploaded files are deleted from disk; built-in pptx files are kept."""
    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Только администратор может удалять шаблоны")

    with open(CATALOG_PATH, encoding="utf-8") as f:
        catalog_data = json.load(f)

    total = len(catalog_data)
    for entry in catalog_data:
        if entry.get("pptx_file", "").startswith("uploads/"):
            pptx_path = TEMPLATES_DIR / entry["pptx_file"]
            pptx_path.unlink(missing_ok=True)

    _save_catalog([])
    logger.info("Deleted all %d templates by user %d", total, current_user.id)
    return {"deleted": total}
