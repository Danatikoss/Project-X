"""
WOPI protocol implementation for Collabora Online integration.

Feature flag: set COLLABORA_URL and WOPI_BASE_URL in .env to enable.
To disable: unset COLLABORA_URL — frontend won't show the Edit button.

Two router objects:
  api_router  — mounted at /api/wopi  (requires user JWT, returns editor URL)
  wopi_router — mounted at /wopi      (called by Collabora, uses WOPI token)
"""
import io
import logging
import urllib.parse
from datetime import datetime, timezone, timedelta
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, Request, Response
from jose import JWTError, jwt
from sqlalchemy.orm import Session

from api.deps import get_current_user
from config import settings
from database import get_db
from models.slide import SlideLibraryEntry
from models.user import User

logger = logging.getLogger(__name__)

api_router = APIRouter()
wopi_router = APIRouter()

_WOPI_TOKEN_EXPIRE_HOURS = 2


# ── Token helpers ─────────────────────────────────────────────────────────────

def _create_wopi_token(slide_id: int, user_id: int) -> str:
    exp = datetime.now(timezone.utc) + timedelta(hours=_WOPI_TOKEN_EXPIRE_HOURS)
    return jwt.encode(
        {"slide_id": slide_id, "user_id": user_id, "exp": exp},
        settings.jwt_secret,
        algorithm=settings.jwt_algorithm,
    )


def _decode_wopi_token(access_token: str) -> dict:
    try:
        payload = jwt.decode(
            access_token, settings.jwt_secret, algorithms=[settings.jwt_algorithm]
        )
        if "slide_id" not in payload or "user_id" not in payload:
            raise ValueError("missing claims")
        return payload
    except (JWTError, ValueError):
        raise HTTPException(status_code=401, detail="Invalid WOPI token")


# ── Single-slide PPTX builder ─────────────────────────────────────────────────

def _edited_pptx_path(slide_id: int) -> Path:
    """Filesystem path where we persist the user-edited version of a slide."""
    d = Path(settings.upload_dir) / "edited"
    d.mkdir(parents=True, exist_ok=True)
    return d / f"slide_{slide_id}.pptx"


def _build_single_slide_pptx(slide: SlideLibraryEntry) -> bytes:
    """
    Return a single-slide PPTX for this slide entry.
    Priority:
      1. User-edited version (if it exists)
      2. Clone from original source PPTX
      3. Fallback: embed thumbnail as a full-page image
    """
    # 1. Edited version
    ep = _edited_pptx_path(slide.id)
    if ep.exists():
        return ep.read_bytes()

    from pptx import Presentation as Prs
    from pptx.util import Inches

    # 2. Clone from source PPTX
    src = slide.source
    if src and src.file_type == "pptx" and src.file_path and Path(src.file_path).exists():
        try:
            from services.export import _clone_slide
            dest = Prs()
            dest.slide_width = Inches(13.33)
            dest.slide_height = Inches(7.5)
            if _clone_slide(dest, src.file_path, slide.slide_index):
                buf = io.BytesIO()
                dest.save(buf)
                return buf.getvalue()
        except Exception as exc:
            logger.warning("Clone failed for slide %d: %s", slide.id, exc)

    # 3. Fallback: thumbnail as image
    dest = Prs()
    dest.slide_width = Inches(13.33)
    dest.slide_height = Inches(7.5)
    sl = dest.slides.add_slide(dest.slide_layouts[6])
    thumb = Path(settings.thumbnail_dir) / (slide.thumbnail_path or "")
    if thumb.exists():
        sl.shapes.add_picture(
            io.BytesIO(thumb.read_bytes()),
            Inches(0), Inches(0),
            dest.slide_width, dest.slide_height,
        )
    buf = io.BytesIO()
    dest.save(buf)
    return buf.getvalue()


# ── API: token / editor URL (authenticated) ───────────────────────────────────

@api_router.get("/token/{slide_id}")
def get_wopi_editor_url(
    slide_id: int,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    """
    Return a short-lived WOPI access token and the Collabora editor URL
    for the given slide.  Frontend opens the URL in an iframe.
    """
    if not settings.collabora_url or not settings.wopi_base_url:
        raise HTTPException(
            status_code=503,
            detail="Collabora integration is not configured (set COLLABORA_URL and WOPI_BASE_URL)",
        )

    slide = db.query(SlideLibraryEntry).get(slide_id)
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    # Ownership check: source owner or admin
    if slide.source and slide.source.owner_id is not None:
        if slide.source.owner_id != user.id and not user.is_admin:
            raise HTTPException(status_code=403, detail="Access denied")

    token = _create_wopi_token(slide_id, user.id)
    wopi_src = f"{settings.wopi_base_url.rstrip('/')}/wopi/files/{slide_id}"
    editor_url = (
        f"{settings.collabora_url.rstrip('/')}/browser/dist/cool.html"
        f"?WOPISrc={urllib.parse.quote(wopi_src, safe='')}"
        f"&access_token={token}"
        f"&lang=ru"
    )
    return {"access_token": token, "editor_url": editor_url}


# ── WOPI callbacks (called by Collabora) ──────────────────────────────────────

@wopi_router.get("/files/{slide_id}")
def wopi_check_file_info(
    slide_id: int,
    access_token: str,
    db: Session = Depends(get_db),
):
    """WOPI CheckFileInfo — returns file metadata."""
    payload = _decode_wopi_token(access_token)
    if payload["slide_id"] != slide_id:
        raise HTTPException(status_code=403, detail="Token mismatch")

    slide = db.query(SlideLibraryEntry).get(slide_id)
    if not slide:
        raise HTTPException(status_code=404)

    user = db.query(User).get(payload["user_id"])
    pptx_bytes = _build_single_slide_pptx(slide)
    ts = (slide.updated_at or slide.created_at or datetime.now(timezone.utc))
    # Collabora expects UTC ISO-8601 with trailing Z
    last_modified = ts.strftime("%Y-%m-%dT%H:%M:%S.0000000Z")

    return {
        "BaseFileName": f"slide_{slide_id}.pptx",
        "Size": len(pptx_bytes),
        "OwnerId": str(payload["user_id"]),
        "UserId": str(payload["user_id"]),
        "UserFriendlyName": (user.name or user.email) if user else "User",
        "UserCanWrite": True,
        "UserCanNotWriteRelative": True,
        "DisablePrint": False,
        "DisableExport": False,
        "LastModifiedTime": last_modified,
    }


@wopi_router.get("/files/{slide_id}/contents")
def wopi_get_file(
    slide_id: int,
    access_token: str,
    db: Session = Depends(get_db),
):
    """WOPI GetFile — returns the PPTX bytes."""
    payload = _decode_wopi_token(access_token)
    if payload["slide_id"] != slide_id:
        raise HTTPException(status_code=403, detail="Token mismatch")

    slide = db.query(SlideLibraryEntry).get(slide_id)
    if not slide:
        raise HTTPException(status_code=404)

    pptx_bytes = _build_single_slide_pptx(slide)
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="slide_{slide_id}.pptx"'},
    )


@wopi_router.post("/files/{slide_id}/contents")
async def wopi_put_file(
    slide_id: int,
    access_token: str,
    request: Request,
    db: Session = Depends(get_db),
):
    """WOPI PutFile — saves the edited PPTX back."""
    payload = _decode_wopi_token(access_token)
    if payload["slide_id"] != slide_id:
        raise HTTPException(status_code=403, detail="Token mismatch")

    slide = db.query(SlideLibraryEntry).get(slide_id)
    if not slide:
        raise HTTPException(status_code=404)

    pptx_bytes = await request.body()
    if not pptx_bytes:
        # Collabora may send an empty PutFile on first open — treat as no-op
        return {"LastModifiedTime": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%S.0000000Z")}

    # Persist edited PPTX
    _edited_pptx_path(slide_id).write_bytes(pptx_bytes)

    # Extract new xml_blob from slide 0
    try:
        import lxml.etree as etree
        from pptx import Presentation as Prs
        prs = Prs(io.BytesIO(pptx_bytes))
        if prs.slides:
            slide.xml_blob = etree.tostring(
                prs.slides[0].shapes._spTree, encoding="unicode"
            )
    except Exception as exc:
        logger.warning("xml_blob extraction failed for slide %d: %s", slide_id, exc)

    slide.updated_at = datetime.now(timezone.utc)
    db.commit()

    logger.info("Slide %d saved via WOPI (user %d)", slide_id, payload["user_id"])
    return {
        "LastModifiedTime": slide.updated_at.strftime("%Y-%m-%dT%H:%M:%S.0000000Z")
    }
