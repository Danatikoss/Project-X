"""
Library API routes.
POST /api/library/upload
GET  /api/library/slides
GET  /api/library/slides/{id}
PATCH /api/library/slides/{id}
DELETE /api/library/slides/{id}
GET  /api/library/sources
"""
import asyncio
import json
import os
import shutil
import uuid
import logging
from pathlib import Path

from fastapi import APIRouter, Depends, File, UploadFile, HTTPException, Query
from sqlalchemy.orm import Session, joinedload

from config import settings
from database import get_db
from models.assembly import AssembledPresentation
from models.slide import SourcePresentation, SlideLibraryEntry, SlideEditVersion
from models.user import User
from api.schemas import (
    UploadResponse, SourcePresentationResponse,
    SlideResponse, SlidePatchRequest, SlideListResponse,
)
from api.utils import slide_to_response
from api.deps import get_current_user
from services.indexing import index_presentation

logger = logging.getLogger(__name__)
router = APIRouter()

# Keep strong references to running indexing tasks to prevent GC
_indexing_tasks: set[asyncio.Task] = set()

ALLOWED_TYPES = {"pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                 "pdf": "application/pdf"}


def _owned_source_ids(db: Session, user_id: int):
    """Return a subquery of source IDs belonging to this user."""
    return db.query(SourcePresentation.id).filter(SourcePresentation.owner_id == user_id).subquery()


def _check_slide_owner(slide: SlideLibraryEntry, user_id: int):
    if slide.source.owner_id != user_id:
        raise HTTPException(403, detail="Нет доступа к этому слайду")


@router.post("/upload", response_model=UploadResponse, status_code=202)
async def upload_presentation(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    filename = file.filename or "upload"
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext not in ("pptx", "pdf"):
        raise HTTPException(400, detail="Поддерживаются только файлы PPTX и PDF")

    max_bytes = settings.max_upload_size_mb * 1024 * 1024

    upload_dir = Path(settings.upload_dir)
    upload_dir.mkdir(parents=True, exist_ok=True)
    unique_name = f"{uuid.uuid4().hex}.{ext}"
    file_path = upload_dir / unique_name

    total = 0
    try:
        with open(file_path, "wb") as f:
            while chunk := await file.read(1024 * 1024):  # 1MB chunks
                total += len(chunk)
                if total > max_bytes:
                    raise HTTPException(
                        413,
                        detail=f"Файл превышает {settings.max_upload_size_mb} МБ",
                    )
                f.write(chunk)
    except HTTPException:
        file_path.unlink(missing_ok=True)
        raise

    source = SourcePresentation(
        owner_id=user.id,
        filename=filename,
        file_path=str(file_path.resolve()),
        file_type=ext,
        status="pending",
    )
    db.add(source)
    db.commit()
    db.refresh(source)

    ws_token = uuid.uuid4().hex

    task = asyncio.create_task(index_presentation(source.id, ws_token))
    _indexing_tasks.add(task)
    task.add_done_callback(_indexing_tasks.discard)

    return UploadResponse(
        source_id=source.id,
        ws_token=ws_token,
        message=f"Файл {filename} загружен. Индексация запущена.",
    )


@router.get("/slides", response_model=SlideListResponse)
def list_slides(
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=24, ge=1, le=100),
    source_id: int | None = Query(default=None),
    source_ids: list[int] = Query(default=[]),
    layout_type: str | None = Query(default=None),
    language: str | None = Query(default=None),
    tag: str | None = Query(default=None),
    label: str | None = Query(default=None),
    is_outdated: bool | None = Query(default=None),
    project_id: int | None = Query(default=None),
    project_ids: list[int] = Query(default=[]),
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    owned = _owned_source_ids(db, user.id)
    query = db.query(SlideLibraryEntry).filter(
        SlideLibraryEntry.source_id.in_(owned),
        SlideLibraryEntry.is_generated == False,  # noqa: E712
    )

    # Support both single and multi-select source filters
    effective_source_ids = list(source_ids)
    if source_id is not None and source_id not in effective_source_ids:
        effective_source_ids.append(source_id)
    if effective_source_ids:
        query = query.filter(SlideLibraryEntry.source_id.in_(effective_source_ids))

    if layout_type:
        query = query.filter(SlideLibraryEntry.layout_type == layout_type)
    if language:
        query = query.filter(SlideLibraryEntry.language == language)
    if is_outdated is not None:
        query = query.filter(SlideLibraryEntry.is_outdated == is_outdated)
    if tag:
        query = query.filter(SlideLibraryEntry.tags_json.ilike(f'%"{tag}"%'))
    if label:
        query = query.filter(SlideLibraryEntry.labels_json.ilike(f'%"{label}"%'))

    # Support both single and multi-select project filters
    effective_project_ids = list(project_ids)
    if project_id is not None and project_id not in effective_project_ids:
        effective_project_ids.append(project_id)
    if effective_project_ids:
        query = query.filter(SlideLibraryEntry.project_id.in_(effective_project_ids))

    total = query.count()
    slides = query.options(joinedload(SlideLibraryEntry.source)) \
                  .order_by(SlideLibraryEntry.id.desc()) \
                  .offset((page - 1) * page_size).limit(page_size).all()

    # Compute usage counts for the returned slides
    slide_ids_on_page = {s.id for s in slides}
    usage_counts: dict[int, int] = {}
    if slide_ids_on_page:
        rows = db.query(AssembledPresentation.slide_ids_json).filter(
            AssembledPresentation.owner_id == user.id
        ).all()
        for (ids_json,) in rows:
            for sid in json.loads(ids_json or "[]"):
                if sid in slide_ids_on_page:
                    usage_counts[sid] = usage_counts.get(sid, 0) + 1

    return SlideListResponse(
        items=[slide_to_response(s, used_in_assemblies=usage_counts.get(s.id, 0)) for s in slides],
        total=total,
        page=page,
        page_size=page_size,
    )


@router.get("/slides/{slide_id}", response_model=SlideResponse)
def get_slide(slide_id: int, db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    slide = db.query(SlideLibraryEntry).options(joinedload(SlideLibraryEntry.source)).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")
    _check_slide_owner(slide, user.id)
    return slide_to_response(slide, db)


@router.patch("/slides/{slide_id}", response_model=SlideResponse)
def update_slide(slide_id: int, body: SlidePatchRequest, db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    slide = db.query(SlideLibraryEntry).options(joinedload(SlideLibraryEntry.source)).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")
    _check_slide_owner(slide, user.id)

    if body.title is not None:
        slide.title = body.title
    if body.summary is not None:
        slide.summary = body.summary
    if body.tags is not None:
        slide.tags_json = json.dumps(body.tags, ensure_ascii=False)
    if body.layout_type is not None:
        slide.layout_type = body.layout_type
    if body.labels is not None:
        slide.labels_json = json.dumps(body.labels, ensure_ascii=False)
    if body.is_outdated is not None:
        slide.is_outdated = body.is_outdated
    if body.access_level is not None:
        slide.access_level = body.access_level
    if "project_id" in body.model_fields_set:
        slide.project_id = body.project_id  # None = unassign, int = assign

    db.commit()
    db.refresh(slide)
    return slide_to_response(slide, db)


@router.delete("/slides/all", status_code=200)
def delete_all_slides(db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    """Delete all library slides belonging to this user (non-generated only)."""
    owned = _owned_source_ids(db, user.id)
    result = db.query(SlideLibraryEntry).filter(
        SlideLibraryEntry.source_id.in_(owned),
        SlideLibraryEntry.is_generated == False,  # noqa: E712
    ).delete(synchronize_session=False)
    db.commit()
    return {"deleted": result}


@router.delete("/slides/{slide_id}", status_code=204)
def delete_slide(slide_id: int, db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    slide = db.query(SlideLibraryEntry).options(joinedload(SlideLibraryEntry.source)).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")
    _check_slide_owner(slide, user.id)
    db.delete(slide)
    db.commit()


@router.delete("/sources/all", status_code=200)
def delete_all_sources(db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    """Delete all sources (and their slides) belonging to this user."""
    sources = db.query(SourcePresentation).filter(
        SourcePresentation.owner_id == user.id
    ).all()
    count = len(sources)
    for source in sources:
        thumb_dir = Path(settings.thumbnail_dir) / str(source.id)
        if thumb_dir.exists():
            shutil.rmtree(thumb_dir)
        db.delete(source)
    db.commit()
    return {"deleted": count}


from pydantic import BaseModel as _BaseModel

class _SaveGeneratedRequest(_BaseModel):
    slide_ids: list[int]


@router.post("/slides/save-generated", status_code=200)
def save_generated_slides(
    body: _SaveGeneratedRequest,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    """Mark AI-generated slides as saved to library (is_generated=False)."""
    owned = _owned_source_ids(db, user.id)
    updated = db.query(SlideLibraryEntry).filter(
        SlideLibraryEntry.id.in_(body.slide_ids),
        SlideLibraryEntry.source_id.in_(owned),
        SlideLibraryEntry.is_generated == True,  # noqa: E712
    ).all()
    for slide in updated:
        slide.is_generated = False
    db.commit()
    return {"saved": len(updated)}


@router.get("/labels", response_model=list[str])
def list_labels(db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    """Return all unique user-defined labels across the user's slides."""
    owned = _owned_source_ids(db, user.id)
    rows = db.query(SlideLibraryEntry.labels_json).filter(
        SlideLibraryEntry.source_id.in_(owned),
        SlideLibraryEntry.labels_json != None,
        SlideLibraryEntry.labels_json != "[]",
    ).all()
    labels: set[str] = set()
    for (lj,) in rows:
        try:
            for lbl in json.loads(lj or "[]"):
                if lbl:
                    labels.add(lbl)
        except Exception:
            pass
    return sorted(labels)


@router.get("/sources", response_model=list[SourcePresentationResponse])
def list_sources(db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    sources = db.query(SourcePresentation).filter(
        SourcePresentation.owner_id == user.id
    ).order_by(SourcePresentation.id.desc()).all()
    return [SourcePresentationResponse.model_validate(s) for s in sources]


@router.post("/sources/{source_id}/extract-media", status_code=200)
def extract_media_from_source(source_id: int, db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    """
    Re-scan a previously indexed PPTX and extract GIF/video files that were
    missed (e.g. slides indexed before media extraction was implemented).
    Only processes slides that currently have gif_path=None or video_path=None.
    """
    import zipfile
    from services.thumbnail import _detect_media_in_pptx_slide, save_media

    source = db.query(SourcePresentation).get(source_id)
    if not source:
        raise HTTPException(404, detail="Источник не найден")
    if source.owner_id != user.id:
        raise HTTPException(403, detail="Нет доступа")
    if source.file_type != "pptx":
        raise HTTPException(400, detail="Только PPTX поддерживается")
    if not Path(source.file_path).exists():
        raise HTTPException(404, detail="Файл источника не найден на диске")

    slides = db.query(SlideLibraryEntry).filter(
        SlideLibraryEntry.source_id == source_id
    ).all()

    updated = 0
    try:
        from pptx import Presentation as _Prs
        from services.thumbnail import _get_slide_dimensions

        # Build correct slide_index → xml_name map using python-pptx (presentation order)
        prs = _Prs(source.file_path)
        slide_xml_by_idx = {
            i: str(s.part.partname).lstrip("/")
            for i, s in enumerate(prs.slides)
        }

        with zipfile.ZipFile(source.file_path, 'r') as zf:
            slide_cx, slide_cy = _get_slide_dimensions(zf)

            for slide in slides:
                needs_gif = not slide.gif_path
                needs_video = not slide.video_path
                needs_rect = not getattr(slide, 'gif_rect_json', None)
                if not (needs_gif or needs_video or needs_rect):
                    continue

                xml_name = slide_xml_by_idx.get(slide.slide_index)
                if not xml_name or xml_name not in zf.namelist():
                    continue

                _, _, gif_bytes, gif_rect, video_bytes, video_ext = \
                    _detect_media_in_pptx_slide(zf, xml_name, slide_cx, slide_cy)

                changed = False
                if needs_gif and gif_bytes:
                    slide.gif_path = save_media(gif_bytes, source_id, slide.slide_index,
                                                settings.thumbnail_dir, 'gif')
                    changed = True
                if gif_rect and needs_rect:
                    slide.gif_rect_json = json.dumps(gif_rect)
                    changed = True
                if needs_video and video_bytes and video_ext:
                    slide.video_path = save_media(video_bytes, source_id, slide.slide_index,
                                                  settings.thumbnail_dir, video_ext)
                    changed = True
                if changed:
                    updated += 1

        db.commit()
    except Exception as e:
        logger.error(f"extract-media failed for source {source_id}: {e}")
        raise HTTPException(500, detail=str(e))

    return {"updated": updated, "total": len(slides)}


@router.delete("/sources/{source_id}", status_code=204)
def delete_source(source_id: int, db: Session = Depends(get_db), user: User = Depends(get_current_user)):
    source = db.query(SourcePresentation).get(source_id)
    if not source:
        raise HTTPException(404, detail="Источник не найден")
    if source.owner_id != user.id:
        raise HTTPException(403, detail="Нет доступа")

    # Delete thumbnail directory
    thumb_dir = Path(settings.thumbnail_dir) / str(source_id)
    if thumb_dir.exists():
        shutil.rmtree(thumb_dir)

    # Cascade deletes slides (defined in model)
    db.delete(source)
    db.commit()


# ── Text editing endpoints ────────────────────────────────────────────────────

def _extract_text_elements(pptx_bytes: bytes) -> list[dict]:
    """Use python-pptx to extract positioned text elements from a 1-slide PPTX."""
    import io
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        logger.warning(f"Cannot open PPTX for text extraction: {e}")
        return []

    if not prs.slides:
        return []

    sl = prs.slides[0]
    slide_w = prs.slide_width or 1
    slide_h = prs.slide_height or 1
    elements = []

    for shape in sl.shapes:
        if not shape.has_text_frame:
            continue
        if shape.left is None or shape.top is None:
            continue

        tf = shape.text_frame
        lines = []
        for para in tf.paragraphs:
            line = "".join(run.text for run in para.runs if run.text)
            if line:
                lines.append(line)
        full_text = "\n".join(lines)
        if not full_text.strip():
            continue

        x_pct = round(shape.left / slide_w * 100, 2)
        y_pct = round(shape.top / slide_h * 100, 2)
        w_pct = round(shape.width / slide_w * 100, 2)
        h_pct = round(shape.height / slide_h * 100, 2)

        font_size = 18
        font_bold = False
        font_color = "#000000"
        font_align = "left"

        try:
            from pptx.enum.text import PP_ALIGN
            for para in tf.paragraphs:
                if not para.runs:
                    continue
                run = para.runs[0]
                if run.font.size:
                    font_size = max(6, int(run.font.size.pt))
                font_bold = bool(run.font.bold)
                try:
                    rgb = run.font.color.rgb
                    font_color = f"#{rgb}"
                except Exception:
                    pass
                if para.alignment == PP_ALIGN.CENTER:
                    font_align = "center"
                elif para.alignment == PP_ALIGN.RIGHT:
                    font_align = "right"
                break
        except Exception:
            pass

        elements.append({
            "id": str(shape.shape_id),
            "name": shape.name,
            "text": full_text,
            "x": x_pct,
            "y": y_pct,
            "w": w_pct,
            "h": h_pct,
            "font_size": font_size,
            "font_bold": font_bold,
            "font_color": font_color,
            "font_align": font_align,
        })

    return elements


def _apply_text_edits(pptx_bytes: bytes, edits: dict) -> bytes:
    """Replace text in specified shapes while preserving paragraph/run structure."""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    if not prs.slides:
        return pptx_bytes

    sl = prs.slides[0]
    for shape in sl.shapes:
        sid = str(shape.shape_id)
        if sid not in edits or not shape.has_text_frame:
            continue

        new_text = str(edits[sid])
        tf = shape.text_frame
        new_lines = new_text.split('\n')

        # Distribute lines only to paragraphs that have runs (skip spacing/empty paragraphs).
        # Use a separate line counter so empty paragraphs don't consume a line slot.
        line_idx = 0
        for para in tf.paragraphs:
            if not para.runs:
                continue
            line = new_lines[line_idx] if line_idx < len(new_lines) else ''
            para.runs[0].text = line
            for run in para.runs[1:]:
                run.text = ''
            line_idx += 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_original_pptx(slide: SlideLibraryEntry) -> bytes:
    """
    Always build a fresh single-slide PPTX from the ORIGINAL source file,
    never from a previously-edited version.  This is used as the clean base
    for applying accumulated text edits so we never compound corruptions.
    """
    import io as _io
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _Inches

    src = slide.source
    if src and src.file_type == "pptx" and src.file_path and Path(src.file_path).exists():
        try:
            from services.export import _clone_slide
            # Preserve original slide dimensions so EMU coordinates map correctly
            _src_prs = _Prs(src.file_path)
            dest = _Prs()
            dest.slide_width = _src_prs.slide_width or _Inches(13.33)
            dest.slide_height = _src_prs.slide_height or _Inches(7.5)
            if _clone_slide(dest, src.file_path, slide.slide_index):
                buf = _io.BytesIO()
                dest.save(buf)
                return buf.getvalue()
        except Exception as exc:
            logger.warning("Original PPTX clone failed for slide %d: %s", slide.id, exc)

    # Fallback: use whatever _build_single_slide_pptx returns
    from api.wopi import _build_single_slide_pptx
    return _build_single_slide_pptx(slide)


@router.get("/slides/{slide_id}/text-elements")
def get_text_elements(
    slide_id: int,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    """Return positioned text elements for a slide (for the inline text editor)."""
    slide = db.query(SlideLibraryEntry).options(
        joinedload(SlideLibraryEntry.source)
    ).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")

    # Always extract structure from ORIGINAL PPTX so positions/paragraphs are clean
    pptx_bytes = _build_original_pptx(slide)
    elements = _extract_text_elements(pptx_bytes)

    # Merge already-saved edits into the element texts
    saved_edits: dict = json.loads(slide.text_edits_json or "{}")
    for el in elements:
        if el["id"] in saved_edits:
            el["text"] = saved_edits[el["id"]]

    return {"elements": elements, "has_edits": bool(saved_edits)}


@router.post("/slides/{slide_id}/text-edits")
def save_text_edits(
    slide_id: int,
    body: dict,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    """
    Save user text edits for a slide.
    body = {"edits": {"shape_id": "new text", ...}}
    Persists a new single-slide PPTX with the edits applied and regenerates thumbnail.
    """
    import tempfile
    import time
    from datetime import datetime, timezone
    from api.wopi import _edited_pptx_path

    slide = db.query(SlideLibraryEntry).options(
        joinedload(SlideLibraryEntry.source)
    ).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")

    edits: dict = body.get("edits", {})
    if not edits:
        return {"ok": True, "edited": 0, "thumb_version": None}

    # Merge with any existing edits
    existing: dict = json.loads(slide.text_edits_json or "{}")
    existing.update(edits)

    # ALWAYS build from original source — never from the previously-edited file.
    # This prevents compounding corruptions when edits are re-applied.
    pptx_bytes = _build_original_pptx(slide)
    updated_bytes = _apply_text_edits(pptx_bytes, existing)

    # Persist edited PPTX (reusing WOPI infrastructure)
    _edited_pptx_path(slide_id).write_bytes(updated_bytes)

    # Update DB
    slide.text_edits_json = json.dumps(existing)
    slide.updated_at = datetime.now(timezone.utc)

    # Save version snapshot
    from sqlalchemy import func
    max_ver = db.query(func.max(SlideEditVersion.version_number)).filter(
        SlideEditVersion.slide_id == slide_id
    ).scalar() or 0
    version_entry = SlideEditVersion(
        slide_id=slide_id,
        version_number=max_ver + 1,
        edits_json=json.dumps(existing),
        created_by_id=user.id,
    )
    db.add(version_entry)
    db.commit()
    db.refresh(version_entry)

    # Regenerate thumbnail from the edited PPTX
    thumb_version: int | None = None
    if slide.thumbnail_path:
        thumb_path = Path(settings.thumbnail_dir) / slide.thumbnail_path
        if thumb_path.parent.exists():
            try:
                import fitz
                from services.thumbnail import _pptx_to_pdf_via_libreoffice

                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                    f.write(updated_bytes)
                    tmp_pptx = f.name

                try:
                    pdf_path = _pptx_to_pdf_via_libreoffice(tmp_pptx)
                    if pdf_path:
                        doc = fitz.open(pdf_path)
                        pix = doc[0].get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                        pix.save(str(thumb_path))
                        try:
                            os.unlink(pdf_path)
                        except Exception:
                            pass
                    else:
                        # Fallback: direct PyMuPDF render
                        doc = fitz.open(tmp_pptx)
                        if doc.page_count > 0:
                            pix = doc[0].get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                            pix.save(str(thumb_path))
                finally:
                    try:
                        os.unlink(tmp_pptx)
                    except Exception:
                        pass

                thumb_version = int(time.time())
            except Exception as e:
                logger.warning(f"Thumbnail regen failed for slide {slide_id}: {e}")

    return {
        "ok": True,
        "edited": len(existing),
        "thumb_version": thumb_version,
        "version_number": version_entry.version_number,
    }


@router.get("/slides/{slide_id}/edit-history")
def get_edit_history(
    slide_id: int,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    """Return the list of saved edit versions for a slide, newest first."""
    slide = db.query(SlideLibraryEntry).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")

    versions = (
        db.query(SlideEditVersion)
        .filter(SlideEditVersion.slide_id == slide_id)
        .order_by(SlideEditVersion.version_number.desc())
        .all()
    )

    result = []
    for v in versions:
        author_name = None
        if v.created_by_id:
            author = db.query(User).get(v.created_by_id)
            author_name = author.name if author else None
        result.append({
            "id": v.id,
            "version_number": v.version_number,
            "created_at": v.created_at.isoformat() + "Z",
            "created_by_name": author_name,
            "edit_count": len(json.loads(v.edits_json or "{}")),
        })

    return {"versions": result}


@router.post("/slides/{slide_id}/edit-history/{version_id}/rollback")
def rollback_edit_version(
    slide_id: int,
    version_id: int,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    """Restore slide text to a previously saved version."""
    import tempfile
    import time
    from datetime import datetime, timezone
    from api.wopi import _edited_pptx_path

    slide = db.query(SlideLibraryEntry).options(
        joinedload(SlideLibraryEntry.source)
    ).get(slide_id)
    if not slide:
        raise HTTPException(404, detail="Слайд не найден")

    version = db.query(SlideEditVersion).get(version_id)
    if not version or version.slide_id != slide_id:
        raise HTTPException(404, detail="Версия не найдена")

    target_edits: dict = json.loads(version.edits_json or "{}")

    # Rebuild PPTX from original and apply target version's edits
    pptx_bytes = _build_original_pptx(slide)
    updated_bytes = _apply_text_edits(pptx_bytes, target_edits)
    _edited_pptx_path(slide_id).write_bytes(updated_bytes)

    # Update slide
    slide.text_edits_json = json.dumps(target_edits)
    slide.updated_at = datetime.now(timezone.utc)

    # Record rollback as a new version entry
    from sqlalchemy import func
    max_ver = db.query(func.max(SlideEditVersion.version_number)).filter(
        SlideEditVersion.slide_id == slide_id
    ).scalar() or 0
    rollback_entry = SlideEditVersion(
        slide_id=slide_id,
        version_number=max_ver + 1,
        edits_json=json.dumps(target_edits),
        created_by_id=user.id,
    )
    db.add(rollback_entry)
    db.commit()
    db.refresh(rollback_entry)

    # Regenerate thumbnail
    thumb_version: int | None = None
    if slide.thumbnail_path:
        thumb_path = Path(settings.thumbnail_dir) / slide.thumbnail_path
        if thumb_path.parent.exists():
            try:
                import fitz
                from services.thumbnail import _pptx_to_pdf_via_libreoffice

                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                    f.write(updated_bytes)
                    tmp_pptx = f.name
                try:
                    pdf_path = _pptx_to_pdf_via_libreoffice(tmp_pptx)
                    if pdf_path:
                        doc = fitz.open(pdf_path)
                        pix = doc[0].get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                        pix.save(str(thumb_path))
                        try:
                            os.unlink(pdf_path)
                        except Exception:
                            pass
                    else:
                        doc = fitz.open(tmp_pptx)
                        if doc.page_count > 0:
                            pix = doc[0].get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                            pix.save(str(thumb_path))
                finally:
                    try:
                        os.unlink(tmp_pptx)
                    except Exception:
                        pass
                thumb_version = int(time.time())
            except Exception as e:
                logger.warning(f"Thumbnail regen failed on rollback for slide {slide_id}: {e}")

    return {
        "ok": True,
        "rolled_back_to_version": version.version_number,
        "new_version_number": rollback_entry.version_number,
        "thumb_version": thumb_version,
    }
