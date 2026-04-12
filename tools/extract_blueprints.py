"""
extract_blueprints.py — analyze PPTX files and extract slide style patterns.

Usage:
    python tools/extract_blueprints.py /path/to/pptx/folder

Output:
    tools/my_style_examples.json   — full per-slide data
    stdout                         — summary statistics
"""

import json
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE as PPT

# Slide canvas dimensions in EMU (standard 13.333" × 7.5")
_SLIDE_H_EMU = 6858000   # 7.5 inches
# Shapes in the top 20% of the slide are candidates for the title
_TITLE_ZONE_EMU = int(_SLIDE_H_EMU * 0.20)   # ≈ 1.5"


def _para_text(para) -> str:
    """Concatenate all run text in a paragraph."""
    return "".join(run.text for run in para.runs) or para.text


def _word_count(text: str) -> int:
    return len(text.split()) if text.strip() else 0


def _text_shapes_sorted(slide) -> list:
    """
    Return all shapes that carry text, sorted by vertical position (top).
    Excludes shapes with empty text.
    """
    shapes = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shapes.append((shape.top or 0, shape, text))
        except Exception:
            pass
    shapes.sort(key=lambda x: x[0])
    return shapes


def _find_title_placeholder(slide) -> tuple[object | None, str]:
    """
    Try placeholder-based title detection first (native PowerPoint files).
    Returns (placeholder_shape_or_None, title_text).
    """
    TITLE_TYPES = {PPT.TITLE, PPT.CENTER_TITLE}
    try:
        ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        ph = ph_map.get(0) or next(
            (p for p in slide.placeholders
             if p.placeholder_format.type in TITLE_TYPES),
            None,
        )
        if ph and ph.has_text_frame:
            text = ph.text_frame.text.strip()
            if text:
                return ph, text
    except Exception:
        pass
    return None, ""


def _find_title_positional(slide) -> tuple[object | None, str]:
    """
    Fallback for Google Slides exports (no placeholders).
    The topmost text shape in the upper 20% of the slide is the title.
    If nothing is that high up, take the single topmost text shape.
    """
    sorted_shapes = _text_shapes_sorted(slide)
    if not sorted_shapes:
        return None, ""

    # Prefer a shape clearly in the title zone
    for top, shape, text in sorted_shapes:
        if top <= _TITLE_ZONE_EMU:
            return shape, text

    # Nothing in title zone — take the very first text shape as a best guess
    _, shape, text = sorted_shapes[0]
    return shape, text


def _extract_slide(slide, filename: str, slide_index: int) -> dict:
    layout_name: str = ""
    try:
        layout_name = slide.slide_layout.name
    except Exception:
        pass

    # ── Title ──────────────────────────────────────────────────────────────────
    title_shape, title_text = _find_title_placeholder(slide)
    if not title_text:
        title_shape, title_text = _find_title_positional(slide)

    # ── Content blocks ─────────────────────────────────────────────────────────
    content_blocks: list[dict] = []
    has_image = False
    total_words = _word_count(title_text)

    for shape in slide.shapes:
        # Images
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                content_blocks.append({"type": "image"})
                continue
        except Exception:
            pass

        # Tables
        try:
            if shape.has_table:
                content_blocks.append({"type": "table"})
                continue
        except Exception:
            pass

        # Must have a text frame
        try:
            if not shape.has_text_frame:
                continue
        except Exception:
            continue

        # Skip the shape we used as the title
        if title_shape is not None and shape is title_shape:
            continue

        # Skip placeholder idx=0 (native title slot, already handled)
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                continue
        except Exception:
            pass

        tf = shape.text_frame
        paragraphs_text: list[str] = []
        for para in tf.paragraphs:
            t = _para_text(para).strip()
            if t:
                paragraphs_text.append(t)

        if not paragraphs_text:
            continue

        full_text = "\n".join(paragraphs_text)
        wc = _word_count(full_text)
        total_words += wc

        content_blocks.append(
            {
                "type": "text",
                "text": full_text,
                "word_count": wc,
                "bullet_count": len(paragraphs_text),
            }
        )

    is_single_statement = bool(
        title_text and (not content_blocks or total_words < 15)
    )

    return {
        "file": filename,
        "slide_index": slide_index,
        "layout_name": layout_name,
        "title": title_text,
        "title_length": len(title_text),
        "content_blocks": content_blocks,
        "has_image": has_image,
        "total_words": total_words,
        "is_single_statement": is_single_statement,
    }


def analyze_folder(folder: Path) -> list[dict]:
    pptx_files = sorted(folder.glob("*.pptx"))
    if not pptx_files:
        print(f"No .pptx files found in {folder}", file=sys.stderr)
        return []

    results: list[dict] = []
    for pptx_path in pptx_files:
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            print(f"  Skipping {pptx_path.name}: {e}", file=sys.stderr)
            continue

        for i, slide in enumerate(prs.slides, start=1):
            try:
                record = _extract_slide(slide, pptx_path.name, i)
                results.append(record)
            except Exception as e:
                print(
                    f"  Error on {pptx_path.name} slide {i}: {e}",
                    file=sys.stderr,
                )

    return results


def print_summary(slides: list[dict]) -> None:
    n = len(slides)
    if n == 0:
        print("No slides found.")
        return

    layout_counts = Counter(s["layout_name"] for s in slides)
    avg_title_len = sum(s["title_length"] for s in slides) / n

    bullet_counts = [
        b["bullet_count"]
        for s in slides
        for b in s["content_blocks"]
        if b["type"] == "text"
    ]
    avg_bullets = sum(bullet_counts) / len(bullet_counts) if bullet_counts else 0.0

    pct_single = 100 * sum(1 for s in slides if s["is_single_statement"]) / n
    pct_images = 100 * sum(1 for s in slides if s["has_image"]) / n

    print(f"\n{'='*50}")
    print(f"  SLIDE STYLE ANALYSIS")
    print(f"{'='*50}")
    print(f"  Total slides analyzed : {n}")
    print(f"  Average title length  : {avg_title_len:.1f} chars")
    print(f"  Average bullets/slide : {avg_bullets:.1f}")
    print(f"  Single-statement      : {pct_single:.1f}%")
    print(f"  Slides with images    : {pct_images:.1f}%")
    print(f"\n  Layout frequency:")
    for layout, count in layout_counts.most_common():
        bar = "█" * min(count, 40)
        label = layout or "(unnamed)"
        print(f"    {label:<35} {count:>4}  {bar}")
    print(f"{'='*50}\n")


def main() -> None:
    if len(sys.argv) != 2:
        print("Usage: python tools/extract_blueprints.py <folder>", file=sys.stderr)
        sys.exit(1)

    folder = Path(sys.argv[1])
    if not folder.is_dir():
        print(f"Not a directory: {folder}", file=sys.stderr)
        sys.exit(1)

    print(f"Scanning {folder} …")
    slides = analyze_folder(folder)

    out_path = Path(__file__).parent / "my_style_examples.json"
    out_path.write_text(json.dumps(slides, ensure_ascii=False, indent=2))
    print(f"Saved {len(slides)} slide records → {out_path}")

    print_summary(slides)


if __name__ == "__main__":
    main()
